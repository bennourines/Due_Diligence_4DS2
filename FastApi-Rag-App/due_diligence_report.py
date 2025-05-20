import os
from typing import List, Dict, Any, Optional
from pydantic import BaseModel
from fastapi import APIRouter, HTTPException
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import tempfile
from fastapi.responses import FileResponse
import matplotlib.pyplot as plt
import io
from PIL import Image
import numpy as np
import pandas as pd
from datetime import datetime

from enhanced_qa import CryptoQASystem

logger = logging.getLogger(__name__)

class ReportQuestion(BaseModel):
    question: str
    slide_title: str

class ReportSection(BaseModel):
    title: str
    questions: List[ReportQuestion]

class GenerateReportRequest(BaseModel):
    crypto_id: str
    report_title: str
    company_name: str = "Crypto Due Diligence"
    logo_path: Optional[str] = None
    sections: List[ReportSection]
    doc_id: str = "all"  # For filtering vector DB retrieval

class ReportResponse(BaseModel):
    status: str
    message: str
    file_path: Optional[str] = None


class DueDiligenceReportGenerator:
    """
    Generates PowerPoint due diligence reports for cryptocurrency projects
    by leveraging the CryptoQASystem to answer predefined questions.
    """
    
    def __init__(self, qa_system: CryptoQASystem):
        """Initialize with a CryptoQASystem instance."""
        self.qa_system = qa_system
        
    def create_presentation(self, 
                           crypto_id: str,
                           report_title: str,
                           company_name: str,
                           logo_path: Optional[str] = None) -> Presentation:
        """Create a new presentation with a title slide."""
        prs = Presentation()
        
        # Set slide dimensions to widescreen (16:9)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = f"{report_title}: {crypto_id.upper()}"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Add subtitle
        subtitle = slide.placeholders[1]
        current_date = datetime.now().strftime("%B %d, %Y")
        subtitle.text = f"Due Diligence Report\n{current_date}\n\n{company_name}"
        
        # Add logo if provided
        if logo_path and os.path.exists(logo_path):
            try:
                left = Inches(9)
                top = Inches(1)
                height = Inches(1.5)
                slide.shapes.add_picture(logo_path, left, top, height=height)
            except Exception as e:
                logger.error(f"Error adding logo: {str(e)}")
        
        return prs
    
    def add_section_title_slide(self, prs: Presentation, section_title: str) -> None:
        """Add a section title slide."""
        section_slide_layout = prs.slide_layouts[2]  # Section header layout
        slide = prs.slides.add_slide(section_slide_layout)
        
        title = slide.shapes.title
        title.text = section_title
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Add a colored rectangle at the bottom for design
        left = Inches(0)
        top = Inches(6)
        width = prs.slide_width
        height = Inches(1.5)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        
        # Set fill color to a blue shade
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
        shape.line.color.rgb = RGBColor(0, 112, 192)
    
    def add_content_slide(self, 
                         prs: Presentation, 
                         slide_title: str, 
                         question: str,
                         answer: str,
                         sources: List[Dict[str, Any]]) -> None:
        """Add a content slide with question, answer, and sources."""
        content_slide_layout = prs.slide_layouts[1]  # Content layout with title and content
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = slide_title
        title.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Get the content placeholder for the main text
        content_placeholder = slide.placeholders[1]
        
        # Add the question and answer
        text_frame = content_placeholder.text_frame
        text_frame.clear()  # Clear any existing text
        
        # Add the question paragraph
        question_p = text_frame.paragraphs[0]
        question_p.text = f"Question: {question}"
        question_p.font.bold = True
        question_p.font.size = Pt(14)
        
        # Add the answer
        answer_p = text_frame.add_paragraph()
        answer_p.text = f"Answer: {self._truncate_text(answer, 1000)}"
        answer_p.font.size = Pt(14)
        answer_p.space_after = Pt(12)
        
        # Add sources if available
        if sources:
            sources_p = text_frame.add_paragraph()
            sources_p.text = "Sources:"
            sources_p.font.bold = True
            sources_p.font.size = Pt(12)
            
            for idx, source in enumerate(sources[:3]):  # Limit to top 3 sources
                source_p = text_frame.add_paragraph()
                doc_id = source.get("doc_id", "Unknown")
                source_file = source.get("source", "Unknown").split("/")[-1]
                source_p.text = f"{idx+1}. {doc_id} ({source_file})"
                source_p.font.size = Pt(12)
                source_p.level = 1  # Indent as a sub-bullet
    
    def add_risk_assessment_slide(self, 
                                 prs: Presentation, 
                                 crypto_id: str,
                                 risk_data: Dict[str, float]) -> None:
        """Add a risk assessment slide with a radar chart."""
        slide_layout = prs.slide_layouts[5]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12), Inches(1)
        )
        title_frame = title_shape.text_frame
        title_p = title_frame.add_paragraph()
        title_p.text = f"Risk Assessment: {crypto_id.upper()}"
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        
        try:
            # Create radar chart using matplotlib
            categories = list(risk_data.keys())
            values = list(risk_data.values())
            
            # Number of variables
            N = len(categories)
            
            # Create angle values for the radar chart
            angles = np.linspace(0, 2*np.pi, N, endpoint=False).tolist()
            values += values[:1]  # Close the loop
            angles += angles[:1]  # Close the loop
            categories += categories[:1]  # Close the loop
            
            # Create plot
            fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(polar=True))
            
            # Plot data
            ax.plot(angles, values, 'o-', linewidth=2, label=crypto_id)
            ax.fill(angles, values, alpha=0.25)
            
            # Set category labels
            ax.set_thetagrids(np.degrees(angles[:-1]), categories[:-1])
            
            # Set y-axis limits
            ax.set_ylim(0, 1)
            ax.set_yticks([0.2, 0.4, 0.6, 0.8, 1.0])
            ax.set_yticklabels(['0.2', '0.4', '0.6', '0.8', '1.0'])
            
            # Add legend
            ax.legend(loc='upper right', bbox_to_anchor=(1.3, 1.1))
            
            # Add title
            plt.title("Risk Factors", size=15, y=1.1)
            
            # Save the chart to a byte buffer
            buf = io.BytesIO()
            plt.savefig(buf, format='png', dpi=300, bbox_inches='tight')
            buf.seek(0)
            
            # Add the image to the slide
            img = Image.open(buf)
            img_bytes = io.BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            
            # Add to slide
            left = Inches(2)
            top = Inches(1.5)
            slide.shapes.add_picture(img_bytes, left, top, height=Inches(5))
            
            buf.close()
            img_bytes.close()
            plt.close(fig)
            
        except Exception as e:
            logger.error(f"Error creating risk chart: {str(e)}")
            # Add error text to slide
            error_shape = slide.shapes.add_textbox(
                Inches(2), Inches(3), Inches(9), Inches(1)
            )
            error_shape.text_frame.text = "Error generating risk chart"
    
    def add_conclusion_slide(self, 
                            prs: Presentation, 
                            crypto_id: str,
                            summary: str) -> None:
        """Add a conclusion slide."""
        slide_layout = prs.slide_layouts[1]  # Content layout with title and content
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = f"Conclusion: {crypto_id.upper()}"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Get the content placeholder for the main text
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        # Add summary
        p = text_frame.paragraphs[0]
        p.text = summary
        p.font.size = Pt(16)
    
    def _truncate_text(self, text: str, max_length: int = 1000) -> str:
        """Truncate text to prevent overflow on slides."""
        if len(text) <= max_length:
            return text
        
        # Truncate at a sentence if possible
        truncated = text[:max_length]
        last_period = truncated.rfind('.')
        
        if last_period > max_length * 0.8:  # Only truncate at sentence if it's not too short
            return truncated[:last_period+1] + " [...]"
        else:
            return truncated + " [...]"
    
    def generate_report(self, request: GenerateReportRequest) -> str:
        """
        Generate a complete due diligence report based on the request.
        
        Returns:
            Path to the generated PPTX file
        """
        logger.info(f"Generating due diligence report for {request.crypto_id}")
        
        try:
            # Create presentation
            prs = self.create_presentation(
                crypto_id=request.crypto_id,
                report_title=request.report_title,
                company_name=request.company_name,
                logo_path=request.logo_path
            )
            
            # Process each section
            for section in request.sections:
                # Add section title slide
                self.add_section_title_slide(prs, section.title)
                
                # Process questions in this section
                for question_item in section.questions:
                    # Get answer from QA system
                    result = self.qa_system.answer_question(
                        question=question_item.question,
                        doc_id=request.doc_id
                    )
                    
                    answer = result.get("answer", "No answer found.")
                    sources = result.get("sources", [])
                    
                    # Add content slide
                    self.add_content_slide(
                        prs,
                        slide_title=question_item.slide_title,
                        question=question_item.question,
                        answer=answer,
                        sources=sources
                    )
            
            # Add risk assessment slide
            # In a real implementation, this would come from a risk scoring system
            # For this example, we'll use dummy data
            risk_data = {
                "Security Risk": 0.65,
                "Market Risk": 0.45,
                "Regulatory Risk": 0.78,
                "Team Risk": 0.3,
                "Technology Risk": 0.5
            }
            self.add_risk_assessment_slide(prs, request.crypto_id, risk_data)
            
            # Generate a summary for the conclusion
            summary_question = f"What is the overall assessment of {request.crypto_id} based on all factors, including security, market potential, regulatory concerns, team quality, and technology?"
            summary_result = self.qa_system.answer_question(
                question=summary_question,
                doc_id=request.doc_id
            )
            summary = summary_result.get("answer", "No summary available.")
            
            # Add conclusion slide
            self.add_conclusion_slide(prs, request.crypto_id, summary)
            
            # Save the presentation
            output_dir = "reports"
            os.makedirs(output_dir, exist_ok=True)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = f"{output_dir}/{request.crypto_id}_due_diligence_{timestamp}.pptx"
            
            prs.save(output_file)
            logger.info(f"Report saved to {output_file}")
            
            return output_file
            
        except Exception as e:
            logger.error(f"Error generating report: {str(e)}")
            raise RuntimeError(f"Failed to generate report: {str(e)}")


# Add to FastAPI app.py - Router for report generation


# Add this to your app.py
# app.include_router(report_router, prefix="/api", tags=["reports"])