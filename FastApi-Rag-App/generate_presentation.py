# generate_presentation.py
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from typing import List, Dict, Optional
from pptx import Presentation 
from pptx.util import Inches
import os
import uuid
from enhanced_qa import CryptoQASystem
from langchain_community.vectorstores import Qdrant
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

class QuestionItem(BaseModel):
    id: str
    text: str
    section_id: int
    section_title: str

class GeneratePresentationRequest(BaseModel):
    selected_questions: List[QuestionItem]
    doc_id: str = "all"  # Default to "all" documents if not specified
    coin_name: str = "Cryptocurrency"

class PresentationResponse(BaseModel):
    file_path: str
    message: str

# Add the missing classes that app.py is trying to import
class GenerateReportRequest(BaseModel):
    crypto_id: str
    sections: List[int] = []
    questions: List[str] = []
    include_market_analysis: bool = True
    include_sentiment_analysis: bool = True
    include_technical_analysis: bool = True
    include_forecast: bool = True
    include_risk_assessment: bool = True
    output_format: str = "pptx"  # pptx, pdf, or both

class ReportResponse(BaseModel):
    status: str
    message: str
    file_path: str

class DueDiligenceReportGenerator:
    """
    Class for generating cryptocurrency due diligence reports.
    """
    def __init__(self, qa_system):
        self.qa_system = qa_system

def add_answer_slides(prs, qa_slide, answer_text, question_text, bg_image_path, max_chars_per_slide=1400, 
                     font_family="Bodoni MT", answer_font_size=Pt(15)):
    """
    Split long answers across multiple slides.
    
    Args:
        prs: The presentation object
        qa_slide: The first slide to add the answer to
        answer_text: The full answer text
        question_text: The question text (for continuation slides)
        bg_image_path: Path to the background image
        max_chars_per_slide: Maximum characters per slide
        font_family: Font family to use
        answer_font_size: Font size for the answer
        
    Returns:
        Number of additional slides created
    """
    # Define styling constants
    WHITE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
    QUESTION_FONT_SIZE = Pt(18)
    
    # Define exact textbox dimensions to ensure text stays within slide boundaries
    textbox_left = Inches(0.5)
    textbox_top = Inches(1.7)
    textbox_width = Inches(12)
    textbox_height = Inches(4.5)  # Ensuring we don't go too close to the slide bottom
    
    # Split the answer into chunks of appropriate size
    # Try to split at sentence boundaries where possible
    chunks = []
    current_chunk = ""
    sentences = answer_text.replace(". ", ".|").replace("! ", "!|").replace("? ", "?|").split("|")
    
    for sentence in sentences:
        if len(current_chunk) + len(sentence) <= max_chars_per_slide:
            current_chunk += sentence + " "
        else:
            # If the current chunk has content, add it to chunks
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
            current_chunk = sentence + " "
    
    # Add the last chunk if it has content
    if current_chunk.strip():
        chunks.append(current_chunk.strip())
    
    # If there's only one chunk, just add it to the original slide
    if len(chunks) <= 1:
        answer_box = qa_slide.shapes.add_textbox(
            textbox_left, textbox_top, textbox_width, textbox_height)
        answer_frame = answer_box.text_frame
        answer_frame.word_wrap = True
        # Set auto-fit to fit text to shape
        answer_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        answer_paragraph = answer_frame.add_paragraph()
        answer_paragraph.text = answer_text
        answer_paragraph.font.color.rgb = WHITE_COLOR
        answer_paragraph.font.size = answer_font_size
        answer_paragraph.font.name = font_family
        return 0  # No additional slides created
    
    # Add first chunk to the original slide with continuation note
    answer_box = qa_slide.shapes.add_textbox(
        textbox_left, textbox_top, textbox_width, textbox_height)
    answer_frame = answer_box.text_frame
    answer_frame.word_wrap = True
    # Set auto-fit to fit text to shape
    answer_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    answer_paragraph = answer_frame.add_paragraph()
    answer_paragraph.text = chunks[0] + "\n\n(Continued on next slide...)"
    answer_paragraph.font.color.rgb = WHITE_COLOR
    answer_paragraph.font.size = answer_font_size
    answer_paragraph.font.name = font_family
    
    # Create additional slides for remaining chunks
    additional_slides_count = 0
    for i, chunk in enumerate(chunks[1:], 1):
        # Create continuation slide
        additional_slides_count += 1
        cont_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use content layout
        
        # Add background
        bg = cont_slide.shapes.add_picture(
            bg_image_path, 0, 0, 
            width=prs.slide_width, height=prs.slide_height)
        # Move to back
        cont_slide.shapes._spTree.remove(bg._element)
        cont_slide.shapes._spTree.insert(0, bg._element)
        
        # Add continuation title
        cont_title_box = cont_slide.shapes.add_textbox(
            textbox_left, Inches(0.5), textbox_width, Inches(1))
        cont_title_frame = cont_title_box.text_frame
        cont_title_para = cont_title_frame.add_paragraph()
        cont_title_para.text = f"(Continued) {question_text}"
        cont_title_para.font.color.rgb = WHITE_COLOR
        cont_title_para.font.size = QUESTION_FONT_SIZE
        cont_title_para.font.bold = True
        cont_title_para.font.name = font_family
        
        # Add continuation text
        cont_answer_box = cont_slide.shapes.add_textbox(
            textbox_left, textbox_top, textbox_width, textbox_height)
        cont_answer_frame = cont_answer_box.text_frame
        cont_answer_frame.word_wrap = True
        # Set auto-fit to fit text to shape
        cont_answer_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        cont_answer_para = cont_answer_frame.add_paragraph()
        
        # Add text with continuation notices
        if i < len(chunks) - 1:
            cont_answer_para.text = f"{chunk}\n\n(Continued on next slide...)"
        else:
            cont_answer_para.text = chunk  # Last chunk doesn't need continuation notice
            
        cont_answer_para.font.color.rgb = WHITE_COLOR
        cont_answer_para.font.size = answer_font_size
        cont_answer_para.font.name = font_family
    
    return additional_slides_count
    
async def generate_presentation(request: GeneratePresentationRequest):
    """
    Generate a PowerPoint presentation based on selected questions with custom styling.
    """
    try:
        # Validate that we have questions
        if not request.selected_questions:
            raise HTTPException(status_code=400, detail="No questions selected")
        
        # Access the global QA system and vectordb
        from app import vectordb, logger
        
        # Create QA system
        qa_system = CryptoQASystem(
            vectordb=vectordb,
            base_model="llama3"
        )
        
        # Define styling
        WHITE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)  # #FFFFFF - white
        QUESTION_FONT_SIZE = Pt(18)   # 18pt for questions
        ANSWER_FONT_SIZE = Pt(15)     # 10pt for answers
        TITLE_FONT_SIZE = Pt(32)      # 32pt for titles
        FONT_FAMILY = "Bodoni MT"     # Font family for all text
        
        # Image paths
        bg_image_path = "TemplateBG.jpg"
        logo_path = "LST.png"
        
        # Validate background image exists with absolute path
        bg_image_abs_path = os.path.abspath(bg_image_path)
        logo_abs_path = os.path.abspath(logo_path)
        
        if not os.path.exists(bg_image_abs_path):
            logger.error(f"BACKGROUND IMAGE NOT FOUND: '{bg_image_abs_path}'")
            raise HTTPException(status_code=500, detail=f"Background image not found at {bg_image_abs_path}")
        else:
            logger.info(f"Background image found at: {bg_image_abs_path}")
        
        # Process answers for each question
        section_answers = {}
        for question in request.selected_questions:
            logger.info(f"Processing question: {question.id} - {question.text}")

            # Add the coin name to the question
            coin_specific_question = f"Regarding {request.coin_name}, {question.text}"

            result = qa_system.answer_question(
                question=coin_specific_question,
                doc_id=request.doc_id
            )

            # Store the answer with the section
            if question.section_id not in section_answers:
                section_answers[question.section_id] = {
                    'title': question.section_title,
                    'questions': []
                }
            
            section_answers[question.section_id]['questions'].append({
                'question_id': question.id,
                'question_text': question.text,
                'answer': result.get('answer', "No answer available")
            })
            
        # Create a new presentation from scratch
        prs = Presentation()
        
        # Set slide dimensions
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Helper function to calculate centered position
        def get_centered_position(object_width):
            return (13.333 - object_width) / 2
        
        # 1. Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Add background
        bg = title_slide.shapes.add_picture(
            bg_image_abs_path,
            0, 0,
            width=prs.slide_width,
            height=prs.slide_height
        )
        # Move to back
        title_slide.shapes._spTree.remove(bg._element)
        title_slide.shapes._spTree.insert(0, bg._element)
        logger.info("Added background to title slide")

        # Logo configuration
        logo_width = 3  # Width in inches
        logo_height = 1.5  # Height proportional to width

        # Title configuration
        title_width = Inches(8)
        title_height = Inches(1.5)

        # Position the logo just above the center point
        center_y = prs.slide_height / 2
        logo_top = center_y - Inches(logo_height) - Inches(1.2)  # Logo positioned above center
        logo_left = Inches(get_centered_position(logo_width))

        # Position the title exactly at the center of the slide
        title_top = center_y - (title_height / 2)  # Center the title box vertically
        title_left = Inches(get_centered_position(8))  # Center the title box horizontally

        # Add logo first (so it appears above title)
        if os.path.exists(logo_abs_path):
            title_slide.shapes.add_picture(logo_abs_path, logo_left, logo_top, width=Inches(logo_width))
            logger.info("Added logo to title slide")

        # Add title in the exact center
        title_box = title_slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_paragraph = title_frame.add_paragraph()
        title_paragraph.text = f"{request.coin_name} Due Diligence Report"
        title_paragraph.font.color.rgb = WHITE_COLOR
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.font.name = FONT_FAMILY  # Set Bodoni MT font
        title_paragraph.alignment = PP_ALIGN.CENTER
        logger.info("Added title text to title slide")
                
        # 2. TOC slide
        toc_slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Add background
        bg = toc_slide.shapes.add_picture(
            bg_image_abs_path,
            0, 0,
            width=prs.slide_width,
            height=prs.slide_height
        )
        # Move to back
        toc_slide.shapes._spTree.remove(bg._element)
        toc_slide.shapes._spTree.insert(0, bg._element)
        logger.info("Added background to TOC slide")
        
        # Add TOC title
        toc_title_left = Inches(get_centered_position(8))
        toc_title_top = Inches(0.5)
        toc_title_width = Inches(8)
        toc_title_height = Inches(1)
        
        toc_title_box = toc_slide.shapes.add_textbox(toc_title_left, toc_title_top, toc_title_width, toc_title_height)
        toc_title_frame = toc_title_box.text_frame
        toc_title_paragraph = toc_title_frame.add_paragraph()
        toc_title_paragraph.text = "Table of Contents"
        toc_title_paragraph.font.color.rgb = WHITE_COLOR
        toc_title_paragraph.font.size = TITLE_FONT_SIZE
        toc_title_paragraph.font.bold = True
        toc_title_paragraph.font.name = FONT_FAMILY  # Set Bodoni MT font
        toc_title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add TOC content
        toc_content_left = Inches(get_centered_position(6))
        toc_content_top = Inches(1.8)
        toc_content_width = Inches(6)
        toc_content_height = Inches(5)
        
        toc_content_box = toc_slide.shapes.add_textbox(toc_content_left, toc_content_top, toc_content_width, toc_content_height)
        toc_content_frame = toc_content_box.text_frame
        
        for section_id, section_data in sorted(section_answers.items()):
            toc_paragraph = toc_content_frame.add_paragraph()
            toc_paragraph.text = f"{section_id}. {section_data['title']}"
            toc_paragraph.font.color.rgb = WHITE_COLOR
            toc_paragraph.font.size = Pt(18)
            toc_paragraph.font.name = FONT_FAMILY  # Set Bodoni MT font
            toc_paragraph.space_after = Pt(12)
        logger.info("Added TOC content")
        
        # Process sections and questions
        for section_id, section_data in sorted(section_answers.items()):
            # Create section slide with background
            section_slide = prs.slides.add_slide(prs.slide_layouts[5])
            # Add background
            bg = section_slide.shapes.add_picture(
                bg_image_abs_path,
                0, 0,
                width=prs.slide_width,
                height=prs.slide_height
            )
            # Move to back
            section_slide.shapes._spTree.remove(bg._element)
            section_slide.shapes._spTree.insert(0, bg._element)
            
            # Add section title
            section_title_left = Inches(get_centered_position(8))
            section_title_top = Inches(2)
            section_title_width = Inches(8)
            section_title_height = Inches(1.5)
            
            section_title_box = section_slide.shapes.add_textbox(section_title_left, section_title_top, section_title_width, section_title_height)
            section_title_frame = section_title_box.text_frame
            section_title_paragraph = section_title_frame.add_paragraph()
            section_title_paragraph.text = f"{section_id}. {section_data['title']}"
            section_title_paragraph.font.color.rgb = WHITE_COLOR
            section_title_paragraph.font.size = TITLE_FONT_SIZE
            section_title_paragraph.font.bold = True
            section_title_paragraph.font.name = FONT_FAMILY  # Set Bodoni MT font
            section_title_paragraph.alignment = PP_ALIGN.CENTER
            logger.info(f"Added section slide for section {section_id}")
            
            # Add slides for each Q&A
            for qa_item in section_data['questions']:
                # Create QA slide with background
                qa_slide = prs.slides.add_slide(prs.slide_layouts[5])
                # Add background
                bg = qa_slide.shapes.add_picture(
                    bg_image_abs_path,
                    0, 0,
                    width=prs.slide_width,
                    height=prs.slide_height
                )
                # Move to back
                qa_slide.shapes._spTree.remove(bg._element)
                qa_slide.shapes._spTree.insert(0, bg._element)
                
                # Add question
                question_left = Inches(0.5)
                question_top = Inches(0.5)
                question_width = Inches(12)
                question_height = Inches(1)
                
                question_box = qa_slide.shapes.add_textbox(question_left, question_top, question_width, question_height)
                question_frame = question_box.text_frame
                question_paragraph = question_frame.add_paragraph()
                question_paragraph.text = qa_item['question_text']
                question_paragraph.font.color.rgb = WHITE_COLOR
                question_paragraph.font.size = QUESTION_FONT_SIZE
                question_paragraph.font.bold = True
                question_paragraph.font.name = FONT_FAMILY  # Set Bodoni MT font
                
                # Add answer spread across multiple slides if needed
                # This replaces the original answer textbox creation code
                additional_slides = add_answer_slides(
                    prs,
                    qa_slide,
                    qa_item['answer'],
                    qa_item['question_text'],
                    bg_image_abs_path,
                    max_chars_per_slide=1400,  # Adjust based on font size and slide layout
                    font_family=FONT_FAMILY,
                    answer_font_size=ANSWER_FONT_SIZE
                )
                
                logger.info(f"Added QA slide for question: {qa_item['question_text'][:30]}... with {additional_slides} additional slides")
        
        # Create a unique filename
        filename = f"{request.coin_name.replace(' ', '_')}_due_diligence_{uuid.uuid4().hex[:8]}.pptx"
        
        # Save the presentation
        output_dir = "reports"
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, filename)
        
        # Save with extra validation
        logger.info(f"Saving presentation to {output_path}")
        try:
            prs.save(output_path)
            if os.path.exists(output_path):
                file_size = os.path.getsize(output_path)
                logger.info(f"Successfully saved: {output_path} ({file_size} bytes)")
            else:
                logger.error(f"File not found after saving: {output_path}")
        except Exception as e:
            logger.error(f"Error saving presentation: {e}")
            raise
        
        return PresentationResponse(
            file_path=output_path,
            message=f"Presentation generated with {len(request.selected_questions)} questions across {len(section_answers)} sections for {request.coin_name} using custom styling."
        )
    
    except Exception as e:
        logger.error(f"Error generating presentation: {str(e)}")
        import traceback
        traceback.print_exc()  # Print full traceback
        raise HTTPException(status_code=500, detail=f"Error generating presentation: {str(e)}")