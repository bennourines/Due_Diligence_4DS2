import os
import glob
import logging
import uuid
import time  
from typing import List
from fastapi import FastAPI, File, UploadFile , APIRouter , HTTPException ,Query
from fastapi.responses import FileResponse
from pydantic import BaseModel
import shutil
from langchain.document_loaders import PyPDFLoader, TextLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.chains import RetrievalQA
from langchain_community.llms import Ollama
from fastapi.middleware.cors import CORSMiddleware
from langchain_community.vectorstores import Qdrant
from qdrant_client import QdrantClient
from qdrant_client.http import models
from qdrant_client.models import Distance, VectorParams
from dotenv import load_dotenv
import requests
from enhanced_qa import CryptoQASystem
from generate_presentation import DueDiligenceReportGenerator, GenerateReportRequest, ReportResponse
from conversation_memory import ConversationManager, ConversationalCryptoQASystem
import pandas as pd
import numpy as np
import aiohttp
import asyncpraw
import requests
from datetime import datetime, timedelta
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.graph_objects as go
import plotly.express as px
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet
from textblob import TextBlob
from langchain_community.vectorstores import FAISS
from langchain_community.docstore.document import Document
from sentence_transformers import SentenceTransformer
from pytrends.request import TrendReq
import statsmodels.api as sm
from sklearn.preprocessing import MinMaxScaler
from tensorflow import keras
from tensorflow.keras.models import Sequential
from tensorflow.keras.layers import LSTM, Dense, Dropout
from sklearn.metrics import mean_absolute_error, mean_squared_error
import traceback
import asyncio
import matplotlib.pyplot as plt
from reportlab.platypus import Image
from pytrends.request import TrendReq
import time
import random
import plotly.graph_objects as go
import torch  # Import the torch module
# Add the necessary import statement for AutoTokenizer
from transformers import AutoTokenizer, AutoModelForSequenceClassification
import emoji
from torch.nn.functional import softmax
import ssl
import aiohttp
from huggingface_hub import login
login(token="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")  # Optional but recommended
import asyncio
import re
import pandas as pd
import aiohttp
import asyncpraw
import emoji
import torch
import torch.nn.functional as F
from transformers import AutoTokenizer, AutoModelForSequenceClassification
import os
from langchain.docstore.document import Document
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
import plotly.graph_objects as go
import matplotlib.pyplot as plt
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
import statsmodels.api as sm
from sklearn.preprocessing import MinMaxScaler
import numpy as np
from datetime import datetime, timedelta
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Use non-GUI backend for saving to PNG
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import logging
from typing import Optional , Dict ,Any 

from generate_presentation import GeneratePresentationRequest, PresentationResponse, generate_presentation


# Load environment variables from .env file
load_dotenv()

# Create a global conversation manager
conversation_manager = ConversationManager()


# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S"
)
logger = logging.getLogger(__name__)

# Get Serper API key from environment variables
SERPER_API_KEY = os.getenv("SERPER_API_KEY")
if not SERPER_API_KEY:
    logger.warning("SERPER_API_KEY not found in environment variables")


# Initialize embeddings
logger.info("Initializing embeddings...")
embeddings = HuggingFaceEmbeddings(
    model_name="sentence-transformers/all-MiniLM-L6-v2",
    model_kwargs={"device": "cpu"},
)

# Initialize Qdrant (local storage)
COLLECTION_NAME = "document_store"
VECTOR_SIZE = 384  # For all-MiniLM-L6-v2

# Create Qdrant client
logger.info("Connecting to Qdrant vector store...")
qdrant_path = "./qdrant_data"
os.makedirs(qdrant_path, exist_ok=True)

client = QdrantClient(path=qdrant_path)

# Check if collection exists, create if it doesn't
try:
    collection_info = client.get_collection(COLLECTION_NAME)
    logger.info(f"Using existing Qdrant collection '{COLLECTION_NAME}'")
except Exception:
    logger.info(f"Creating new Qdrant collection '{COLLECTION_NAME}'")
    client.create_collection(
        collection_name=COLLECTION_NAME,
        vectors_config=VectorParams(
            size=VECTOR_SIZE,
            distance=Distance.COSINE,
        ),
    )

# Initialize Qdrant wrapper with UUID generated IDs
vectordb = Qdrant(
    client=client,
    collection_name=COLLECTION_NAME,
    embeddings=embeddings,
)

# Shared text splitter for all documents
splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],  # Update with your Next.js frontend URL
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["Content-Disposition"],  # Important for downloads
)

app.add_route("/generate-presentation", generate_presentation)

class ProcessedFile(BaseModel):
    filename: str
    chunks_added: int

class IngestResponse(BaseModel):
    status: str  # "ok", "error", or "partial"
    message: Optional[str] = None
    chunks_added: int = 0
    processed_files: Optional[List[ProcessedFile]] = None
    errors: Optional[List[str]] = None
class QARequest(BaseModel):
    id: str
    question: str


# Define the request model
class QATestQuestion(BaseModel):
    question: str
    reference_answer: str
    doc_id: str = "all"

class QAPerformanceRequest(BaseModel):
    test_questions: List[QATestQuestion]
    threshold: float = 0.7  # Minimum acceptable F1 score

class SearchRequest(BaseModel):
    keyword: str
    num_results: int = 10  # Default number of results to return

class SearchResponse(BaseModel):
    status: str
    results: List[dict] = []
    message: str = ""

class GenerateQuestionsRequest(BaseModel):
    file_id: str  # The document ID to generate questions from
class Question(BaseModel):
    question: str


class GenerateQuestionsResponse(BaseModel):
    status: str  # "ok" or "error"
    message: Optional[str] = None
    questions: Optional[List[Question]] = None
    document_info: Optional[Dict[str, Any]] = None  # Added to store document metadata



# First, let's define the request and response models
class MultipleQARequest(BaseModel):
    questions: List[str]
    doc_id: str = "all"  # Default to "all" documents if not specified

class QuestionAnswer(BaseModel):
    question: str
    answer: str
    sources: List[dict] = []

class MultipleQAResponse(BaseModel):
    results: List[QuestionAnswer]
    status: str = "success"
    error: str = ""


@app.on_event("startup")
def load_cleaned_store():

    # Check if the function is disabled via environment variable
    disable_startup_load = os.getenv("DISABLE_STARTUP_LOAD", "false").lower() == "true"
    
    if disable_startup_load:
        logger.info("Startup data loading disabled by DISABLE_STARTUP_LOAD environment variable")
        return

    text_dir = "cleaned_text"  # Your folder name
    logger.info(f"Loading files from '{text_dir}'...")
    
    # Check if directory exists
    if not os.path.exists(text_dir):
        logger.warning(f"Directory '{text_dir}' does not exist. Skipping.")
        return
    
    # Get all files in the directory (supporting both .txt and .pdf)
    all_files = glob.glob(os.path.join(text_dir, "*.txt")) + glob.glob(os.path.join(text_dir, "*.pdf"))
    logger.info(f"Found {len(all_files)} files.")
    
    for file_path in all_files:
        doc_id = os.path.splitext(os.path.basename(file_path))[0]
        file_ext = os.path.splitext(file_path)[1].lower()
        
        logger.info(f"Processing file '{file_path}' as doc_id '{doc_id}'...")
        
        # Check if file is empty
        if os.path.getsize(file_path) == 0:
            logger.warning(f"File '{file_path}' is empty. Skipping.")
            continue
            
        try:
            # Load based on file extension
            if file_ext == '.pdf':
                loader = PyPDFLoader(file_path)
                file_type = "pdf"
            else:  # Default to text loader
                loader = TextLoader(file_path)
                file_type = "text"
                
            docs = loader.load()
            logger.info(f"Loaded {len(docs)} raw documents from '{file_path}'.")
            
            # Add document source metadata
            for doc in docs:
                doc.metadata["source"] = file_path
                doc.metadata["doc_id"] = doc_id
                doc.metadata["file_type"] = file_type
            
            chunks = splitter.split_documents(docs)
            logger.info(f"Split into {len(chunks)} chunks.")
            
            # Skip if no chunks were created
            if not chunks:
                logger.warning(f"No chunks were created from '{file_path}'. Skipping.")
                continue
            
            # Create a list of UUIDs for document IDs
            ids = [str(uuid.uuid4()) for _ in range(len(chunks))]
            
            # Add documents to Qdrant with UUID ids
            vectordb.add_documents(documents=chunks, ids=ids)
            logger.info(f"Added {len(chunks)} chunks to vector store from '{doc_id}'.")
            
        except Exception as e:
            logger.error(f"Error processing file '{file_path}': {str(e)}")
            continue
    
    logger.info(f"Completed loading all files from '{text_dir}' into Qdrant.")


@app.post("/ingest", response_model=IngestResponse)
async def ingest(files: List[UploadFile] = File(...)):
    """
    Ingest multiple files: load, split, and add to vector store.
    """
    logger.info(f"Received ingest request for {len(files)} files")
    
    total_chunks_added = 0
    processed_files = []
    errors = []
    
    try:
        # Create a temporary directory if it doesn't exist
        temp_dir = "temp_uploads"
        os.makedirs(temp_dir, exist_ok=True)
        
        for file in files:
            try:
                # Save the uploaded file to a temporary location
                file_path = os.path.join(temp_dir, file.filename)
                with open(file_path, "wb") as buffer:
                    shutil.copyfileobj(file.file, buffer)
                
                # Get file extension
                file_ext = os.path.splitext(file.filename)[1].lower()
                
                # Load based on file extension
                if file_ext == '.pdf':
                    loader = PyPDFLoader(file_path)
                    file_type = "pdf"
                else:  # Default to text loader
                    loader = TextLoader(file_path)
                    file_type = "text"
                    
                docs = loader.load()
                logger.info(f"Loaded {len(docs)} documents from '{file.filename}'.")
                
                # Add document source metadata
                for doc in docs:
                    doc.metadata["source"] = file.filename
                    doc.metadata["doc_id"] = os.path.splitext(file.filename)[0]
                    doc.metadata["file_type"] = file_type
                
                chunks = splitter.split_documents(docs)
                logger.info(f"Split into {len(chunks)} chunks.")
                
                # Generate UUID ids for each chunk
                ids = [str(uuid.uuid4()) for _ in range(len(chunks))]
                
                # Add to vector store
                vectordb.add_documents(documents=chunks, ids=ids)
                logger.info(f"Added {len(chunks)} chunks to vector store from '{file.filename}'.")
                
                total_chunks_added += len(chunks)
                processed_files.append({
                    "filename": file.filename,
                    "chunks_added": len(chunks)
                })
                
                # Clean up the temporary file
                os.remove(file_path)
                
            except Exception as e:
                error_msg = f"Error processing file '{file.filename}': {str(e)}"
                logger.error(error_msg)
                errors.append(error_msg)
                
                # Clean up the temporary file if it exists
                if 'file_path' in locals() and os.path.exists(file_path):
                    os.remove(file_path)
                    
            finally:
                # Ensure file is closed
                await file.close()
                
        # Determine overall status
        if len(errors) == 0:
            status = "ok"
        elif len(errors) < len(files):
            status = "partial"
        else:
            status = "error"
            
        return IngestResponse(
            status=status,
            chunks_added=total_chunks_added,
            processed_files=processed_files,
            errors=errors if errors else None
        )
        
    except Exception as e:
        logger.error(f"Error during batch ingest: {str(e)}")
        return IngestResponse(status="error", message=str(e))

@app.post("/qa")
def qa(req: QARequest):
    """
    Retrieve and answer questions about cryptocurrency using enhanced QA with reasoning capabilities.
    """
    logger.info(f"Received crypto QA request: id={req.id}, question='{req.question}'")
    try:
        # Create enhanced crypto QA system
        qa_system = CryptoQASystem(
            vectordb=vectordb,
            base_model="llama3"
        )
        
        # Get enhanced answer
        result = qa_system.answer_question(
            question=req.question,
            doc_id=req.id
        )
        
        return result
    except Exception as e:
        logger.error(f"Error in crypto QA endpoint: {str(e)}")
        return {"answer": f"Error: {str(e)}"}

@app.get("/status")
def get_status():
    """
    Get information about the vector store status.
    """
    try:
        collection_info = client.get_collection(COLLECTION_NAME)
        points_count = collection_info.points_count
        
        return {
            "status": "ok",
            "collection_name": COLLECTION_NAME,
            "documents_count": points_count,
            "vector_size": VECTOR_SIZE
        }
    except Exception as e:
        logger.error(f"Error getting status: {str(e)}")
        return {"status": "error", "message": str(e)}

@app.delete("/clear")
def clear_collection():
    """
    Clear all documents from the collection.
    """
    try:
        client.delete_collection(collection_name=COLLECTION_NAME)
        logger.info(f"Collection '{COLLECTION_NAME}' deleted.")
        
        # Recreate the collection
        client.create_collection(
            collection_name=COLLECTION_NAME,
            vectors_config=VectorParams(
                size=VECTOR_SIZE,
                distance=Distance.COSINE,
            ),
        )
        logger.info(f"Collection '{COLLECTION_NAME}' recreated.")
        
        return {"status": "ok", "message": f"Collection '{COLLECTION_NAME}' cleared and recreated."}
    except Exception as e:
        logger.error(f"Error clearing collection: {str(e)}")
        return {"status": "error", "message": str(e)}

@app.post("/search", response_model=SearchResponse)
def search(req: SearchRequest):
    """
    Perform a web search using Serper API for the given keyword.
    """
    logger.info(f"Received search request for keyword: {req.keyword}")
    
    if not SERPER_API_KEY:
        return SearchResponse(
            status="error",
            message="Serper API key not configured"
        )
    
    try:
        # Prepare the request to Serper API
        url = "https://google.serper.dev/search"
        headers = {
            'X-API-KEY': SERPER_API_KEY,
            'Content-Type': 'application/json'
        }
        payload = {
            'q': req.keyword,
            'num': req.num_results
        }
        
        # Make the request to Serper API
        response = requests.post(url, headers=headers, json=payload)
        response.raise_for_status()  # Raise exception for non-200 status codes
        
        # Process the response
        search_results = response.json()
        
        # Extract organic search results
        organic_results = search_results.get('organic', [])
        processed_results = [{
            'title': result.get('title', ''),
            'link': result.get('link', ''),
            'snippet': result.get('snippet', '')
        } for result in organic_results[:req.num_results]]
        
        logger.info(f"Successfully retrieved {len(processed_results)} search results")
        
        return SearchResponse(
            status="ok",
            results=processed_results
        )
        
    except requests.exceptions.RequestException as e:
        error_msg = f"Error making request to Serper API: {str(e)}"
        logger.error(error_msg)
        return SearchResponse(
            status="error",
            message=error_msg
        )
    except Exception as e:
        error_msg = f"Unexpected error during search: {str(e)}"
        logger.error(error_msg)
        return SearchResponse(
            status="error",
            message=error_msg
        )


@app.post("/qa/performance")
def qa_performance_metrics(req: QAPerformanceRequest):
    """
    Test performance metrics and scores for the QA system against reference answers.
    Returns precision, recall, F1 score, latency, and confidence metrics.
    """
    logger.info(f"Running performance test on QA system with {len(req.test_questions)} questions")
    
    try:
        # Initialize metrics
        results = {
            "overall": {
                "precision": 0.0,
                "recall": 0.0,
                "f1_score": 0.0,
                "latency_ms": 0.0,
                "confidence": 0.0
            },
            "questions": []
        }
        
        # Initialize LLM for evaluation
        try:
            eval_llm = Ollama(model="llama3", temperature=0.0)
            logger.info("Initialized evaluation LLM")
        except Exception as e:
            logger.error(f"Error initializing evaluation LLM: {str(e)}")
            return {"error": "Could not initialize evaluation model"}
        
        total_latency = 0
        
        # Test each question
        for idx, test_item in enumerate(req.test_questions):
            question = test_item.question
            reference_answer = test_item.reference_answer
            doc_id = test_item.doc_id if hasattr(test_item, "doc_id") else "all"
            
            logger.info(f"Testing question {idx+1}/{len(req.test_questions)}: '{question}'")
            
            # Measure latency
            start_time = time.time()
            
            # Create a retriever with optional filtering
            search_filter = None
            if doc_id != "all" and doc_id != "":
                search_filter = models.Filter(
                    must=[
                        models.FieldCondition(
                            key="metadata.doc_id",
                            match=models.MatchValue(value=doc_id)
                        )
                    ]
                )
                
            retriever = vectordb.as_retriever(
                search_type="similarity",
                search_kwargs={
                    "k": 4,
                    "filter": search_filter
                }
            )
            
            # Create and run the chain
            chain = RetrievalQA.from_chain_type(
                llm=Ollama(model="llama3", temperature=0.1),
                chain_type="stuff",
                retriever=retriever,
                return_source_documents=True
            )
            
            result = chain({"query": question})
            answer = result["result"]
            
            # Calculate latency
            end_time = time.time()
            latency_ms = (end_time - start_time) * 1000
            total_latency += latency_ms
            
            # Use evaluation LLM to calculate metrics
            eval_prompt = f"""
            You are an expert at evaluating question answering systems. 
            Compare the system's answer to the reference answer and score on a scale of 0 to 1:
            
            Question: {question}
            
            Reference Answer: {reference_answer}
            
            System Answer: {answer}
            
            Please provide numerical scores for:
            1. Precision (How accurate is the provided information? 0-1)
            2. Recall (How complete is the answer compared to reference? 0-1)
            3. Confidence (How certain can we be about this answer? 0-1)
            
            Also calculate F1 score as the harmonic mean of precision and recall.
            Return your evaluation as a JSON object with these metrics as float values.
            """
            
            try:
                eval_result = eval_llm.invoke(eval_prompt)
                
                # Parse the JSON metrics from the response
                # Find JSON content between ```json and ``` if present
                import re
                import json
                
                json_pattern = r'```json\s*([\s\S]*?)\s*```'
                json_match = re.search(json_pattern, eval_result)
                
                if json_match:
                    metrics_str = json_match.group(1)
                else:
                    # Try to find any JSON-like content
                    json_pattern = r'\{[\s\S]*?\}'
                    json_match = re.search(json_pattern, eval_result)
                    metrics_str = json_match.group(0) if json_match else "{}"
                
                try:
                    metrics = json.loads(metrics_str)
                except:
                    # Fallback if JSON parsing fails
                    logger.warning("JSON parsing failed, using default metrics")
                    metrics = {
                        "precision": 0.5,
                        "recall": 0.5,
                        "f1_score": 0.5,
                        "confidence": 0.5
                    }
                    
                # Calculate F1 if not provided
                if "f1_score" not in metrics and "precision" in metrics and "recall" in metrics:
                    p = metrics["precision"]
                    r = metrics["recall"]
                    metrics["f1_score"] = 2 * (p * r) / (p + r) if (p + r) > 0 else 0
            
            except Exception as e:
                logger.error(f"Error in evaluation: {str(e)}")
                metrics = {
                    "precision": 0.5,
                    "recall": 0.5,
                    "f1_score": 0.5,
                    "confidence": 0.5
                }
            
            # Extract sources
            sources = []
            for doc in result.get("source_documents", []):
                source_info = {
                    "source": doc.metadata.get("source", "Unknown"),
                    "doc_id": doc.metadata.get("doc_id", "Unknown"),
                    "page": doc.metadata.get("page", 0),
                    "file_type": doc.metadata.get("file_type", "Unknown")
                }
                sources.append(source_info)
            
            # Add question results
            question_result = {
                "question": question,
                "reference_answer": reference_answer,
                "system_answer": answer,
                "metrics": metrics,
                "latency_ms": latency_ms,
                "sources": sources
            }
            
            results["questions"].append(question_result)
            
            logger.info(f"Question {idx+1} evaluated - F1: {metrics.get('f1_score', 0):.2f}, Latency: {latency_ms:.2f}ms")
        
        # Calculate overall metrics
        if req.test_questions:
            avg_precision = sum(q["metrics"].get("precision", 0) for q in results["questions"]) / len(req.test_questions)
            avg_recall = sum(q["metrics"].get("recall", 0) for q in results["questions"]) / len(req.test_questions)
            avg_f1 = sum(q["metrics"].get("f1_score", 0) for q in results["questions"]) / len(req.test_questions)
            avg_confidence = sum(q["metrics"].get("confidence", 0) for q in results["questions"]) / len(req.test_questions)
            avg_latency = total_latency / len(req.test_questions)
            
            results["overall"] = {
                "precision": avg_precision,
                "recall": avg_recall,
                "f1_score": avg_f1,
                "latency_ms": avg_latency,
                "confidence": avg_confidence,
                "questions_evaluated": len(req.test_questions)
            }
        
        logger.info(f"Performance test completed. Overall F1: {results['overall']['f1_score']:.2f}")
        return results
    
    except Exception as e:
        logger.error(f"Error in QA performance endpoint: {str(e)}")
        return {"error": f"Error in performance testing: {str(e)}"}

@app.post("/generate-questions", response_model=GenerateQuestionsResponse)
async def generate_questions(req: GenerateQuestionsRequest):
    """
    Generate 5 questions from the specified document.
    """
    logger.info(f"Received request to generate questions for document: {req.file_id}")
    
    try:
        # Ensure file_id is provided
        if not req.file_id:
            return GenerateQuestionsResponse(
                status="error",
                message="Error: file_id is required"
            )
        
        # Create a search filter for the specific document
        search_filter = models.Filter(
            must=[
                models.FieldCondition(
                    key="metadata.doc_id",
                    match=models.MatchValue(value=req.file_id)
                )
            ]
        )
        
        # Get relevant document chunks from the vector store
        retriever = vectordb.as_retriever(
            search_type="similarity",
            search_kwargs={
                "k": 10,  # Retrieve more chunks to get comprehensive content
                "filter": search_filter
            }
        )
        
        # Get document chunks
        query_result = retriever.get_relevant_documents("What is this document about?")
        
        if not query_result:
            return GenerateQuestionsResponse(
                status="error",
                message=f"No content found for document ID: {req.file_id}"
            )
        
        # Combine document content for context
        document_content = "\n\n".join([doc.page_content for doc in query_result])
        
        # Truncate content if it's too large
        if len(document_content) > 12000:
            document_content = document_content[:12000] + "..."
            logger.info(f"Document content truncated to 12000 characters")
        
        # Initialize Ollama LLM
        try:
            llm = Ollama(model="llama3", temperature=0.7)  # Higher temperature for creative questions
            logger.info("Initialized Ollama LLM for question generation")
        except Exception as e:
            logger.error(f"Error initializing Ollama: {str(e)}")
            return GenerateQuestionsResponse(
                status="error",
                message="Error: Could not initialize language model. Please check if Ollama is running."
            )
        
        # Create prompt for generating questions
        prompt = f"""
        Based on the following document content, generate 5 diverse and specific questions that would test someone's understanding of the material.
        
        The questions should:
        1. Be answerable from the provided content
        2. Cover different aspects of the document
        3. Include both factual and analytical questions
        4. Be clear and specific
        5. Vary in difficulty level
        
        Document content:
        {document_content}
        
        Return ONLY the 5 questions, each on a new line starting with 'Question: '.
        """
        
        # Generate questions
        logger.info("Generating questions from document content...")
        result = llm.invoke(prompt)
        
        # Process the result to extract questions
        questions = []
        
        # Parse questions from the response
        for line in result.split('\n'):
            line = line.strip()
            if line.startswith('Question:') or line.startswith('Q:') or line.startswith(f'1.') or line.startswith(f'- '):
                # Clean up the question format
                question_text = line
                for prefix in ['Question:', 'Q:', '1.', '2.', '3.', '4.', '5.', '- ']:
                    if question_text.startswith(prefix):
                        question_text = question_text[len(prefix):].strip()
                        break
                
                questions.append(Question(question=question_text))
        
        # Handle case where parsing might not have worked
        if not questions:
            # Try another parsing approach - just split by numbers
            import re
            q_matches = re.findall(r'(?:\d+[\)\.]\s*|\bQuestion\s*\d*[\:\.]\s*)(.*?)(?=\d+[\)\.]\s*|\bQuestion\s*\d*[\:\.]\s*|$)', result, re.DOTALL)
            if q_matches:
                questions = [Question(question=q.strip()) for q in q_matches if q.strip()]
        
        # If we still don't have enough questions, just get the first 5 sentences
        if len(questions) < 5:
            # Get first 5 sentences that end with question marks
            import re
            q_sentences = re.findall(r'[^.!?]*\?', result)
            for q in q_sentences:
                if len(questions) >= 5:
                    break
                if q.strip() and len(q.strip()) > 10:  # Minimum question length
                    questions.append(Question(question=q.strip()))
        
        # Limit to 5 questions
        questions = questions[:5]
        
        # Add document metadata to response
        doc_metadata = {}
        if query_result and len(query_result) > 0:
            # Extract filename from the first document's metadata
            doc_metadata = {
                "filename": query_result[0].metadata.get("source", "Unknown"),
                "doc_id": req.file_id,
                "doc_type": query_result[0].metadata.get("file_type", "Unknown")
            }
        
        # Return the questions with document info
        logger.info(f"Generated {len(questions)} questions for document ID: {req.file_id}")
        return GenerateQuestionsResponse(
            status="ok",
            questions=questions,
            document_info=doc_metadata
        )
        
    except Exception as e:
        logger.error(f"Error generating questions: {str(e)}")
        return GenerateQuestionsResponse(
            status="error",
            message=f"Error generating questions: {str(e)}"
        )
    
    
presentation_router = APIRouter()

@presentation_router.post("/generate-presentation", response_model=PresentationResponse)
async def api_generate_presentation(request: GeneratePresentationRequest):
    return await generate_presentation(request)

# Include the router in your app
app.include_router(presentation_router)

@app.post("/multiple-qa", response_model=MultipleQAResponse)
def multiple_qa(req: MultipleQARequest):
    """
    Process multiple cryptocurrency questions at once using enhanced QA with reasoning capabilities.
    """
    logger.info(f"Received multiple QA request: doc_id={req.doc_id}, questions count={len(req.questions)}")
    
    try:
        # Create enhanced crypto QA system
        qa_system = CryptoQASystem(
            vectordb=vectordb,
            base_model="llama3"
        )
        
        results = []
        
        # Process each question
        for question in req.questions:
            logger.info(f"Processing question: '{question}'")
            
            # Get enhanced answer for the current question
            result = qa_system.answer_question(
                question=question,
                doc_id=req.doc_id
            )
            
            # Add to results
            qa_item = QuestionAnswer(
                question=question,
                answer=result.get("answer", "Error processing question"),
                sources=result.get("sources", [])
            )
            results.append(qa_item)
        
        return MultipleQAResponse(results=results)
        
    except Exception as e:
        logger.error(f"Error in multiple QA endpoint: {str(e)}")
        return MultipleQAResponse(
            results=[],
            status="error",
            error=f"Error processing multiple questions: {str(e)}"
        )

# Add a health check endpoint
@app.get("/health")
async def health_check():
    """Health check endpoint to verify the server is running"""
    return {"status": "ok", "message": "Server is running"}
    

# Add a route to download the presentation
@app.get("/download-pres/{filename}")
async def download_presentation(filename: str):
    """
    Download a generated presentation.
    """
    try:
        # Log the request for debugging
        print(f"Received download request for file: {filename}")
        
        file_path = f"reports/{filename}"
        print(f"Looking for file at path: {file_path}")
        
        # Check if file exists
        if not os.path.exists(file_path):
            print(f"File not found: {file_path}")
            
            # List available files in the directory for debugging
            print("Available files in reports directory:")
            if os.path.exists("reports"):
                files = os.listdir("reports")
                for f in files:
                    print(f"  - {f}")
            else:
                print("Reports directory does not exist")
            
            raise HTTPException(
                status_code=404,
                detail=f"Report {filename} not found"
            )
        
        # Log file details
        file_size = os.path.getsize(file_path)
        print(f"File found: {filename}, Size: {file_size} bytes")
        
        # Return the file
        return FileResponse(
            path=file_path,
            filename=filename,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        print(f"Error in download_presentation: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Error processing download: {str(e)}"
        )
    

class CryptoAnalyzer:
    def __init__(self, api_keys):
        self.data_folder = r"Data"
        self.api_keys = api_keys
        self.data = None
        self.sentiment_data = None
        self.device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        self.sentiment_tokenizer = AutoTokenizer.from_pretrained("cardiffnlp/twitter-roberta-base-sentiment")
        self.sentiment_model = AutoModelForSequenceClassification.from_pretrained("cardiffnlp/twitter-roberta-base-sentiment").to(self.device)
        self.coin_csv_mapping = self._build_coin_mapping()
        self.coin_name = None
        self.coin_symbol = None
        self.technical_indicators = None
        self.event_summary_df = None
        self.forecasts = {'EMA': None, 'ARIMA': None, 'LSTM': None, 'Combined': None}

        # Initialize logger
        self.logger = logging.getLogger(__name__)
        if not self.logger.handlers:
            self.logger.setLevel(logging.INFO)
            handler = logging.StreamHandler()
            formatter = logging.Formatter('%(levelname)s: %(message)s')
            handler.setFormatter(formatter)
            self.logger.addHandler(handler)

    def _build_coin_mapping(self):
        mapping = {}
        if not os.path.exists(self.data_folder):
            print("[WARNING] Data folder {self.data_folder} does not exist.")
            return mapping
        for filename in os.listdir(self.data_folder):
            if filename.endswith(".csv"):
                file_path = os.path.join(self.data_folder, filename)
                try:
                    df = pd.read_csv(file_path, nrows=1)
                    coin_name = df['Coin Name'].iloc[0].lower()
                    coin_symbol = df['Coin Symbol'].iloc[0].lower()
                    mapping[coin_name] = file_path
                    mapping[coin_symbol] = file_path
                except Exception as e:
                    print("[WARNING] Could not read {filename}: {e}")
        print(f"[INFO] Loaded mapping for {len(mapping)//2} coins.")
        return mapping

    def _detect_coin_from_question(self, question):
        question = question.lower()
        for coin in self.coin_csv_mapping:
            if coin in question:
                return coin
        return None

    def load_data_from_csv(self, file_path):
        try:
            print(f"[INFO] Loading data from {file_path}")
            self.data = pd.read_csv(file_path)
            expected_columns = ['Coin Name', 'Coin Symbol', 'Date', 'Open', 'High', 'Low', 'Close', 'Volume', 'Market Cap']
            if not all(col in self.data.columns for col in expected_columns):
                raise ValueError("CSV must contain: {expected_columns}")

            self.coin_name = self.data['Coin Name'].iloc[0].lower()
            self.coin_symbol = self.data['Coin Symbol'].iloc[0].lower()
            self.data['Date'] = pd.to_datetime(self.data['Date'], errors='coerce')
            self.data = self.data.sort_values('Date')
            self.data.set_index('Date', inplace=True)

            numeric_columns = ['Open', 'High', 'Low', 'Close', 'Volume', 'Market Cap']
            for col in numeric_columns:
                self.data[col] = pd.to_numeric(self.data[col].replace('[\$,]', '', regex=True), errors='coerce')

            if self.data[numeric_columns].isna().any().any():
                print(f"[WARNING] Non-numeric values found in numeric columns. Dropping invalid rows.")
                self.data = self.data.dropna(subset=numeric_columns)

            if self.data.empty:
                print(f"[ERROR] No valid data remains after cleaning for {file_path}")
                self.data = None
                return None

            print(f"[INFO] Loaded data for {self.coin_name} ({self.coin_symbol})")
            print(f"[INFO] Date range: {self.data.index.min()} to {self.data.index.max()}")
            print(f"[INFO] Total entries: {len(self.data)}")

            self.data['Return'] = self.data['Close'].pct_change()
            self.data.dropna(inplace=True)
            return self.data
        except Exception as e:
            print(f"[ERROR] Error loading CSV: {e}")
            self.data = None
            return None

    async def analyze_text_sentiment(self, text):
        if not text or not isinstance(text, str) or text.strip() == "":
            return 0.0

        try:
            text = emoji.demojize(text, delimiters=(" :", ": "))
            text = text.lower().strip()
            text = re.sub(r'[^a-zA-Z\s]', ' ', text)
            text = re.sub(r'\s+', ' ', text)
            print(f"Processing text: '{text[:50]}...'")

            encoded_input = self.sentiment_tokenizer(
                text,
                return_tensors='pt',
                truncation=True,
                max_length=512,
                padding='max_length'
            ).to(self.device)

            with torch.no_grad():
                output = self.sentiment_model(**encoded_input)

            scores = F.softmax(output.logits, dim=1).cpu().numpy()[0]
            compound = scores[2] - scores[0]
            print(f"Sentiment score for '{text[:50]}...': {compound}")

            del encoded_input, output
            if torch.cuda.is_available():
                torch.cuda.empty_cache()

            return round(float(compound), 4)

        except Exception as e:
            print(f"[ERROR] Error analyzing sentiment: {e}")
            return 0.0

    async def fetch_reddit_data(self, session, start_date, end_date):
        try:
            print(f"[INFO] Fetching Reddit data for {self.coin_name or 'unknown coin'}...")

            reddit = asyncpraw.Reddit(
                client_id=self.api_keys['reddit_client_id'],
                client_secret=self.api_keys['reddit_client_secret'],
                user_agent=self.api_keys['reddit_user_agent']
            )

            posts = []
            subreddits = ['cryptocurrency', 'crypto', 'cryptomarkets', 'cryptotechnology']
            query = f"{self.coin_symbol.upper()} OR {self.coin_name.lower()} OR #{self.coin_name.lower()} OR cryptocurrency"
            MAX_POSTS = 500

            for subreddit_name in subreddits:
                subreddit = await reddit.subreddit(subreddit_name)
                print(f"[INFO] Searching subreddit: {subreddit_name}")
                async for submission in subreddit.search(query, time_filter='all', limit=MAX_POSTS):
                    post_date = pd.to_datetime(submission.created_utc, unit='s').tz_localize(None)
                    if start_date <= post_date <= end_date:
                        text = f"{submission.title} {submission.selftext}".strip()
                        if text:
                            posts.append({
                                'date': post_date,
                                'text': text
                            })
                    await asyncio.sleep(0.5)

            await reddit.close()

            if not posts:
                print("[INFO] No Reddit posts found in date range")
                return pd.DataFrame(columns=['date', 'text'])

            df = pd.DataFrame(posts)
            df['date'] = pd.to_datetime(df['date']).dt.tz_localize(None)
            print(f"[INFO] Fetched {len(df)} Reddit posts")
            return df

        except Exception as e:
            print(f"[ERROR] Error fetching Reddit data: {e}")
            return pd.DataFrame(columns=['date', 'text'])

    async def fetch_cryptocompare_data(self, session, start_date, end_date):
        try:
            print("[INFO] Fetching CryptoCompare news data...")
            api_key = self.api_keys.get('cryptocompare_api_key')

            if not api_key:
                print("[ERROR] No CryptoCompare API key found")
                return pd.DataFrame(columns=['date', 'title', 'body'])

            url = "https://min-api.cryptocompare.com/data/v2/news/"
            params = {
                'api_key': api_key,
                'categories': self.coin_symbol.upper(),
                'lTs': int((start_date - pd.Timedelta(days=60)).timestamp()),
                'toTs': int(end_date.timestamp()),
                'limit': 200
            }

            async with session.get(url, params=params, timeout=15) as resp:
                if resp.status != 200:
                    print(f"[ERROR] CryptoCompare news fetch status {resp.status}")
                    raise Exception(f"HTTP {resp.status}")

                data = await resp.json()
                items = data.get('Data', [])

            rows = []
            for item in items:
                try:
                    published_ts = item.get('published_on')
                    if not isinstance(published_ts, (int, float)) or published_ts <= 0:
                        continue

                    dt = pd.to_datetime(published_ts, unit='s', utc=True).tz_convert(None)
                    if start_date <= dt <= end_date:
                        title = str(item.get('title', '') or '')
                        body = str(item.get('body', '') or '')
                        if title or body:
                            rows.append({
                                'date': dt,
                                'title': title,
                                'body': body
                            })
                except Exception as e:
                    print(f"[WARNING] Failed to process CryptoCompare news item {item.get('id')}: {e}")
                    continue

            if rows:
                df = pd.DataFrame(rows)
                df['date'] = pd.to_datetime(df['date'])
                df = df.set_index('date').sort_index()
                print(f"[INFO] CryptoCompare fetched {len(df)} news items")
                return df

            print("[INFO] No news data found, falling back to social metrics...")
            url = f"https://min-api.cryptocompare.com/data/social/coin/latest?symbol={self.coin_symbol.upper()}"
            headers = {"authorization": f"Apikey {api_key}"}

            async with session.get(url, headers=headers, timeout=15) as resp:
                if resp.status != 200:
                    print(f"[ERROR] CryptoCompare social fetch status {resp.status}")
                    raise Exception(f"HTTP {resp.status}")

                data = await resp.json()
                if 'Data' not in data or 'General' not in data['Data']:
                    print("[ERROR] Invalid CryptoCompare social response format")
                    return pd.DataFrame(columns=['date', 'text'])

            reddit = data['Data'].get('Reddit', {})
            twitter = data['Data'].get('Twitter', {})
            sentiment_text = (
                f"Reddit subscribers: {reddit.get('subscribers', 0)}, posts/hour: {reddit.get('posts_per_hour', 0)}. "
                f"Twitter followers: {twitter.get('followers', 0)}, tweets/hour: {twitter.get('status_per_hour', 0) or 0}."
            )

            sentiment_score = await self.analyze_text_sentiment(sentiment_text)
            sentiment_dates = pd.date_range(start=start_date, end=end_date, freq='D')
            df = pd.DataFrame({
                'date': sentiment_dates,
                'text': [sentiment_text] * len(sentiment_dates),
                'sentiment': [sentiment_score] * len(sentiment_dates)
            })
            df['date'] = pd.to_datetime(df['date'])
            df = df.set_index('date').sort_index()
            print(f"[WARNING] Using CryptoCompare social data snapshot for {len(df)} days")
            return df

        except Exception as e:
            print(f"[ERROR] CryptoCompare fetch failed: {e}")
            return pd.DataFrame(columns=['date', 'title', 'body'])

    async def fetch_cryptopanic_data(self, session, start_date, end_date):
        try:
            print("[INFO] Fetching CryptoPanic news data...")
            api_key = self.api_keys.get('cryptopanic_api_key')
            if not api_key:
                print("[ERROR] No CryptoPanic API key found")
                return pd.DataFrame(columns=['date', 'title'])

            start_str = start_date.strftime('%Y-%m-%d')
            end_str = end_date.strftime('%Y-%m-%d')
            base_url = "https://cryptopanic.com/api/v1/posts/"
            params = {
                'auth_token': api_key,
                'currencies': self.coin_symbol.upper(),
                'kind': 'news',
                'from': start_str,
                'to': end_str,
                'public': 'true'
            }

            async with session.get(base_url, params=params, timeout=15) as resp:
                if resp.status != 200:
                    print(f"[ERROR] CryptoPanic fetch status {resp.status}: {await resp.text()}")
                    raise Exception(f"HTTP {resp.status}")

                data = await resp.json()
                items = data.get('results', [])

            rows = []
            for item in items:
                try:
                    published_dt = pd.to_datetime(item.get('published_at')).tz_localize(None)
                    if start_date <= published_dt <= end_date:
                        title = str(item.get('title', '') or '')
                        if title:
                            rows.append({
                                'date': published_dt,
                                'title': title
                            })
                except Exception as e:
                    print(f"[WARNING] Failed to process CryptoPanic news item: {e}")
                    continue

            if rows:
                df = pd.DataFrame(rows)
                df['date'] = pd.to_datetime(df['date'])
                df = df.set_index('date').sort_index()
                print(f"[INFO] CryptoPanic fetched {len(df)} news items")
                return df

            print("[INFO] No CryptoPanic news data found, defaulting to 0.0 sentiment")
            sentiment_dates = pd.date_range(start=start_date, end=end_date, freq='D')
            df = pd.DataFrame({'date': sentiment_dates, 'title': ['No data'] * len(sentiment_dates)})
            df['date'] = pd.to_datetime(df['date'])
            df = df.set_index('date').sort_index()
            return df

        except Exception as e:
            print(f"[ERROR] CryptoPanic fetch failed: {e}")
            return pd.DataFrame(columns=['date', 'title'])

    async def analyze_sentiment(self, start_date, end_date, session):
        try:
            start_date = pd.to_datetime(start_date, errors='coerce')
            end_date = pd.to_datetime(end_date, errors='coerce')
            if pd.isna(start_date) or pd.isna(end_date):
                print(f"[ERROR] Invalid date format: start_date={start_date}, end_date={end_date}")
                return None
            if start_date > end_date:
                print(f"[ERROR] start_date ({start_date}) must be before end_date ({end_date})")
                return None
            if not self.coin_symbol or not self.coin_name:
                print("[ERROR] Coin name or symbol not set")
                return None

            if self.data is not None and not self.data.empty:
                max_data_date = self.data.index.max()
                if end_date > max_data_date:
                    print(f"[INFO] Adjusting end_date from {end_date} to {max_data_date} to match available data")
                    end_date = max_data_date

            sentiment_dates = pd.date_range(start=start_date, end=end_date, freq='D')
            sentiment = pd.DataFrame(
                index=sentiment_dates,
                columns=['Reddit_Sentiment', 'CryptoPanic_Sentiment', 'CryptoCompare_Sentiment', 'Combined_Sentiment'],
                dtype=float
            )
            sentiment.index.name = 'Date'

            reddit_df = await self.fetch_reddit_data(session, start_date, end_date)
            if not reddit_df.empty:
                texts = [text for text in reddit_df['text'] if isinstance(text, str) and text.strip()]
                if texts:
                    sentiments = await asyncio.gather(*[self.analyze_text_sentiment(text) for text in texts])
                    reddit_df['sentiment'] = sentiments
                    reddit_daily = reddit_df.groupby(reddit_df['date'].dt.date)['sentiment'].mean()
                    reddit_daily.index = pd.to_datetime(reddit_daily.index)
                    sentiment['Reddit_Sentiment'] = reddit_daily.reindex(sentiment_dates, fill_value=0.0)
                    print(f"[INFO] Reddit sentiment processed: {len(texts)} posts")

            cryptopanic_df = await self.fetch_cryptopanic_data(session, start_date, end_date)
            if not cryptopanic_df.empty:
                texts = [title for title in cryptopanic_df['title'] if isinstance(title, str) and title.strip()]
                if texts:
                    sentiments = await asyncio.gather(*[self.analyze_text_sentiment(text) for text in texts])
                    cryptopanic_df['sentiment'] = sentiments
                    cryptopanic_daily = cryptopanic_df['sentiment'].groupby(cryptopanic_df.index.date).mean()
                    cryptopanic_daily.index = pd.to_datetime(cryptopanic_daily.index)
                    sentiment['CryptoPanic_Sentiment'] = cryptopanic_daily.reindex(sentiment_dates, fill_value=0.0)
                    print(f"[INFO] CryptoPanic sentiment processed: {len(texts)} news items")
                else:
                    sentiment['CryptoPanic_Sentiment'] = 0.0
            else:
                sentiment['CryptoPanic_Sentiment'] = 0.0

            cryptocompare_df = await self.fetch_cryptocompare_data(session, start_date, end_date)
            if not cryptocompare_df.empty:
                if 'title' in cryptocompare_df.columns:
                    texts = [f"{t} {b}".strip() for t, b in zip(cryptocompare_df['title'], cryptocompare_df['body']) if isinstance(t, str) and isinstance(b, str)]
                    if texts:
                        sentiments = await asyncio.gather(*[self.analyze_text_sentiment(text) for text in texts])
                        cryptocompare_df['sentiment'] = sentiments
                        cryptocompare_daily = cryptocompare_df['sentiment'].groupby(cryptocompare_df.index.date).mean()
                        cryptocompare_daily.index = pd.to_datetime(cryptocompare_daily.index)
                        sentiment['CryptoCompare_Sentiment'] = cryptocompare_daily.reindex(sentiment_dates, fill_value=0.0)
                        print(f"[INFO] CryptoCompare sentiment processed: {len(texts)} news items")
                else:
                    sentiment['CryptoCompare_Sentiment'] = cryptocompare_df['sentiment'].reindex(sentiment_dates, fill_value=0.0)
                    print(f"[INFO] CryptoCompare social sentiment processed: {len(cryptocompare_df)} days")

            for col in ['Reddit_Sentiment', 'CryptoPanic_Sentiment', 'CryptoCompare_Sentiment']:
                sentiment[col] = sentiment[col].ffill()

            weights = {'Reddit': 0.4, 'CryptoPanic': 0.3, 'CryptoCompare': 0.3}
            valid_sources = []
            reddit_count = len(reddit_df) if not reddit_df.empty else 0
            cryptopanic_count = len(cryptopanic_df) if not cryptopanic_df.empty else 0
            if reddit_count > 0:
                weights['Reddit'] = min(0.6, 0.4 + 0.05 * (reddit_count / 10))
            if cryptopanic_count > 0:
                weights['CryptoPanic'] = min(0.5, 0.3 + 0.05 * (cryptopanic_count / 10))
            weights['CryptoCompare'] = 1.0 - weights['Reddit'] - weights['CryptoPanic']
            print(f"[DEBUG] Weights: Reddit={weights['Reddit']}, CryptoPanic={weights['CryptoPanic']}, CryptoCompare={weights['CryptoCompare']}")

            if sentiment['Reddit_Sentiment'].abs().sum() > 0:
                valid_sources.append(('Reddit', weights['Reddit']))
            if sentiment['CryptoPanic_Sentiment'].abs().sum() > 0:
                valid_sources.append(('CryptoPanic', weights['CryptoPanic']))
            if sentiment['CryptoCompare_Sentiment'].abs().sum() > 0:
                valid_sources.append(('CryptoCompare', weights['CryptoCompare']))
            print(f"[DEBUG] Valid sources: {valid_sources}")

            if valid_sources:
                total_weight = sum(w for _, w in valid_sources)
                sentiment['Combined_Sentiment'] = 0.0
                for source, weight in valid_sources:
                    sentiment['Combined_Sentiment'] += sentiment[f'{source}_Sentiment'] * weight
                if total_weight > 0:
                    sentiment['Combined_Sentiment'] /= total_weight
            else:
                sentiment['Combined_Sentiment'] = 0.0

            sentiment['Combined_Sentiment'] = sentiment['Combined_Sentiment'].fillna(0.0)
            self.sentiment_data = sentiment
            print("\nSentiment Table:")
            print(sentiment.to_string())
            print(f"[INFO] Sentiment analysis completed for {self.coin_symbol} ({len(sentiment)} days)")
            return sentiment

        except Exception as e:
            print(f"[ERROR] Error in sentiment analysis: {e}")
            import traceback
            traceback.print_exc()
            return None

    async def analyze_price_events(self, coin_symbol, start_date=None, end_date=None, session=None):
        try:
            file_path = self.coin_csv_mapping.get(coin_symbol.lower())
            if not file_path:
                print(f"[ERROR] No historical price data file found for {coin_symbol}")
                return None

            self.load_data_from_csv(file_path)
            if self.data is None:
                print(f"[ERROR] Failed to load historical price data for {coin_symbol}")
                return None

            price_data = self.data.copy()
            price_data.reset_index(inplace=True)

            if start_date and end_date:
                start_date = pd.to_datetime(start_date)
                end_date = pd.to_datetime(end_date)
                price_data = price_data[(price_data['Date'] >= start_date) & (price_data['Date'] <= end_date)]

            price_data['Daily_Return'] = price_data['Close'].pct_change() * 100

            CRASH_THRESHOLD = -10.0
            CORRECTION_THRESHOLD = -5.0
            SIGNIFICANT_DROP_THRESHOLD = -3.0

            events_df = pd.DataFrame(columns=[
                'Date', 'Event_Type', 'Price_Change_Pct',
                'Reddit_Sentiment', 'CryptoPanic_Sentiment',
                'CryptoCompare_Sentiment', 'Google_Trends_Sentiment',
                'Related_News'
            ])

            for _, row in price_data.iterrows():
                date = row['Date']
                if pd.isna(row['Daily_Return']):
                    continue

                event_type = None
                if row['Daily_Return'] <= CRASH_THRESHOLD:
                    event_type = 'Market Crash'
                elif row['Daily_Return'] <= CORRECTION_THRESHOLD:
                    event_type = 'Market Correction'
                elif row['Daily_Return'] <= SIGNIFICANT_DROP_THRESHOLD:
                    event_type = 'Significant Drop'

                if event_type:
                    event_date = date
                    sentiment_period_start = event_date - pd.Timedelta(days=1)
                    sentiment_period_end = event_date + pd.Timedelta(days=1)

                    if not hasattr(self, 'sentiment_data') or self.sentiment_data is None:
                        await self.analyze_sentiment(sentiment_period_start, sentiment_period_end, session)

                    event_sentiment = self.sentiment_data.loc[sentiment_period_start:sentiment_period_end] \
                        if self.sentiment_data is not None else pd.DataFrame()

                    news_headlines = await self.fetch_news_for_period(
                        coin_symbol, sentiment_period_start, sentiment_period_end, session
                    )

                    events_df = pd.concat([events_df, pd.DataFrame([{
                        'Date': event_date,
                        'Event_Type': event_type,
                        'Price_Change_Pct': row['Daily_Return'],
                        'Reddit_Sentiment': event_sentiment['Reddit_Sentiment'].mean() if not event_sentiment.empty else 0,
                        'CryptoPanic_Sentiment': event_sentiment['CryptoPanic_Sentiment'].mean() if not event_sentiment.empty else 0,
                        'CryptoCompare_Sentiment': event_sentiment['CryptoCompare_Sentiment'].mean() if not event_sentiment.empty else 0,
                        'Google_Trends_Sentiment': 0,
                        'Related_News': news_headlines
                    }])], ignore_index=True)

            events_df.sort_values('Date', inplace=True)
            self.event_summary_df = events_df
            print(f"[INFO] Found {len(events_df)} significant price events for {coin_symbol}")
            return events_df

        except Exception as e:
            print(f"[ERROR] Error analyzing price events: {e}")
            import traceback
            traceback.print_exc()
            return None

    async def fetch_news_for_period(self, coin_symbol, start_date, end_date, session):
        try:
            news_items = []
            crypto_panic_news = await self.fetch_cryptopanic_headlines(coin_symbol, start_date, end_date, session)
            if crypto_panic_news:
                news_items.extend([f"CryptoPanic: {headline}" for headline in crypto_panic_news])

            crypto_compare_news = await self.fetch_cryptocompare_headlines(coin_symbol, start_date, end_date, session)
            if crypto_compare_news:
                news_items.extend([f"CryptoCompare: {headline}" for headline in crypto_compare_news])

            return news_items
        except Exception as e:
            print(f"[ERROR] Error fetching news for period: {e}")
            return []

    async def fetch_cryptopanic_headlines(self, coin_symbol, start_date, end_date, session):
        try:
            start_str = start_date.strftime('%Y-%m-%d')
            end_str = end_date.strftime('%Y-%m-%d')
            api_key = self.api_keys.get('cryptopanic_api_key')
            if not api_key:
                return []

            base_url = "https://cryptopanic.com/api/v1/posts/"
            params = {
                'auth_token': api_key,
                'currencies': coin_symbol,
                'kind': 'news',
                'from': start_str,
                'to': end_str,
                'public': 'true'
            }

            async with session.get(base_url, params=params) as response:
                if response.status == 200:
                    data = await response.json()
                    results = data.get('results', [])
                    return [item.get('title', 'No Title') for item in results]
            return []
        except Exception as e:
            print(f"[ERROR] Error fetching CryptoPanic headlines: {e}")
            return []

    async def fetch_cryptocompare_headlines(self, coin_symbol, start_date, end_date, session):
        try:
            api_key = self.api_keys.get('cryptocompare_api_key')
            if not api_key:
                return []

            url = "https://min-api.cryptocompare.com/data/v2/news/"
            params = {
                'categories': coin_symbol,
                'api_key': api_key
            }

            async with session.get(url, params=params) as response:
                if response.status == 200:
                    data = await response.json()
                    news_items = data.get('Data', [])
                    start_ts = int(start_date.timestamp())
                    end_ts = int(end_date.timestamp())
                    return [article.get('title', 'No Title') for article in news_items if start_ts <= article.get('published_on', 0) <= end_ts]
            return []
        except Exception as e:
            print(f"[ERROR] Error fetching CryptoCompare headlines: {e}")
            return []

    def calculate_technical_indicators(self):
        try:
            if self.data is None or self.data.empty:
                print("[ERROR] No data available for technical indicators")
                return None

            data = self.data.copy()

            delta = data['Close'].diff()
            gain = (delta.where(delta > 0, 0)).rolling(window=14).mean()
            loss = (-delta.where(delta < 0, 0)).rolling(window=14).mean()
            rs = gain / loss
            data['RSI'] = 100 - (100 / (1 + rs))

            ema12 = data['Close'].ewm(span=12, adjust=False).mean()
            ema26 = data['Close'].ewm(span=26, adjust=False).mean()
            data['MACD'] = ema12 - ema26
            data['MACD_Signal'] = data['MACD'].ewm(span=9, adjust=False).mean()

            data['SMA_20'] = data['Close'].rolling(window=20).mean()
            data['EMA_26'] = data['Close'].ewm(span=26, adjust=False).mean()

            data['Market_Regime'] = 'Neutral'
            data.loc[data['RSI'] > 70, 'Market_Regime'] = 'Overbought'
            data.loc[data['RSI'] < 30, 'Market_Regime'] = 'Oversold'

            self.technical_indicators = data
            print("[INFO] Technical indicators calculated")
            return data
        except Exception as e:
            print(f"[ERROR] Error calculating technical indicators: {e}")
            return None

    async def detect_events(self):
        try:
            if self.data is None:
                raise ValueError("No data loaded")

            data = self.calculate_technical_indicators()

            price_change_threshold = 0.05
            volume_change_threshold = 0.5
            rsi_extreme = 70

            events = []

            for i in range(1, len(data)):
                prev_row = data.iloc[i-1]
                current_row = data.iloc[i]

                price_change = (current_row['Close'] - prev_row['Close']) / prev_row['Close']
                volume_change = (current_row['Volume'] - prev_row['Volume']) / prev_row['Volume'] if prev_row['Volume'] != 0 else 0

                if abs(price_change) >= price_change_threshold:
                    event_type = "Price Spike" if price_change > 0 else "Price Drop"
                    events.append({
                        'date': current_row.name,
                        'event_type': event_type,
                        'price_change': price_change,
                        'volume_change': volume_change,
                        'rsi': current_row['RSI'],
                        'sentiment': self.sentiment_data.loc[current_row.name, 'Combined_Sentiment'] if self.sentiment_data is not None and current_row.name in self.sentiment_data.index else 0
                    })

                elif abs(volume_change) >= volume_change_threshold:
                    events.append({
                        'date': current_row.name,
                        'event_type': "Volume Spike",
                        'price_change': price_change,
                        'volume_change': volume_change,
                        'rsi': current_row['RSI'],
                        'sentiment': self.sentiment_data.loc[current_row.name, 'Combined_Sentiment'] if self.sentiment_data is not None and current_row.name in self.sentiment_data.index else 0
                    })

                elif current_row['RSI'] >= rsi_extreme or current_row['RSI'] <= (100 - rsi_extreme):
                    event_type = "Overbought" if current_row['RSI'] >= rsi_extreme else "Oversold"
                    events.append({
                        'date': current_row.name,
                        'event_type': event_type,
                        'price_change': price_change,
                        'volume_change': volume_change,
                        'rsi': current_row['RSI'],
                        'sentiment': self.sentiment_data.loc[current_row.name, 'Combined_Sentiment'] if self.sentiment_data is not None and current_row.name in self.sentiment_data.index else 0
                    })

            if events:
                self.event_summary_df = pd.DataFrame(events)
                print(f"[INFO] Detected {len(events)} significant events")
                return self.event_summary_df
            else:
                print("[INFO] No significant events detected")
                return None

        except Exception as e:
            print(f"[ERROR] Error detecting events: {e}")
            return None

    def forecast_ema(self, days=30):
        try:
            if self.data is None or self.sentiment_data is None:
                raise ValueError("Historical or sentiment data missing")

            data = self.calculate_technical_indicators()
            last_date = data.index[-1]
            forecast_dates = pd.date_range(start=last_date + timedelta(days=1), periods=days)
            forecast = pd.DataFrame(index=forecast_dates, columns=['Forecasted_Price'])

            last_ema = data['EMA_26'].iloc[-1]
            last_price = data['Close'].iloc[-1]
            avg_sentiment = self.sentiment_data['Combined_Sentiment'].iloc[-10:].mean() if not self.sentiment_data.empty else 0

            for i, date in enumerate(forecast_dates):
                sentiment_factor = 1 + (avg_sentiment * 0.1)
                last_ema = last_ema + (last_price - last_ema) * (2 / (26 + 1))
                forecast.loc[date, 'Forecasted_Price'] = last_ema * sentiment_factor
                last_price = forecast.loc[date, 'Forecasted_Price']

            self.forecasts['EMA'] = forecast
            return forecast
        except Exception as e:
            print(f"[ERROR] Error forecasting with EMA: {e}")
            return None

    def forecast_arima(self, days=30):
        try:
            if self.data is None:
                raise ValueError("Historical data missing")

            prices = self.data['Close'].values
            model = sm.tsa.ARIMA(prices, order=(5, 1, 0))
            model_fit = model.fit()
            forecast = model_fit.forecast(steps=days)

            last_date = self.data.index[-1]
            forecast_dates = pd.date_range(start=last_date + timedelta(days=1), periods=days)
            forecast_df = pd.DataFrame({'Forecasted_Price': forecast}, index=forecast_dates)

            if self.sentiment_data is not None and not self.sentiment_data.empty:
                avg_sentiment = self.sentiment_data['Combined_Sentiment'].iloc[-10:].mean()
                sentiment_factor = 1 + (avg_sentiment * 0.1)
                forecast_df['Forecasted_Price'] *= sentiment_factor

            self.forecasts['ARIMA'] = forecast_df
            return forecast_df
        except Exception as e:
            print(f"[ERROR] Error forecasting with ARIMA: {e}")
            return None

    def forecast_lstm(self, days=30, look_back=60):
        try:
            if self.data is None:
                raise ValueError("Historical data missing")

            prices = self.data['Close'].values
            scaler = MinMaxScaler(feature_range=(0, 1))
            scaled_prices = scaler.fit_transform(prices.reshape(-1, 1))

            X, y = [], []
            for i in range(len(scaled_prices) - look_back):
                X.append(scaled_prices[i:(i + look_back), 0])
                y.append(scaled_prices[i + look_back, 0])
            X, y = np.array(X), np.array(y)
            X = np.reshape(X, (X.shape[0], X.shape[1], 1))

            model = Sequential()
            model.add(LSTM(units=50, return_sequences=True, input_shape=(look_back, 1)))
            model.add(Dropout(0.2))
            model.add(LSTM(units=50))
            model.add(Dropout(0.2))
            model.add(Dense(units=1))
            model.compile(optimizer='adam', loss='mean_squared_error')
            model.fit(X, y, epochs=20, batch_size=32, verbose=0)

            last_sequence = scaled_prices[-look_back:]
            forecast = []
            current_sequence = last_sequence.copy()

            for _ in range(days):
                current_sequence_reshaped = current_sequence.reshape((1, look_back, 1))
                next_pred = model.predict(current_sequence_reshaped, verbose=0)
                forecast.append(next_pred[0, 0])
                current_sequence = np.roll(current_sequence, -1)
                current_sequence[-1] = next_pred[0, 0]

            forecast = scaler.inverse_transform(np.array(forecast).reshape(-1, 1)).flatten()
            last_date = self.data.index[-1]
            forecast_dates = pd.date_range(start=last_date + timedelta(days=1), periods=days)
            forecast_df = pd.DataFrame({'Forecasted_Price': forecast}, index=forecast_dates)

            if self.sentiment_data is not None and not self.sentiment_data.empty:
                avg_sentiment = self.sentiment_data['Combined_Sentiment'].iloc[-10:].mean()
                sentiment_factor = 1 + (avg_sentiment * 0.1)
                forecast_df['Forecasted_Price'] *= sentiment_factor

            self.forecasts['LSTM'] = forecast_df
            return forecast_df
        except Exception as e:
            print(f"[ERROR] Error forecasting with LSTM: {e}")
            return None

    def forecast_prices(self, days=60):
        try:
            ema_forecast = self.forecast_ema(days=days)
            arima_forecast = self.forecast_arima(days=days)
            lstm_forecast = self.forecast_lstm(days=days)

            forecast_dates = pd.date_range(start=self.data.index[-1] + timedelta(days=1), periods=days, freq='D')
            combined_forecast = pd.DataFrame(index=forecast_dates, columns=['Forecasted_Price'])

            for date in forecast_dates:
                forecasts = []
                if ema_forecast is not None and date in ema_forecast.index:
                    forecasts.append(ema_forecast.loc[date, 'Forecasted_Price'])
                if arima_forecast is not None and date in arima_forecast.index:
                    forecasts.append(arima_forecast.loc[date, 'Forecasted_Price'])
                if lstm_forecast is not None and date in lstm_forecast.index:
                    forecasts.append(lstm_forecast.loc[date, 'Forecasted_Price'])
                combined_forecast.loc[date, 'Forecasted_Price'] = np.mean(forecasts) if forecasts else np.nan

            combined_forecast.dropna(inplace=True)
            self.forecasts['Combined'] = combined_forecast
            print(f"[INFO] Generated combined forecast for {days} days")
            return combined_forecast
        except Exception as e:
            print(f"[ERROR] Error combining forecasts: {e}")
            return None





    async def fetch_protocol_events(self, coin_symbol, start_date, end_date):
        events = []
        crypto_repos = {
            'BTC': {'name': 'Bitcoin', 'repos': ['bitcoin/bitcoin']},
            'ETH': {'name': 'Ethereum', 'repos': ['ethereum/go-ethereum', 'ethereum/EIPs']},
            'SOL': {'name': 'Solana', 'repos': ['solana-labs/solana']},
            'ADA': {'name': 'Cardano', 'repos': ['input-output-hk/cardano-node']},
            'DOT': {'name': 'Polkadot', 'repos': ['paritytech/polkadot']}
        }
        coin_info = crypto_repos.get(coin_symbol.upper())
        if not coin_info:
            self.logger.info(f"No GitHub repositories mapped for {coin_symbol}")
            return events
        try:
            headers = {'Authorization': f"token {self.api_keys['github_token']}" if 'github_token' in self.api_keys else {}}
            async with aiohttp.ClientSession() as session:
                for repo in coin_info['repos']:
                    release_url = f"https://api.github.com/repos/{repo}/releases"
                    async with session.get(release_url, headers=headers) as response:
                        if response.status == 200:
                            releases = await response.json()
                            for release in releases:
                                if 'published_at' in release and release['published_at']:
                                    release_date = pd.to_datetime(release['published_at']).tz_localize(None)
                                    if start_date <= release_date <= end_date:
                                        events.append({
                                            'date': release_date,
                                            'name': f"{coin_info['name']} Release: {release['name']}",
                                            'impact': 'Moderate Positive',
                                            'source': 'GitHub',
                                            'description': release['body'][:100] + '...' if release.get('body') and len(release['body']) > 100 else release.get('body', 'No description')
                                        })
                        else:
                            self.logger.warning(f"GitHub API (releases) returned status {response.status} for {repo}")
                    milestone_url = f"https://api.github.com/repos/{repo}/milestones"
                    async with session.get(milestone_url, headers=headers) as response:
                        if response.status == 200:
                            milestones = await response.json()
                            for milestone in milestones:
                                if 'due_on' in milestone and milestone['due_on']:
                                    due_date = pd.to_datetime(milestone['due_on']).tz_localize(None)
                                    if start_date <= due_date <= end_date:
                                        events.append({
                                            'date': due_date,
                                            'name': f"{coin_info['name']} Milestone: {milestone['title']}",
                                            'impact': 'Moderate Positive',
                                            'source': 'GitHub',
                                            'description': milestone['description'][:100] + '...' if milestone.get('description') and len(milestone['description']) > 100 else milestone.get('description', 'No description')
                                        })
                        else:
                            self.logger.warning(f"GitHub API (milestones) returned status {response.status} for {repo}")
            self.logger.info(f"Fetched {len(events)} protocol events for {coin_symbol}")
            return events
        except Exception as e:
            self.logger.error(f"Error fetching protocol events: {e}")
            return []



    async def check_upcoming_events(self, coin_symbol, forecast_days=30):
        events = []
        try:
            today = pd.Timestamp.now().normalize()
            start_date = today
            end_date = today + pd.Timedelta(days=forecast_days)

            # Fetching various event types and adding event_type and impact_classification
            protocol_events = await self.fetch_protocol_events(coin_symbol, start_date, end_date)
            for event in protocol_events:
                event['event_type'] = 'protocol'
                impact = event.get('impact', 'unknown').lower()
                if 'positive' in impact or impact == 'high':
                    event['impact_classification'] = 'good'
                elif impact == 'medium':
                    event['impact_classification'] = 'neutral'
                elif impact == 'low':
                    event['impact_classification'] = 'bad'
                else:
                    event['impact_classification'] = 'unknown'
            self.logger.info(f"Protocol Events: {len(protocol_events)}")
            for event in protocol_events:
                self.logger.info(f"Protocol: {event['name']}, Date: {event['date']}, Impact: {event['impact']}, Classification: {event['impact_classification']}")

            token_unlock_events = await self.fetch_token_unlock_events(coin_symbol, start_date, end_date)
            for event in token_unlock_events:
                event['event_type'] = 'token_unlock'
                impact = event.get('impact', 'unknown').lower()
                if impact == 'high':
                    event['impact_classification'] = 'good'
                elif impact == 'medium':
                    event['impact_classification'] = 'neutral'
                elif impact == 'low':
                    event['impact_classification'] = 'bad'
                else:
                    event['impact_classification'] = 'unknown'
            self.logger.info(f"Token Unlock Events: {len(token_unlock_events)}")
            for event in token_unlock_events:
                self.logger.info(f"Token Unlock: {event['name']}, Date: {event['date']}, Impact: {event['impact']}, Classification: {event['impact_classification']}")

            economic_events = await self.fetch_economic_events(start_date, end_date)
            for event in economic_events:
                event['event_type'] = 'economic'
                impact = event.get('impact', 'unknown').lower()
                if impact == 'high':
                    event['impact_classification'] = 'good'
                elif impact == 'medium':
                    event['impact_classification'] = 'neutral'
                elif impact == 'low':
                    event['impact_classification'] = 'bad'
                else:
                    event['impact_classification'] = 'unknown'
            self.logger.info(f"Economic Events: {len(economic_events)}")
            for event in economic_events:
                self.logger.info(f"Economic: {event['name']}, Date: {event['date']}, Impact: {event['impact']}, Classification: {event['impact_classification']}")

            crypto_specific_events = await self.check_crypto_specific_events(coin_symbol, start_date, end_date)
            for event in crypto_specific_events:
                event['event_type'] = 'crypto_specific'
                impact = event.get('impact', 'unknown').lower()
                if impact == 'high':
                    event['impact_classification'] = 'good'
                elif impact == 'medium':
                    event['impact_classification'] = 'neutral'
                elif impact == 'low':
                    event['impact_classification'] = 'bad'
                else:
                    event['impact_classification'] = 'unknown'
            self.logger.info(f"Crypto-Specific Events: {len(crypto_specific_events)}")
            for event in crypto_specific_events:
                self.logger.info(f"Crypto-Specific: {event['name']}, Date: {event['date']}, Impact: {event['impact']}, Classification: {event['impact_classification']}")

            regulatory_events = await self.fetch_regulatory_events(coin_symbol, start_date, end_date)
            for event in regulatory_events:
                event['event_type'] = 'regulatory'
                impact = event.get('impact', 'unknown').lower()
                if impact == 'high':
                    event['impact_classification'] = 'good'  # Assume high regulatory impact is positive (e.g., clarity)
                elif impact == 'medium':
                    event['impact_classification'] = 'neutral'
                elif impact == 'low':
                    event['impact_classification'] = 'bad'
                else:
                    event['impact_classification'] = 'unknown'
            self.logger.info(f"Regulatory Events: {len(regulatory_events)}")
            for event in regulatory_events:
                self.logger.info(f"Regulatory: {event['name']}, Date: {event['date']}, Impact: {event['impact']}, Classification: {event['impact_classification']}")

            exchange_events = await self.fetch_exchange_events(coin_symbol, start_date, end_date)
            for event in exchange_events:
                event['event_type'] = 'exchange'
                impact = event.get('impact', 'unknown').lower()
                if impact == 'high':
                    event['impact_classification'] = 'good'
                elif impact == 'medium':
                    event['impact_classification'] = 'neutral'
                elif impact == 'low':
                    event['impact_classification'] = 'bad'
                else:
                    event['impact_classification'] = 'unknown'
            self.logger.info(f"Exchange Events: {len(exchange_events)}")
            for event in exchange_events:
                self.logger.info(f"Exchange: {event['name']}, Date: {event['date']}, Impact: {event['impact']}, Classification: {event['impact_classification']}")

            # Collecting all events
            events.extend(protocol_events)
            events.extend(token_unlock_events)
            events.extend(economic_events)
            events.extend(crypto_specific_events)
            events.extend(regulatory_events)
            events.extend(exchange_events)

            # Sorting events by date
            events.sort(key=lambda x: x['date'])

            self.logger.info(f"Total Upcoming Events: {len(events)}")
            return events

        except Exception as e:
            self.logger.error(f"Error checking upcoming events: {e}")
            return []







    async def fetch_token_unlock_events(self, coin_symbol, start_date, end_date):
        events = []
        try:
            # Map coin symbol to CoinGecko ID
            coin_ids = {
                'BTC': 'bitcoin',
                'ETH': 'ethereum',
                'SOL': 'solana'
            }
            coin_id = coin_ids.get(coin_symbol.upper())
            if not coin_id:
                self.logger.info(f"No CoinGecko ID mapped for {coin_symbol}")
                return events

            # Fetch coin metadata to get contract address (if available)
            async with aiohttp.ClientSession() as session:
                url = f"https://api.coingecko.com/api/v3/coins/{coin_id}"
                async with session.get(url) as response:
                    if response.status == 200:
                        coin_data = await response.json()
                        contract_address = coin_data.get('contract_address')  # May be None for native coins
                    else:
                        self.logger.warning(f"CoinGecko API returned status {response.status} for {coin_id}")
                        contract_address = None

            # Fetch market data to detect unlock events via volume spikes
            url = f"https://api.coingecko.com/api/v3/coins/{coin_id}/market_chart/range"
            params = {
                'vs_currency': 'usd',
                'from': int(start_date.timestamp()),
                'to': int(end_date.timestamp())
            }

            async with aiohttp.ClientSession() as session:
                for attempt in range(3):
                    async with session.get(url, params=params) as response:
                        if response.status == 200:
                            market_data = await response.json()
                            volumes = market_data.get('total_volumes', [])

                            if volumes:
                                baseline_volume = sum(v[1] for v in volumes) / len(volumes) if volumes else 0
                                for i in range(1, len(volumes)):
                                    current_date = pd.to_datetime(volumes[i][0] / 1000, unit='s')
                                    if start_date <= current_date <= end_date:
                                        prev_volume = volumes[i - 1][1] if i > 0 else baseline_volume
                                        current_volume = volumes[i][1]
                                        volume_change = (current_volume - prev_volume) / prev_volume if prev_volume else 0
                                        if volume_change > 1.0:
                                            impact = 'High' if volume_change > 2.0 else 'Medium'
                                            events.append({
                                                'date': current_date,
                                                'name': f"{coin_symbol} Potential Token Unlock Event",
                                                'impact': impact,
                                                'source': 'CoinGecko Market Data',
                                                'description': f"Volume spike of {volume_change:.2%} detected"
                                            })
                            break
                        elif response.status == 429:
                            self.logger.warning(f"CoinGecko API rate limit exceeded, retrying ({attempt + 1}/3)...")
                            await asyncio.sleep(2 ** attempt)
                        else:
                            self.logger.warning(f"CoinGecko API returned status {response.status} for market data")
                            break

            # Predefined token unlock events
            known_unlocks = {
                'SOL': [
                    {'date': '2025-06-01', 'name': 'Solana Validator Unlock', 'percentage': 3.5, 'amount': 1000000},
                    {'date': '2025-07-15', 'name': 'Solana Team Unlock', 'percentage': 2.0, 'amount': 500000}
                ],
                'ETH': [
                    {'date': '2025-08-15', 'name': 'Ethereum Foundation Unlock', 'percentage': 1.2, 'amount': 200000},
                    {'date': '2025-09-01', 'name': 'Ethereum Staking Unlock', 'percentage': 0.8, 'amount': 150000}
                ],
                'BTC': [
                    {'date': '2025-06-15', 'name': 'Bitcoin Halving Aftermath Unlock', 'percentage': 1.0, 'amount': 900000}
                ]
            }

            if coin_symbol.upper() in known_unlocks:
                for unlock in known_unlocks[coin_symbol.upper()]:
                    unlock_date = pd.to_datetime(unlock['date'])
                    if start_date <= unlock_date <= end_date:
                        unlock_pct = unlock['percentage']
                        impact = 'High' if unlock_pct > 10 else 'Medium' if unlock_pct > 5 else 'Low'
                        events.append({
                            'date': unlock_date,
                            'name': unlock['name'],
                            'impact': impact,
                            'source': 'Predefined Unlocks Database',
                            'description': f"Unlock of {unlock['amount']} tokens ({unlock_pct}% of supply)"
                        })

            self.logger.info(f"Fetched {len(events)} token unlock events for {coin_symbol}")
            return events
        except Exception as e:
            self.logger.error(f"Error fetching token unlock events: {e}")
            return []















    async def fetch_economic_events(self, start_date, end_date):
        events = []
        try:
            # Create a custom SSL context with TLS 1.2
            ssl_context = ssl.create_default_context()
            ssl_context.minimum_version = ssl.TLSVersion.TLSv1_2

            # Fetch economic events from FRED API
            async with aiohttp.ClientSession(
                connector=aiohttp.TCPConnector(ssl=ssl_context)
            ) as session:
                releases = [
                    {'release_id': '53', 'name': 'Consumer Price Index (CPI)', 'impact': 'High'},
                    {'release_id': '10', 'name': 'Employment Situation (Non-Farm Payrolls)', 'impact': 'High'},
                    {'release_id': '46', 'name': 'Gross Domestic Product (GDP)', 'impact': 'Medium'},
                    {'release_id': '14', 'name': 'Unemployment Rate', 'impact': 'High'},
                    {'release_id': '82', 'name': 'Federal Funds Rate (FOMC)', 'impact': 'High'}
                ]

                for release in releases:
                    release_id = release['release_id']
                    url = "https://api.stlouisfed.org/fred/releases/dates"
                    params = {
                        'release_id': release_id,
                        'api_key': self.api_keys.get('fred_api_key', 'YOUR_FRED_API_KEY'),
                        'file_type': 'json',
                        'realtime_start': start_date.strftime('%Y-%m-%d'),
                        'realtime_end': end_date.strftime('%Y-%m-%d')
                    }
                    for attempt in range(3):
                        try:
                            async with session.get(url, params=params) as response:
                                if response.status == 200:
                                    data = await response.json()
                                    if 'release_dates' not in data:
                                        self.logger.warning(f"FRED API returned no release dates for release ID {release_id}")
                                        break
                                    for entry in data['release_dates']:
                                        event_date = pd.to_datetime(entry['date'])
                                        if start_date <= event_date <= end_date:
                                            events.append({
                                                'date': event_date,
                                                'name': f"{release['name']} Release",
                                                'impact': release['impact'],
                                                'source': 'FRED API',
                                                'description': f"{release['name']} Release"
                                            })
                                    break
                                elif response.status == 429:
                                    self.logger.warning(f"FRED API rate limit exceeded for release ID {release_id}, retrying ({attempt + 1}/3)...")
                                    await asyncio.sleep(2 ** (attempt + 1))
                                else:
                                    self.logger.warning(f"FRED API returned status {response.status} for release ID {release_id}")
                                    break
                        except Exception as e:
                            self.logger.warning(f"Request failed for FRED release ID {release_id}: {e}")
                            break
                    await asyncio.sleep(0.5)

            # Predefined economic events as fallback (unchanged)
            known_economic_events = [
                {'date': '2025-05-07', 'name': 'US Non-Farm Payrolls', 'type': 'Employment Report', 'impact': 'High'},
                {'date': '2025-05-15', 'name': 'FOMC Interest Rate Decision', 'type': 'FOMC Meeting', 'impact': 'High'},
                {'date': '2025-05-20', 'name': 'US CPI Data Release', 'type': 'CPI Report', 'impact': 'High'},
                {'date': '2025-06-01', 'name': 'EU GDP Release', 'type': 'GDP Report', 'impact': 'Medium'},
                {'date': '2025-06-15', 'name': 'ECB Monetary Policy Meeting', 'type': 'ECB Meeting', 'impact': 'Medium'},
                {'date': '2025-07-01', 'name': 'US Retail Sales', 'type': 'Retail Sales Report', 'impact': 'Medium'},
                {'date': '2025-07-10', 'name': 'US PPI Release', 'type': 'PPI Report', 'impact': 'Medium'},
                {'date': '2025-08-01', 'name': 'US Unemployment Rate', 'type': 'Employment Report', 'impact': 'High'},
                {'date': '2025-08-15', 'name': 'China Industrial Production', 'type': 'Industrial Data', 'impact': 'Medium'},
                {'date': '2025-09-01', 'name': 'US Labor Day (Market Holiday)', 'type': 'Market Holiday', 'impact': 'Low'}
            ]

            for event in known_economic_events:
                event_date = pd.to_datetime(event['date'])
                if start_date <= event_date <= end_date:
                    events.append({
                        'date': event_date,
                        'name': event['name'],
                        'impact': event['impact'],
                        'source': 'Predefined Economic Calendar',
                        'description': f"{event['type']} - {event['name']}"
                    })

            events_df = pd.DataFrame(events)
            if not events_df.empty:
                events_df = events_df.drop_duplicates(subset=['date', 'name'], keep='first')
                events = events_df.to_dict('records')

            print(f"[INFO] Fetched {len(events)} economic events")
            return events
        except Exception as e:
            print(f"[ERROR] Error fetching economic events: {e}")
            return []

    async def check_crypto_specific_events(self, coin_symbol, start_date, end_date):
        events = []
        try:
            # Bitcoin Halving
            if coin_symbol.upper() == 'BTC':
                next_halving = pd.Timestamp('2028-04-20')
                if start_date <= next_halving <= end_date:
                    events.append({
                        'date': next_halving,
                        'name': 'Bitcoin Halving',
                        'impact': 'High',
                        'source': 'Bitcoin Protocol',
                        'description': 'Block reward will be halved, reducing new Bitcoin supply'
                    })

            # Ethereum-specific events
            elif coin_symbol.upper() == 'ETH':
                eth_events = [
                    {'date': '2025-07-15', 'name': 'EIP-7600 Implementation', 'impact': 'Medium'},
                    {'date': '2025-09-30', 'name': 'Ethereum Layer 2 Integration', 'impact': 'High'}
                ]
                for event in eth_events:
                    event_date = pd.to_datetime(event['date'])
                    if start_date <= event_date <= end_date:
                        events.append({
                            'date': event_date,
                            'name': event['name'],
                            'impact': event['impact'],
                            'source': 'Ethereum Roadmap',
                            'description': f"Scheduled Ethereum upgrade: {event['name']}"
                        })

            # Solana-specific events
            elif coin_symbol.upper() == 'SOL':
                sol_events = [
                    {'date': '2025-05-20', 'name': 'Solana Network Upgrade v2.0', 'impact': 'Medium'}
                ]
                for event in sol_events:
                    event_date = pd.to_datetime(event['date'])
                    if start_date <= event_date <= end_date:
                        events.append({
                            'date': event_date,
                            'name': event['name'],
                            'impact': event['impact'],
                            'source': 'Solana Roadmap',
                            'description': f"Scheduled Solana upgrade: {event['name']}"
                        })

            print(f"[INFO] Fetched {len(events)} crypto-specific events for {coin_symbol}")
            return events
        except Exception as e:
            print(f"[ERROR] Error checking crypto-specific events: {e}")
            return []

    async def fetch_regulatory_events(self, coin_symbol, start_date, end_date):
        events = []
        try:
            async with aiohttp.ClientSession() as session:
                url = "https://www.federalregister.gov/api/v1/documents"
                params = {
                    'conditions[term]': 'cryptocurrency',
                    'per_page': 20,
                    'order': 'newest'
                }
                async with session.get(url, params=params) as response:
                    if response.status == 200:
                        data = await response.json()
                        for doc in data.get('results', []):
                            event_date = pd.to_datetime(doc.get('publication_date'))
                            if start_date <= event_date <= end_date:
                                title = doc.get('title', 'Unknown Regulatory Event')
                                impact = 'High' if 'regulation' in title.lower() or 'sec' in title.lower() else 'Medium'
                                events.append({
                                    'date': event_date,
                                    'name': title,
                                    'impact': impact,
                                    'source': 'Federal Register API',
                                    'description': doc.get('abstract', 'No description available')
                                })
                    else:
                        print(f"[WARNING] Federal Register API returned status {response.status}")

            print(f"[INFO] Fetched {len(events)} regulatory events for {coin_symbol}")
            return events
        except Exception as e:
            print(f"[ERROR] Error fetching regulatory events: {e}")
            return []








    async def fetch_exchange_events(self, coin_symbol, start_date, end_date):
        events = []
        try:
            # Map coin symbol to CoinGecko ID
            coin_ids = {
                'BTC': 'bitcoin',
                'ETH': 'ethereum',
                'SOL': 'solana'
            }
            coin_id = coin_ids.get(coin_symbol.upper())
            if not coin_id:
                self.logger.info(f"No CoinGecko ID mapped for {coin_symbol}")
                return events

            # Reduce to one exchange to minimize API calls
            exchanges_to_check = ['binance']

            # Fetch tickers from major exchanges to detect new trading pairs
            async with aiohttp.ClientSession() as session:
                headers = {'x-cg-api-key': self.api_keys.get('coingecko_api_key', '')}
                for exchange_id in exchanges_to_check:
                    url = f"https://api.coingecko.com/api/v3/exchanges/{exchange_id}/tickers"
                    params = {'coin_ids': coin_id}
                    for attempt in range(3):  # Retry up to 3 times for rate limits
                        try:
                            async with session.get(url, headers=headers, params=params) as response:
                                if response.status == 200:
                                    data = await response.json()
                                    tickers = data.get('tickers', [])
                                    for ticker in tickers:
                                        if ticker.get('coin_id') == coin_id:
                                            last_trade_timestamp = ticker.get('timestamp')
                                            if last_trade_timestamp:
                                                event_date = pd.to_datetime(last_trade_timestamp).tz_localize(None)
                                                if start_date <= event_date <= end_date:
                                                    base = ticker.get('base', coin_symbol.upper())
                                                    target = ticker.get('target', 'Unknown')
                                                    events.append({
                                                        'date': event_date,
                                                        'name': f"{coin_symbol} Listed on {exchange_id.replace('-', ' ').title()} ({base}/{target})",
                                                        'impact': 'Medium',
                                                        'source': 'CoinGecko API',
                                                        'description': f"New trading pair {base}/{target} detected on {exchange_id}"
                                                    })
                                    break  # Exit retry loop on success
                                elif response.status == 429:
                                    self.logger.warning(f"Rate limit exceeded for {exchange_id}, retrying ({attempt + 1}/3)...")
                                    await asyncio.sleep(2 ** (attempt + 1))  # Exponential backoff: 2s, 4s, 8s
                                else:
                                    self.logger.warning(f"API error {response.status} from {exchange_id}")
                                    break
                        except Exception as e:
                            self.logger.warning(f"Request failed for {exchange_id}: {e}")
                            break

            # Predefined exchange events
            known_exchange_events = [
                {'date': '2025-05-10', 'name': f"{coin_symbol} Listed on Coinbase", 'impact': 'Medium'},
                {'date': '2025-05-15', 'name': f"{coin_symbol} Staking on Binance", 'impact': 'Medium'}
            ]
            for event in known_exchange_events:
                event_date = pd.to_datetime(event['date'])
                if start_date <= event_date <= end_date:
                    events.append({
                        'date': event_date,
                        'name': event['name'],
                        'impact': event['impact'],
                        'source': 'Predefined Exchange Events',
                        'description': event['name']
                    })

            self.logger.info(f"Fetched {len(events)} exchange events for {coin_symbol}")
            return events
        except Exception as e:
            self.logger.error(f"Error fetching exchange events: {e}")
            return []











    def setup_rag_system(self):
        try:
            documents = [
                f"{self.coin_name} ({self.coin_symbol}) is a cryptocurrency influenced by market sentiment, news, and events like tariff augmentations, elections, Bitcoin halvings, and regulatory changes.",
                f"Historical price data includes Open, High, Low, Close, Volume, and Market Cap, used to analyze past event impacts and forecast future prices."
            ]

            if self.data is not None and not self.data.empty:
                latest = self.data.iloc[-1]
                prev_day = self.data.iloc[-2] if len(self.data) > 1 else None
                price_change = (latest['Close'] - prev_day['Close']) / prev_day['Close'] * 100 if prev_day is not None else 0
                documents.extend([
                    f"Latest price for {self.coin_symbol}: ${latest['Close']:.2f} on {self.data.index[-1].strftime('%Y-%m-%d')}.",
                    f"24h price change: {price_change:.2f}%.",
                    f"Latest market cap: ${latest['Market Cap']:.2f}.",
                    f"24h return: {self.data['Return'].iloc[-1]*100:.2f}%."
                ])

            if self.sentiment_data is not None and not self.sentiment_data.empty:
                reddit_last_nonzero = self.sentiment_data['Reddit_Sentiment'][self.sentiment_data['Reddit_Sentiment'] != 0].last_valid_index()
                reddit_sentiment = self.sentiment_data['Reddit_Sentiment'].loc[reddit_last_nonzero] if reddit_last_nonzero is not None else 0
                cryptopanic_last_nonzero = self.sentiment_data['CryptoPanic_Sentiment'][self.sentiment_data['CryptoPanic_Sentiment'] != 0].last_valid_index()
                cryptopanic_sentiment = self.sentiment_data['CryptoPanic_Sentiment'].loc[cryptopanic_last_nonzero] if cryptopanic_last_nonzero is not None else 0
                cryptocompare_last_nonzero = self.sentiment_data['CryptoCompare_Sentiment'][self.sentiment_data['CryptoCompare_Sentiment'] != 0].last_valid_index()
                cryptocompare_sentiment = self.sentiment_data['CryptoCompare_Sentiment'].loc[cryptocompare_last_nonzero] if cryptocompare_last_nonzero is not None else 0
                combined_last_nonzero = self.sentiment_data['Combined_Sentiment'][self.sentiment_data['Combined_Sentiment'] != 0].last_valid_index()
                combined_sentiment = self.sentiment_data['Combined_Sentiment'].loc[combined_last_nonzero] if combined_last_nonzero is not None else 0

                avg_sentiment = self.sentiment_data['Combined_Sentiment'].iloc[-10:].mean()
                sentiment_desc = 'positive' if combined_sentiment > 0.2 else 'negative' if combined_sentiment < -0.2 else 'neutral'
                documents.extend([
                    f"Latest sentiment for {self.coin_symbol} on {self.sentiment_data.index[-1].strftime('%Y-%m-%d')}:",
                    f"Reddit sentiment: {reddit_sentiment:.2f} ({'no data' if reddit_sentiment == 0 else sentiment_desc}).",
                    f"CryptoPanic sentiment: {cryptopanic_sentiment:.2f} ({'no data' if cryptopanic_sentiment == 0 else sentiment_desc}).",
                    f"CryptoCompare sentiment: {cryptocompare_sentiment:.2f} ({'no data' if cryptocompare_sentiment == 0 else sentiment_desc}).",
                    f"Combined sentiment: {combined_sentiment:.2f} ({sentiment_desc}).",
                    f"Average combined sentiment over last 10 days: {avg_sentiment:.2f} ({'positive' if avg_sentiment > 0.2 else 'negative' if avg_sentiment < -0.2 else 'neutral'})."
                ])

            if self.event_summary_df is not None and not self.event_summary_df.empty:
                for _, row in self.event_summary_df.iterrows():
                    sentiment_desc = 'positive' if row.get('sentiment', 0) > 0.2 else 'negative' if row.get('sentiment', 0) < -0.2 else 'neutral'
                    documents.append(
                        f"Event on {row['date'].strftime('%Y-%m-%d')}: {row['event_type']}, price changed by {row['price_change']*100:.2f}% with {sentiment_desc} sentiment."
                    )

            docs = [Document(page_content=doc) for doc in documents]
            embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
            self.vectorstore = FAISS.from_documents(docs, embedding_model)
            print("[INFO] RAG system set up successfully.")
            return self.vectorstore

        except Exception as e:
            print(f"[ERROR] Failed to set up RAG: {e}")
            return None



    def generate_interactive_plots(self, output_dir='plots'):
        try:
            os.makedirs(output_dir, exist_ok=True)

            # Forecast Plot
            if hasattr(self, 'forecasts') and self.forecasts.get('Combined') is not None and not self.forecasts['Combined'].empty:
                historical_data = self.data.tail(100)  # Limit to last 100 days
                fig = go.Figure()
                fig.add_trace(go.Scatter(x=historical_data.index, y=historical_data['Close'], name='Historical Price', line=dict(color='blue')))
                fig.add_trace(go.Scatter(x=self.forecasts['Combined'].index, y=self.forecasts['Combined']['Forecasted_Price'], name='Combined Forecast', line=dict(color='red', dash='dash')))
                fig.update_layout(title=f"{self.coin_symbol.upper()} Price Forecast", xaxis_title="Date", yaxis_title="Price (USD)", template="plotly_dark")
                fig.write_html(os.path.join(output_dir, f"{self.coin_symbol}_forecast.html"))
                self.logger.info(f"Forecast plot saved to {output_dir}/{self.coin_symbol}_forecast.html")

            # Sentiment Plot
            if self.sentiment_data is not None and not self.sentiment_data.empty:
                sentiment_data = self.sentiment_data.tail(100)  # Limit to last 100 days
                fig = go.Figure()
                for col in ['Combined_Sentiment']:  # Focus on combined sentiment
                    if col in sentiment_data.columns:
                        fig.add_trace(go.Scatter(x=sentiment_data.index, y=sentiment_data[col], name=col.replace('_', ' '), mode='lines'))
                fig.update_layout(title=f"{self.coin_symbol.upper()} Sentiment Analysis", xaxis_title="Date", yaxis_title="Sentiment Score (-1 to 1)", template="plotly_dark")
                fig.write_html(os.path.join(output_dir, f"{self.coin_symbol}_sentiment.html"))
                self.logger.info(f"Sentiment plot saved to {output_dir}/{self.coin_symbol}_sentiment.html")

            # Event Plot
            self.logger.info("Starting Event Plot generation...")
            if self.event_summary_df is not None and not self.event_summary_df.empty:
                self.logger.info(f"Total events in event_summary_df: {len(self.event_summary_df)}")

                price_change_threshold = 0.05  # 5% price change
                important_events = self.event_summary_df[
                    self.event_summary_df['price_change'].abs() > price_change_threshold
                ].copy()

                max_events = 10
                important_events = important_events.tail(max_events)

                self.logger.info(f"Filtered to {len(important_events)} important events (price change > {price_change_threshold*100}%)")
                if len(important_events) == 0:
                    self.logger.info("No important events to plot.")
                    return True

                self.logger.info(f"Sample of important events:\n{important_events.head().to_string()}")

                data = self.data.tail(100)  # Limit to last 100 days
                fig = go.Figure()
                fig.add_trace(go.Scatter(x=data.index, y=data['Close'], name='Price', line=dict(color='blue')))

                for idx, event in important_events.iterrows():
                    try:
                        event_date = pd.to_datetime(event['date'])
                        if event_date.tzinfo is None:
                            event_date = event_date.tz_localize('UTC')
                        else:
                            event_date = event_date.tz_convert('UTC')

                        event_date_ms = int(event_date.timestamp() * 1000)
                        event_type = str(event['event_type']) if pd.notna(event['event_type']) else "Unknown Event"

                        self.logger.debug(f"Processing event {idx + 1}/{len(important_events)}: {event_type} at {event_date} (price change: {event['price_change']*100:.2f}%)")

                        fig.add_vline(
                            x=event_date_ms,
                            line_dash="dash",
                            line_color="red",
                            annotation_text=event_type,
                            annotation_position="top"
                        )
                        self.logger.debug(f"Added vline for event {idx + 1}/{len(important_events)}")
                    except Exception as e:
                        self.logger.warning(f"Failed to plot event at index {idx}: {e}")
                        continue

                self.logger.info("Updating plot layout...")
                fig.update_layout(
                    title=f"{self.coin_symbol.upper()} Price with Significant Events",
                    xaxis_title="Date",
                    yaxis_title="Price (USD)",
                    template="plotly_dark",
                    xaxis=dict(type='date')
                )

                self.logger.info("Writing event plot to file...")
                fig.write_html(os.path.join(output_dir, f"{self.coin_symbol}_events.html"))
                self.logger.info(f"Event plot saved to {output_dir}/{self.coin_symbol}_events.html")

            else:
                self.logger.info("No events to plot (event_summary_df is empty).")

            return True
        except Exception as e:
            self.logger.error(f"Error generating interactive plots: {e}")
            return None



    def export_analysis_to_csv(self, output_path="crypto_analysis_export.csv"):
        try:
            export_data = self.data.copy() if self.data is not None else pd.DataFrame()

            if self.sentiment_data is not None:
                export_data = export_data.join(
                    self.sentiment_data, how='outer', lsuffix='_data', rsuffix='_sentiment'
                )

            if hasattr(self, 'technical_indicators') and self.technical_indicators is not None:
                export_data = export_data.join(
                    self.technical_indicators.drop(
                        columns=['Open', 'High', 'Low', 'Close', 'Volume', 'Market Cap', 'Return'],
                        errors='ignore'
                    ),
                    how='outer',
                    lsuffix='_data',
                    rsuffix='_indicators'
                )

            if hasattr(self, 'forecasts') and self.forecasts.get('Combined') is not None:
                export_data = pd.concat([export_data, self.forecasts['Combined']], axis=0)

            export_data.to_csv(output_path)
            print(f"[INFO] Analysis exported to {output_path}")
            return output_path

        except Exception as e:
            print(f"[ERROR] Error exporting to CSV: {e}")
            return None




















    async def generate_pdf_report(self, pdf_filepath='crypto_report.pdf', upcoming_events=None):
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
            from reportlab.lib.styles import getSampleStyleSheet
            import matplotlib.pyplot as plt
            import os
            from datetime import datetime

            def download_pixabay_image(api_key, query, output_path):
                """Download an image from Pixabay API."""
                url = f"https://pixabay.com/api/?key={api_key}&q={query}&image_type=photo&per_page=3"
                response = requests.get(url)
                if response.status_code == 200:
                    data = response.json()
                    if data['hits']:
                        image_url = data['hits'][0]['largeImageURL']
                        img_response = requests.get(image_url)
                        if img_response.status_code == 200:
                            with open(output_path, 'wb') as f:
                                f.write(img_response.content)
                            return True
                return False

            # Pixabay API key
            pixabay_api_key = "50210954-a7aca149c8b1407da7cd5e6bc"

            # Generate PDF
            doc = SimpleDocTemplate(pdf_filepath, pagesize=letter)
            styles = getSampleStyleSheet()
            pdf_elements = []

            # Title (PDF)
            pdf_elements.append(Paragraph(f"Cryptocurrency Analysis Report: {self.coin_name} ({self.coin_symbol})", styles['Title']))
            pdf_elements.append(Spacer(1, 12))
            pdf_elements.append(Paragraph(f"Report generated on {datetime.now().strftime('%Y-%m-%d')}", styles['Normal']))
            pdf_elements.append(Spacer(1, 24))

            # Price Chart (PDF)
            pdf_elements.append(Paragraph("Price Chart", styles['Heading2']))
            fig, ax = plt.subplots(figsize=(8, 4))
            self.data['Close'].plot(ax=ax)
            ax.set_title(f"{self.coin_symbol} Historical Price")
            ax.set_xlabel('Date')
            ax.set_ylabel('Price (USD)')
            chart_path = 'price_chart.png'
            plt.savefig(chart_path)
            plt.close()
            pdf_elements.append(Image(chart_path, width=400, height=200))
            pdf_elements.append(Spacer(1, 12))

            # Price and Sentiment Combined Chart (PDF)
            pdf_elements.append(Paragraph("Price and Sentiment Combined Chart", styles['Heading2']))
            if self.sentiment_data is not None and not self.sentiment_data.empty and self.data is not None and not self.data.empty:
                combined_data = self.data[['Close']].join(self.sentiment_data[['Combined_Sentiment']], how='inner')
                if not combined_data.empty:
                    fig, ax1 = plt.subplots(figsize=(8, 4))
                    ax1.plot(combined_data.index, combined_data['Close'], color='blue', label='Price (USD)')
                    ax1.set_xlabel('Date')
                    ax1.set_ylabel('Price (USD)', color='blue')
                    ax1.tick_params(axis='y', labelcolor='blue')
                    ax2 = ax1.twinx()
                    ax2.plot(combined_data.index, combined_data['Combined_Sentiment'], color='orange', label='Combined Sentiment')
                    ax2.set_ylabel('Sentiment Score (-1 to 1)', color='orange')
                    ax2.tick_params(axis='y', labelcolor='orange')
                    ax2.axhline(y=0, color='gray', linestyle='--')
                    plt.title(f"{self.coin_symbol} Price and Sentiment Over Time")
                    fig.tight_layout()
                    lines1, labels1 = ax1.get_legend_handles_labels()
                    lines2, labels2 = ax2.get_legend_handles_labels()
                    ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper left')
                    combined_chart_path = 'combined_price_sentiment_chart.png'
                    plt.savefig(combined_chart_path)
                    plt.close()
                    pdf_elements.append(Image(combined_chart_path, width=400, height=200))
                else:
                    pdf_elements.append(Paragraph("No overlapping price and sentiment data available.", styles['Normal']))
            else:
                pdf_elements.append(Paragraph("Combined Chart: Sentiment or price data missing.", styles['Normal']))
            pdf_elements.append(Spacer(1, 12))

            # Forecast Chart (PDF)
            pdf_elements.append(Paragraph("Price Forecast", styles['Heading2']))
            if self.forecasts.get('Combined') is not None and not self.forecasts['Combined'].empty:
                forecast_end_date = pd.to_datetime('2025-05-31')
                last_historical_date = self.data.index[-1]
                forecast_days = (forecast_end_date - last_historical_date).days + 1
                self.forecast_prices(days=forecast_days)
                fig, ax = plt.subplots(figsize=(10, 6))
                self.data['Close'].plot(ax=ax, label='Historical Price', color='blue')
                self.forecasts['Combined']['Forecasted_Price'].plot(ax=ax, label='Forecast', color='red', linestyle='--')
                ax.set_title(f"{self.coin_symbol} Price Forecast (Up to May 2025)")
                ax.set_xlabel('Date')
                ax.set_ylabel('Price (USD)')
                ax.legend()
                forecast_chart_path = 'forecast_chart.png'
                plt.savefig(forecast_chart_path)
                plt.close()
                pdf_elements.append(Image(forecast_chart_path, width=500, height=300))
            else:
                pdf_elements.append(Paragraph("Price Forecast: No data available.", styles['Normal']))
            pdf_elements.append(Spacer(1, 12))

            # Market Summary (PDF)
            pdf_elements.append(Paragraph("Market Summary", styles['Heading2']))
            current_price = self.data['Close'].iloc[-1]
            price_change = self.data['Close'].iloc[-1] - self.data['Close'].iloc[-2] if len(self.data) > 1 else 0
            price_change_pct = (price_change / self.data['Close'].iloc[-2]) * 100 if len(self.data) > 1 and self.data['Close'].iloc[-2] != 0 else 0
            pdf_elements.append(Paragraph(f"Current Price: ${current_price:.2f} USD", styles['Normal']))
            pdf_elements.append(Paragraph(f"24h Change: ${price_change:.2f} USD ({price_change_pct:.2f}%)", styles['Normal']))
            if 'Market Cap' in self.data.columns:
                pdf_elements.append(Paragraph(f"Market Cap: ${self.data['Market Cap'].iloc[-1]:.2f} USD", styles['Normal']))
            pdf_elements.append(Spacer(1, 12))

            # Sentiment Analysis (PDF)
            pdf_elements.append(Paragraph("Sentiment Analysis", styles['Heading2']))
            if self.sentiment_data is not None and not self.sentiment_data.empty:
                def get_last_sentiment(col):
                    idx = self.sentiment_data[col][self.sentiment_data[col] != 0].last_valid_index()
                    return self.sentiment_data[col].loc[idx] if idx is not None else 0

                reddit_sentiment = get_last_sentiment('Reddit_Sentiment')
                cryptopanic_sentiment = get_last_sentiment('CryptoPanic_Sentiment')
                cryptocompare_sentiment = get_last_sentiment('CryptoCompare_Sentiment')
                combined_sentiment = get_last_sentiment('Combined_Sentiment')

                sentiment_range = self.sentiment_data['Combined_Sentiment'].abs().max()
                threshold = max(0.1, sentiment_range * 0.3)
                sentiment_desc = 'Positive' if combined_sentiment > threshold else 'Negative' if combined_sentiment < -threshold else 'Neutral'

                pdf_elements.append(Paragraph(f"Market Sentiment: {sentiment_desc}", styles['Normal']))
                pdf_elements.append(Paragraph(f"Reddit Sentiment: {reddit_sentiment:.2f}", styles['Normal']))
                pdf_elements.append(Paragraph(f"CryptoPanic Sentiment: {cryptopanic_sentiment:.2f}", styles['Normal']))
                pdf_elements.append(Paragraph(f"CryptoCompare Sentiment: {cryptocompare_sentiment:.2f}", styles['Normal']))
                pdf_elements.append(Paragraph(f"Combined Sentiment: {combined_sentiment:.2f}", styles['Normal']))
                pdf_elements.append(Spacer(1, 12))

                fig, ax = plt.subplots(figsize=(8, 4))
                self.sentiment_data.plot(ax=ax)
                ax.set_title(f"{self.coin_symbol} Sentiment Analysis")
                ax.set_xlabel('Date')
                ax.set_ylabel('Sentiment Score')
                ax.axhline(y=0, color='r', linestyle='--')
                sentiment_chart_path = 'sentiment_chart.png'
                plt.savefig(sentiment_chart_path)
                plt.close()
                pdf_elements.append(Image(sentiment_chart_path, width=400, height=200))
            else:
                pdf_elements.append(Paragraph("Sentiment Analysis: No data available.", styles['Normal']))
            pdf_elements.append(Spacer(1, 12))

            # Event Analysis (PDF)
            pdf_elements.append(Paragraph("Event Analysis", styles['Heading2']))
            if self.event_summary_df is not None and not self.event_summary_df.empty:
                filtered_events = self.event_summary_df[self.event_summary_df['sentiment'] != 0]
                pdf_elements.append(Paragraph(f"Detected {len(filtered_events)} significant events with non-zero sentiment:", styles['Normal']))

                # Event table with sentiment-based coloring
                event_data = [['Date', 'Event Type', 'Price Change', 'Sentiment']]
                for _, event in filtered_events.iterrows():
                    sentiment_color = colors.green if event['sentiment'] > 0 else colors.red if event['sentiment'] < 0 else colors.grey
                    event_data.append([
                        event['date'].strftime('%Y-%m-%d'),
                        event['event_type'],
                        f"{event['price_change']*100:.2f}%",
                        f"{event['sentiment']:.2f}"
                    ])

                event_table = Table(event_data)
                event_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                    ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                    ('FONTSIZE', (0, 1), (-1, -1), 10),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('BACKGROUND', (0, 1), (0, -1), colors.lightgrey),
                    ('BACKGROUND', (1, 1), (1, -1), colors.lightgrey),
                    ('BACKGROUND', (2, 1), (2, -1), colors.lightgrey),
                    ('BACKGROUND', (3, 1), (3, -1), colors.white),
                ]))

                # Apply sentiment-based background color to the Sentiment column
                for i in range(1, len(event_data)):
                    sentiment_value = float(event_data[i][3])
                    color = colors.green if sentiment_value > 0 else colors.red if sentiment_value < 0 else colors.grey
                    event_table.setStyle(TableStyle([
                        ('BACKGROUND', (3, i), (3, i), color)
                    ]))

                pdf_elements.append(event_table)
            else:
                pdf_elements.append(Paragraph("Event Analysis: No significant events detected.", styles['Normal']))
            pdf_elements.append(Spacer(1, 12))

            # Upcoming Events (PDF)
            pdf_elements.append(Paragraph("Upcoming Events", styles['Heading2']))
            if upcoming_events:
                pdf_elements.append(Paragraph(f"Identified {len(upcoming_events)} upcoming events:", styles['Normal']))
                upcoming_data = [['Date', 'Event Name', 'Expected Impact', 'Source']]
                for event in upcoming_events:
                    upcoming_data.append([
                        event['date'].strftime('%Y-%m-%d'),
                        event['name'],
                        event['impact'],
                        event['source']
                    ])
                upcoming_table = Table(upcoming_data)
                upcoming_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                    ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                    ('FONTSIZE', (0, 1), (-1, -1), 10),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                pdf_elements.append(upcoming_table)
            else:
                pdf_elements.append(Paragraph("Upcoming Events: No upcoming events detected.", styles['Normal']))
            pdf_elements.append(Spacer(1, 12))

            # Investment Recommendation (PDF)
            pdf_elements.append(Paragraph("Investment Recommendations", styles['Heading2']))
            recommendation = await self._generate_investment_recommendation()
            pdf_elements.append(Paragraph("Recommendation: " + recommendation['action'], styles['Normal']))
            pdf_elements.append(Paragraph("Reasoning:", styles['Normal']))
            for reason in recommendation['reasoning']:
                pdf_elements.append(Paragraph(f"- {reason}", styles['Normal']))
            pdf_elements.append(Spacer(1, 12))

            # Build PDF
            doc.build(pdf_elements)
            self.logger.info(f"Report generated at {pdf_filepath}")

            for chart_file in ['price_chart.png', 'combined_price_sentiment_chart.png', 'forecast_chart.png', 'sentiment_chart.png']:
                if os.path.exists(chart_file):
                    os.remove(chart_file)

            return pdf_filepath

        except Exception as e:
            self.logger.error(f"Error generating PDF report: {e}")
            return None







    async def generate_pptx_report(self, ppt_filepath='crypto_presentation.pptx', upcoming_events=None):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            import os
            import requests
            from datetime import datetime
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            import matplotlib.pyplot as plt
            import pandas as pd

            def download_pixabay_image(api_key, query, output_path):
                """Download an image from Pixabay API."""
                url = f"https://pixabay.com/api/?key={api_key}&q={query}&image_type=photo&per_page=3"
                response = requests.get(url)
                if response.status_code == 200:
                    data = response.json()
                    if data['hits']:
                        image_url = data['hits'][0]['largeImageURL']
                        img_response = requests.get(image_url)
                        if img_response.status_code == 200:
                            with open(output_path, 'wb') as f:
                                f.write(img_response.content)
                            return True
                return False

            # Pixabay API key and background image path
            pixabay_api_key = "50210954-a7aca149c8b1407da7cd5e6bc"
            background_img_path = "photo.jfif"

            # PowerPoint Setup
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            def set_slide_background(slide, img_path):
                """Set background image for a PowerPoint slide."""
                if os.path.exists(img_path):
                    left = top = Inches(0)
                    slide.shapes.add_picture(img_path, left, top, width=Inches(13.333), height=Inches(7.5))

            def get_centered_position(object_width):
                """Calculate centered position for an object."""
                return (13.333 - object_width) / 2

            title_width = 11
            image_width = 7
            table_width = 6
            title_left = Inches(0)
            table_left = Inches(get_centered_position(table_width))

            # Apply background to all slides
            if os.path.exists(background_img_path):
                for slide in prs.slides._sldIdLst:
                    slide_layout = prs.slides._sldIdLst.index(slide)
                    slide = prs.slides[slide_layout]
                    set_slide_background(slide, background_img_path)
            else:
                self.logger.warning("Background image not found: photo.jfif")

            
            

            # Price and Sentiment Combined Chart Slide
            if self.sentiment_data is not None and not self.sentiment_data.empty and self.data is not None and not self.data.empty:
                combined_data = self.data[['Close']].join(self.sentiment_data[['Combined_Sentiment']], how='inner')
                if not combined_data.empty:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    set_slide_background(slide, background_img_path)
                    box = slide.shapes.add_textbox(title_left, Inches(0.7), Inches(title_width), Inches(1.5))
                    tf = box.text_frame
                    title = tf.add_paragraph()
                    title.text = "Price and Sentiment Combined Chart"
                    title.font.size = Pt(36)
                    title.font.bold = True
                    title.font.color.rgb = RGBColor(255, 255, 255)
                    try:
                        title.alignment = PP_ALIGN.CENTER
                    except:
                        title.alignment = 2
                    fig, ax1 = plt.subplots(figsize=(10, 6))
                    ax1.plot(combined_data.index, combined_data['Close'], color='blue', label='Price (USD)')
                    ax1.set_xlabel('Date')
                    ax1.set_ylabel('Price (USD)', color='blue')
                    ax1.tick_params(axis='y', labelcolor='blue')
                    ax2 = ax1.twinx()
                    ax2.plot(combined_data.index, combined_data['Combined_Sentiment'], color='orange', label='Combined Sentiment')
                    ax2.set_ylabel('Sentiment Score (-1 to 1)', color='orange')
                    ax2.tick_params(axis='y', labelcolor='orange')
                    ax2.axhline(y=0, color='gray', linestyle='--')
                    plt.title(f"{self.coin_symbol} Price and Sentiment Over Time")
                    fig.tight_layout()
                    lines1, labels1 = ax1.get_legend_handles_labels()
                    lines2, labels2 = ax2.get_legend_handles_labels()
                    ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper left')
                    combined_chart_path = 'combined_price_sentiment_chart.png'
                    plt.savefig(combined_chart_path)
                    plt.close()
                    slide.shapes.add_picture(combined_chart_path, Inches(get_centered_position(image_width)), Inches(2.3), width=Inches(image_width))
                    os.remove(combined_chart_path)
            else:
                self.logger.warning("No sentiment or price data available for combined chart.")

            
            # Market Summary Slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_background(slide, background_img_path)
            box = slide.shapes.add_textbox(title_left, Inches(0.7), Inches(title_width), Inches(1.5))
            tf = box.text_frame
            title = tf.add_paragraph()
            title.text = "Market Summary"
            title.font.size = Pt(36)
            title.font.bold = True
            title.font.color.rgb = RGBColor(255, 255, 255)
            try:
                title.alignment = PP_ALIGN.CENTER
            except:
                title.alignment = 2
            content_box = slide.shapes.add_textbox(title_left, Inches(2.5), Inches(title_width), Inches(3))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            if self.data is not None and not self.data.empty:
                current_price = self.data['Close'].iloc[-1]
                price_change = self.data['Close'].iloc[-1] - self.data['Close'].iloc[-2] if len(self.data) > 1 else 0
                price_change_pct = (price_change / self.data['Close'].iloc[-2]) * 100 if len(self.data) > 1 and self.data['Close'].iloc[-2] != 0 else 0
                text_lines = [f"• Current Price: ${current_price:.2f} USD",
                              f"• 24h Change: ${price_change:.2f} USD ({price_change_pct:.2f}%)"]
                if 'Market Cap' in self.data.columns:
                    text_lines.append(f"• Market Cap: ${self.data['Market Cap'].iloc[-1]:.2f} USD")
                for line in text_lines:
                    para = content_frame.add_paragraph()
                    para.text = line
                    para.font.size = Pt(24)
                    para.font.color.rgb = RGBColor(230, 230, 230)
                    para.space_after = Pt(20)
            else:
                self.logger.warning("No market data available for summary.")

            
            # Event Analysis Slide (Split across slides if needed)
            if self.event_summary_df is not None and not self.event_summary_df.empty:
                filtered_events = self.event_summary_df[self.event_summary_df['sentiment'] != 0]
                content_box = slide.shapes.add_textbox(title_left, Inches(2.5), Inches(title_width), Inches(0.5))
                content_frame = content_box.text_frame
                para = content_frame.add_paragraph()
                para.font.size = Pt(24)
                para.font.color.rgb = RGBColor(230, 230, 230)
                para.space_after = Pt(20)
                rows_per_slide = 10  # Adjust this number based on your slide height and row height
                for i in range(0, len(filtered_events), rows_per_slide):
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    set_slide_background(slide, background_img_path)
                    box = slide.shapes.add_textbox(title_left, Inches(0.7), Inches(title_width), Inches(1.5))
                    tf = box.text_frame
                    title = tf.add_paragraph()
                    title.text = "Event Analysis"
                    title.font.size = Pt(36)
                    title.font.bold = True
                    title.font.color.rgb = RGBColor(255, 255, 255)
                    try:
                        title.alignment = PP_ALIGN.CENTER
                    except:
                        title.alignment = 2
                    batch_events = filtered_events.iloc[i:i + rows_per_slide]
                    rows = len(batch_events) + 1
                    cols = 4
                    table = slide.shapes.add_table(rows, cols, table_left, Inches(2.5), Inches(table_width), Inches(rows * 0.25)).table
                    table.cell(0, 0).text = "Date"
                    table.cell(0, 1).text = "Event Type"
                    table.cell(0, 2).text = "Price Change"
                    table.cell(0, 3).text = "Sentiment"
                    for cell in table.rows[0].cells:
                        cell.text_frame.paragraphs[0].font.bold = True
                        cell.text_frame.paragraphs[0].font.size = Pt(18)
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                        try:
                            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        except:
                            cell.text_frame.paragraphs[0].alignment = 2
                    for j, (_, event) in enumerate(batch_events.iterrows()):
                        table.cell(j + 1, 0).text = event['date'].strftime('%Y-%m-%d')
                        table.cell(j + 1, 1).text = event['event_type']
                        table.cell(j + 1, 2).text = f"{event['price_change']*100:.2f}%"
                        sentiment = event['sentiment']
                        if sentiment > 0:
                            table.cell(j + 1, 3).text = f"{sentiment:.2f}"
                            table.cell(j + 1, 3).fill.solid()
                            table.cell(j + 1, 3).fill.fore_color.rgb = RGBColor(0, 128, 0)  # Green
                        elif sentiment < 0:
                            table.cell(j + 1, 3).text = f"{sentiment:.2f}"
                            table.cell(j + 1, 3).fill.solid()
                            table.cell(j + 1, 3).fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red
                        else:
                            table.cell(j + 1, 3).text = f"{sentiment:.2f}"
                            table.cell(j + 1, 3).fill.solid()
                            table.cell(j + 1, 3).fill.fore_color.rgb = RGBColor(128, 128, 128)  # Gray
                        for k in range(cols):
                            cell = table.cell(j + 1, k)
                            cell.text_frame.paragraphs[0].font.size = Pt(16)
                            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                            try:
                                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            except:
                                cell.text_frame.paragraphs[0].alignment = 2

            # Upcoming Events Slide (Split across slides if needed)
            if upcoming_events:
                content_box = slide.shapes.add_textbox(title_left, Inches(2.5), Inches(title_width), Inches(0.5))
                content_frame = content_box.text_frame
                para = content_frame.add_paragraph()
                para.font.size = Pt(24)
                para.font.color.rgb = RGBColor(230, 230, 230)
                para.space_after = Pt(20)
                upcoming_data = [['Date', 'Event Name', 'Expected Impact', 'event_type']]
                for event in upcoming_events:
                    upcoming_data.append([
                        event['date'].strftime('%Y-%m-%d'),
                        event['name'],
                        event['impact_classification'],
                        event['event_type']
                    ])
                rows_per_slide = 5  # Adjust this number based on your slide height and row height
                for i in range(0, len(upcoming_data), rows_per_slide):
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    set_slide_background(slide, background_img_path)
                    box = slide.shapes.add_textbox(title_left, Inches(0.7), Inches(title_width), Inches(1.5))
                    tf = box.text_frame
                    title = tf.add_paragraph()
                    title.text = "Upcoming Events"
                    title.font.size = Pt(36)
                    title.font.bold = True
                    title.font.color.rgb = RGBColor(255, 255, 255)
                    try:
                        title.alignment = PP_ALIGN.CENTER
                    except:
                        title.alignment = 2
                    batch_data = upcoming_data[i:i + rows_per_slide]
                    rows = len(batch_data)
                    cols = 4
                    table = slide.shapes.add_table(rows, cols, table_left, Inches(2.5), Inches(table_width), Inches(rows * 0.25)).table
                    for j, header in enumerate(batch_data[0]):
                        table.cell(0, j).text = header
                    for cell in table.rows[0].cells:
                        cell.text_frame.paragraphs[0].font.bold = True
                        cell.text_frame.paragraphs[0].font.size = Pt(18)
                        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                        try:
                            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        except:
                            cell.text_frame.paragraphs[0].alignment = 2
                    for j in range(1, rows):
                        for k in range(cols):
                            table.cell(j, k).text = str(batch_data[j][k])
                            cell = table.cell(j, k)
                            cell.text_frame.paragraphs[0].font.size = Pt(16)
                            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(217, 226, 243) if j % 2 == 0 else RGBColor(242, 242, 242)
                            try:
                                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            except:
                                cell.text_frame.paragraphs[0].alignment = 2

            # Investment Recommendation Slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_background(slide, background_img_path)
            box = slide.shapes.add_textbox(title_left, Inches(0.7), Inches(title_width), Inches(1.5))
            tf = box.text_frame
            title = tf.add_paragraph()
            title.text = "Investment Recommendations"
            title.font.size = Pt(36)
            title.font.bold = True
            title.font.color.rgb = RGBColor(255, 255, 255)
            try:
                title.alignment = PP_ALIGN.CENTER
            except:
                title.alignment = 2
            content_box = slide.shapes.add_textbox(title_left, Inches(1.2), Inches(title_width), Inches(2.1))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            recommendation = await self._generate_investment_recommendation()
            para = content_frame.add_paragraph()
            para.text = f"• Recommendation: {recommendation['action']}"
            para.font.size = Pt(24)
            para.font.color.rgb = RGBColor(230, 230, 230)
            para.space_after = Pt(20)
            para = content_frame.add_paragraph()
            para.text = "• Reasoning:"
            para.font.size = Pt(24)
            para.font.color.rgb = RGBColor(230, 230, 230)
            para.space_after = Pt(20)
            for reason in recommendation['reasoning']:
                para = content_frame.add_paragraph()
                para.text = f"  - {reason}"
                para.font.size = Pt(20)
                para.font.color.rgb = RGBColor(230, 230, 230)
                para.space_after = Pt(20)

            # Save PowerPoint
            prs.save(ppt_filepath)
            self.logger.info(f"Presentation generated at {ppt_filepath}")
            return ppt_filepath

        except Exception as e:
            self.logger.error(f"Error generating PPTX report: {e}")
            return None

















    async def _generate_investment_recommendation(self):
        try:
            action = "Hold"
            reasoning = []
            score = 0  # Score to determine recommendation (-100 to 100)

            # 1. Sentiment Analysis
            if self.sentiment_data is not None and not self.sentiment_data.empty:
                avg_sentiment = self.sentiment_data['Combined_Sentiment'].iloc[-10:].mean()
                sentiment_range = self.sentiment_data['Combined_Sentiment'].abs().max()
                threshold = max(0.1, sentiment_range * 0.3)
                sentiment_score = avg_sentiment * 40  # Sentiment contributes up to 40 points
                score += sentiment_score
                if avg_sentiment > threshold:
                    reasoning.append(f"Positive market sentiment (score: {avg_sentiment:.2f}) suggests potential upward momentum.")
                elif avg_sentiment < -threshold:
                    reasoning.append(f"Negative market sentiment (score: {avg_sentiment:.2f}) indicates potential downward pressure.")
                else:
                    reasoning.append(f"Neutral market sentiment (score: {avg_sentiment:.2f}) provides no strong directional signal.")

            # 2. Price Forecast
            if self.forecasts.get('Combined') is not None and not self.forecasts['Combined'].empty:
                current_price = self.data['Close'].iloc[-1]
                forecast_avg = self.forecasts['Combined']['Forecasted_Price'].mean()
                price_change_pct = ((forecast_avg - current_price) / current_price) * 100
                forecast_score = min(max(price_change_pct * 2, -40), 40)  # Forecast contributes up to 40 points
                score += forecast_score
                if price_change_pct > 5:
                    reasoning.append(f"Forecasted price increase of {price_change_pct:.2f}% suggests a potential buying opportunity.")
                elif price_change_pct < -5:
                    reasoning.append(f"Forecasted price decrease of {price_change_pct:.2f}% suggests caution or potential selling.")
                else:
                    reasoning.append(f"Forecasted price change of {price_change_pct:.2f}% indicates stability.")

            # 3. Upcoming Events (Updated to use impact_classification)
            upcoming_events = await self.check_upcoming_events(self.coin_symbol, forecast_days=30)
            if upcoming_events:
                good_events = sum(1 for event in upcoming_events if event.get('impact_classification') == 'good')
                bad_events = sum(1 for event in upcoming_events if event.get('impact_classification') == 'bad')
                neutral_events = sum(1 for event in upcoming_events if event.get('impact_classification') == 'neutral')
                unknown_events = sum(1 for event in upcoming_events if event.get('impact_classification') == 'unknown')

                # Adjust scoring based on event classification
                event_score = (good_events * 10) - (bad_events * 10) + (neutral_events * 2)
                score += event_score

                # Build reasoning based on event counts
                event_reasoning = []
                if good_events > 0:
                    event_reasoning.append(f"{good_events} upcoming positive events (e.g., protocol upgrades, listings) may drive price upward.")
                if bad_events > 0:
                    event_reasoning.append(f"{bad_events} upcoming negative events (e.g., low-impact holidays, minor unlocks) may exert downward pressure.")
                if neutral_events > 0:
                    event_reasoning.append(f"{neutral_events} upcoming neutral events (e.g., moderate-impact economic releases) may increase volatility.")
                if unknown_events > 0:
                    event_reasoning.append(f"{unknown_events} events with unknown impact add uncertainty.")
                if event_reasoning:
                    reasoning.extend(event_reasoning)
                else:
                    reasoning.append("Upcoming events have no clear directional impact.")
            else:
                reasoning.append("No significant upcoming events detected.")

            # 4. Technical Indicators
            if self.technical_indicators is not None and not self.technical_indicators.empty:
                latest_indicators = self.technical_indicators.iloc[-1]
                rsi = latest_indicators['RSI']
                macd = latest_indicators['MACD']
                macd_signal = latest_indicators['MACD_Signal']
                market_regime = latest_indicators['Market_Regime']

                # RSI Analysis
                if rsi > 70:
                    score -= 20
                    reasoning.append(f"RSI ({rsi:.2f}) indicates overbought conditions, suggesting a potential pullback.")
                elif rsi < 30:
                    score += 20
                    reasoning.append(f"RSI ({rsi:.2f}) indicates oversold conditions, suggesting a potential buying opportunity.")
                else:
                    reasoning.append(f"RSI ({rsi:.2f}) is neutral, indicating no extreme conditions.")

                # MACD Analysis
                if macd > macd_signal and macd > 0:
                    score += 20
                    reasoning.append("MACD is above the signal line and positive, indicating bullish momentum.")
                elif macd < macd_signal and macd < 0:
                    score -= 20
                    reasoning.append("MACD is below the signal line and negative, indicating bearish momentum.")
                else:
                    reasoning.append("MACD shows no strong directional signal.")

                reasoning.append(f"Market Regime: {market_regime}")

            # Determine Recommendation
            if score > 50:
                action = "Buy"
                reasoning.append(f"Overall score ({score:.2f}) is strongly positive, recommending a Buy.")
            elif score < -50:
                action = "Sell"
                reasoning.append(f"Overall score ({score:.2f}) is strongly negative, recommending a Sell.")
            else:
                action = "Hold"
                reasoning.append(f"Overall score ({score:.2f}) is neutral, recommending a Hold.")

            return {
                'action': action,
                'reasoning': reasoning,
                'score': score
            }

        except Exception as e:
            print(f"[ERROR] Error generating investment recommendation: {e}")
            return {
                'action': "Hold",
                'reasoning': ["Unable to generate recommendation due to insufficient data or errors."],
                'score': 0
            }





    async def run_analysis(
        self,
        question,
        start_date=None,
        end_date=None,
        output_dir='plots',
        csv_output='crypto_analysis_export.csv',
        pdf_output='crypto_report.pdf',
        ppt_output='crypto_presentation.pptx'
    ):
        try:
            # Step 1: Set default dates if not provided
            if end_date is None:
                end_date = datetime.now()
            if start_date is None:
                start_date = end_date - timedelta(days=90)

            self.logger.info(f"Running analysis for question: '{question}'")
            self.logger.info(f"Date range: {start_date} to {end_date}")

            # Step 2: Detect the coin from the question
            coin = self._detect_coin_from_question(question)
            if not coin:
                self.logger.error("Could not detect a cryptocurrency in the question.")
                return {"error": "Could not detect a cryptocurrency in the question."}

            file_path = self.coin_csv_mapping.get(coin)
            if not file_path:
                self.logger.error(f"No data file found for {coin}.")
                return {"error": f"No data file found for {coin}."}

            # Step 3: Load historical data
            self.load_data_from_csv(file_path)
            if self.data is None or self.data.empty:
                self.logger.error(f"Failed to load data for {coin}.")
                return {"error": f"Failed to load data for {coin}."}

            # Step 4: Perform sentiment analysis
            async with aiohttp.ClientSession() as session:
                sentiment_result = await self.analyze_sentiment(start_date, end_date, session)
                if sentiment_result is None:
                    self.logger.warning("Sentiment analysis failed. Proceeding without sentiment data.")

            # Step 5: Detect events
            event_result = await self.detect_events()
            if event_result is None:
                self.logger.warning("Event detection failed. Proceeding without event data.")

            # Step 6: Calculate technical indicators
            tech_indicators = self.calculate_technical_indicators()
            if tech_indicators is None:
                self.logger.warning("Technical indicator calculation failed. Proceeding without indicators.")

            # Step 7: Generate price forecasts
            forecast_result = self.forecast_prices(days=60)
            if forecast_result is None:
                self.logger.warning("Price forecasting failed. Proceeding without forecasts.")

            # Step 8: Check upcoming events ONCE and store the result
            self.upcoming_events = await self.check_upcoming_events(self.coin_symbol, forecast_days=30)
            if not self.upcoming_events:
                self.logger.info("No upcoming events detected.")

            # Step 9: Set up RAG system
            rag_result = self.setup_rag_system()
            if rag_result is None:
                self.logger.warning("Failed to set up RAG system.")

            # Step 10: Generate interactive plots
            plot_result = self.generate_interactive_plots(output_dir=output_dir)
            if not plot_result:
                self.logger.warning("Failed to generate interactive plots.")

            # Step 11: Export analysis to CSV
            csv_result = self.export_analysis_to_csv(output_path=csv_output)
            if not csv_result:
                self.logger.warning("Failed to export analysis to CSV.")

            # Step 12: Generate PDF and PPT report using stored upcoming events
            pdf_path = await self.generate_pdf_report(pdf_filepath=pdf_output, upcoming_events=self.upcoming_events)
            ppt_path = await self.generate_pptx_report(ppt_filepath=ppt_output, upcoming_events=self.upcoming_events)
            if not pdf_path or not ppt_path:
                self.logger.warning("Failed to generate PDF or PPT report.")

            # Step 13: Summarize the results
            summary = {
                "coin_name": self.coin_name,
                "coin_symbol": self.coin_symbol,
                "date_range": f"{start_date.strftime('%Y-%m-%d')} to {end_date.strftime('%Y-%m-%d')}",
                "sentiment_completed": self.sentiment_data is not None and not self.sentiment_data.empty,
                "events_detected": len(self.event_summary_df) if self.event_summary_df is not None else 0,
                "upcoming_events": len(self.upcoming_events) if self.upcoming_events else 0,
                "forecast_completed": self.forecasts.get('Combined') is not None and not self.forecasts['Combined'].empty,
                "technical_indicators_completed": self.technical_indicators is not None and not self.technical_indicators.empty,
                "plots_generated": plot_result is not None,
                "csv_exported": csv_result is not None,
                "pdf_generated": pdf_path is not None,
                "ppt_generated": ppt_path is not None,
            }

            # Add investment recommendation to summary
            recommendation = await self._generate_investment_recommendation()
            summary["recommendation"] = {
                "action": recommendation['action'],
                "reasoning": recommendation['reasoning'],
                "score": recommendation['score']
            }

            self.logger.info(f"Analysis completed successfully for {self.coin_symbol}.")
            self.logger.info("Summary:")
            for key, value in summary.items():
                if key != "recommendation":
                    self.logger.info(f"  {key}: {value}")
            self.logger.info("  recommendation:")
            self.logger.info(f"    action: {summary['recommendation']['action']}")
            for reason in summary['recommendation']['reasoning']:
                self.logger.info(f"      - {reason}")
            self.logger.info(f"    score: {summary['recommendation']['score']}")

            return summary

        except Exception as e:
            self.logger.error(f"Error running analysis: {e}")
            return {"error": str(e)}

"""ahmed'scode

"""

import os
import glob
import sys
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from prophet import Prophet

# Additional imports for handling file permission issues
from datetime import datetime
import time
import random

# Make psutil optional
try:
    import psutil
except ImportError:
    # Simply continue without psutil - we'll handle this gracefully
    pass

class CryptoForecastGenerator:
    """
    A class to generate cryptocurrency price forecasts and create PowerPoint presentations.
    """

    def __init__(self):
        """Initialize the crypto forecast generator."""
        # Constants
        self.FUTURE_DAYS = 100
        self.PLOT_DPI = 300
        self.SLIDE_WIDTH = 13.333  # inches
        self.SLIDE_HEIGHT = 7.5    # inches

        # Ensure necessary directories exist
        os.makedirs("plots", exist_ok=True)
        os.makedirs("pptx", exist_ok=True)

        # Paths
        self.search_dir = "Data"
        

        self.background_img = 'photo.jfif'
        self.forecast_plot = os.path.join('plots', 'forecast_plot.png')
        self.indicators_plot = os.path.join('plots', 'indicators.png')

        # PowerPoint settings
        self.title_width = 11      # inches
        self.image_width = 7       # inches for charts (changed from 10 to 7)
        self.table_width = 9       # inches for tables

        # Table colors (matching the background theme)
        self.header_color = RGBColor(68, 114, 196)     # Medium blue (#4472C4)
        self.row_color1 = RGBColor(217, 226, 243)      # Light blue (#D9E2F3)
        self.row_color2 = RGBColor(242, 242, 242)      # Very light blue/white (#F2F2F2)
        self.header_text_color = RGBColor(255, 255, 255)  # White
        self.row_text_color = RGBColor(0, 0, 0)        # Black

        # Initialize other attributes
        self.df = None
        self.forecast = None
        self.coin_name = None
        self.matched_file = None

    def get_centered_position(self, object_width):
        """Calculate centered position for an object."""
        return (self.SLIDE_WIDTH - object_width) / 2

    def find_crypto_file(self, keyword):
        """Find a cryptocurrency CSV file based on a keyword."""
        if not os.path.exists(self.search_dir):
            os.makedirs(self.search_dir, exist_ok=True)
            print(f"⚠️ Directory '{self.search_dir}' created. Please place your CSV files there.")
            raise FileNotFoundError(f"No CSV files found in {self.search_dir}")

        csv_files = glob.glob(os.path.join(self.search_dir, "*.csv"))
        if not csv_files:
            raise FileNotFoundError(f"No CSV files found in {self.search_dir}")

        matched_file = None
        for f in csv_files:
            if keyword.lower() in f.lower():
                matched_file = f
                break

        if not matched_file:
            raise FileNotFoundError(f"No file found for keyword '{keyword}' in {self.search_dir}")

        self.matched_file = matched_file
        self.coin_name = os.path.basename(matched_file).replace(" Historical Data.csv", "")
        print(f"✅ Using file: {matched_file}")
        return matched_file

    def load_and_preprocess_data(self):
        """Load and preprocess cryptocurrency data."""
        try:
            df = pd.read_csv(self.matched_file)

            # Clean column names
            df.columns = [col.strip() for col in df.columns]

            # Parse date
            df['Date'] = pd.to_datetime(df['Date'], errors='coerce')
            df = df.dropna(subset=['Date']).sort_values('Date')

            # Remove $ and commas then convert to numeric
            numeric_cols = ['Open', 'High', 'Low', 'Close', 'Volume']
            # Add Market Cap if it exists
            if 'Market Cap' in df.columns:
                numeric_cols.append('Market Cap')

            for col in numeric_cols:
                if col in df.columns:
                    df[col] = df[col].replace({'\$': '', ',': ''}, regex=True)
                    df[col] = pd.to_numeric(df[col], errors='coerce')

            # Rename Close to Price to match old naming if needed
            if 'Close' in df.columns and 'Price' not in df.columns:
                df.rename(columns={'Close': 'Price'}, inplace=True)

            # Extract Coin Name from the first row if available
            if 'Coin Name' in df.columns:
                self.coin_name = df['Coin Name'].iloc[0]

            # Keep only necessary columns
            required_cols = ['Date', 'Price', 'Open', 'High', 'Low']
            # Add Volume if it exists
            if 'Volume' in df.columns:
                required_cols.append('Volume')

            # Check if all required columns are present
            missing_cols = [col for col in required_cols if col not in df.columns]
            if missing_cols:
                raise ValueError(f"Missing required columns: {', '.join(missing_cols)}")

            df = df[required_cols].dropna(subset=['Price'])

            # Add Volume if missing (with zeros)
            if 'Volume' not in df.columns:
                df['Volume'] = 0

            self.df = df
            return df

        except Exception as e:
            print(f"⚠️ Error loading data: {str(e)}")
            raise

    def compute_rsi(self, series, period=14):
        """Compute Relative Strength Index for a price series."""
        delta = series.diff()
        gain = delta.where(delta > 0, 0)
        loss = -delta.where(delta < 0, 0)
        avg_gain = gain.rolling(window=period).mean()
        avg_loss = loss.rolling(window=period).mean()
        rs = avg_gain / (avg_loss + 1e-6)  # Adding small epsilon to avoid division by zero
        return 100 - (100 / (1 + rs))

    def generate_features(self):
        """Generate technical features from price data."""
        df = self.df

        # Price-based features
        df['Return'] = df['Price'].pct_change()
        df['LogReturn'] = np.log(df['Price'] / df['Price'].shift(1))
        df['MA_7'] = df['Price'].rolling(7).mean()
        df['MA_30'] = df['Price'].rolling(30).mean()
        df['STD_7'] = df['Price'].rolling(7).std()

        # Momentum and volatility
        df['Momentum_10'] = df['Price'] - df['Price'].shift(10)
        df['Volatility_7d'] = df['Return'].rolling(7).std()

        # Candlestick features
        df['Price_Range'] = df['High'] - df['Low']
        df['Candle_Body'] = abs(df['Price'] - df['Open'])
        df['Upper_Shadow'] = df['High'] - df[['Price', 'Open']].max(axis=1)
        df['Lower_Shadow'] = df[['Price', 'Open']].min(axis=1) - df['Low']

        # Volume-based features
        df['Volume_Change'] = df['Volume'].pct_change()
        df['Volatility_to_Volume'] = df['Volatility_7d'] / (df['Volume'] + 1e-6)

        # Technical indicators
        df['RSI_14'] = self.compute_rsi(df['Price'], 14)

        # Lag features
        df['Lag_1'] = df['Price'].shift(1)
        df['Lag_7'] = df['Price'].shift(7)

        # Ratio features
        df['Price_to_MA7'] = df['Price'] / (df['MA_7'] + 1e-6)
        df['Momentum_to_Volatility'] = df['Momentum_10'] / (df['Volatility_7d'] + 1e-6)

        self.df = df
        return df

    def create_forecast(self):
        """Create price forecast using Prophet model with log-transform to avoid negative prices."""
        try:
        # Prepare data for Prophet with log transformation
            df_prophet = self.df[['Date', 'Price']].dropna().copy()
            df_prophet = df_prophet[df_prophet['Price'] > 0]  # Ensure no non-positive values
            df_prophet['y'] = np.log(df_prophet['Price'])
            df_prophet.rename(columns={'Date': 'ds'}, inplace=True)

        # Train the Prophet model
            print("Training forecasting model (log scale)...")
            model = Prophet(daily_seasonality=True)
            model.fit(df_prophet)

        # Forecast
            future = model.make_future_dataframe(periods=self.FUTURE_DAYS)
            forecast = model.predict(future)

        # Exponentiate predictions to return to original scale
            forecast['yhat'] = np.exp(forecast['yhat'])
            forecast['yhat_lower'] = np.exp(forecast['yhat_lower'])
            forecast['yhat_upper'] = np.exp(forecast['yhat_upper'])

            self.forecast = forecast

        # Plot forecast
            print("Generating forecast plot (back-transformed)...")
            fig1, ax = plt.subplots(figsize=(10, 6))
            ax.plot(df_prophet['ds'], np.exp(df_prophet['y']), label='Observed')
            ax.plot(forecast['ds'], forecast['yhat'], label='Forecast')
            ax.fill_between(forecast['ds'], forecast['yhat_lower'], forecast['yhat_upper'], alpha=0.3)
            ax.set_title('Forecasted Price')
            ax.grid(True, alpha=0.3)
            plt.tight_layout()
            plt.savefig(self.forecast_plot, dpi=self.PLOT_DPI)
            plt.close()

            return forecast

        except Exception as e:
            print(f"⚠️ Error creating forecast: {str(e)}")
            raise

    def create_indicator_plots(self):
        """Create technical indicator plots."""
        try:
            print("Generating technical indicators plot...")
            # Create a 2x2 subplot of key indicators
            fig2, axs = plt.subplots(2, 2, figsize=(12, 8))

            # Plot each indicator
            axs[0, 0].plot(self.df['Date'], self.df['Price'], color='#1f77b4')
            axs[0, 0].set_title("Price")
            axs[0, 0].grid(True, alpha=0.3)

            axs[0, 1].plot(self.df['Date'], self.df['Volatility_7d'], color='#ff7f0e')
            axs[0, 1].set_title("Volatility (7d)")
            axs[0, 1].grid(True, alpha=0.3)

            axs[1, 0].plot(self.df['Date'], self.df['RSI_14'], color='#2ca02c')
            axs[1, 0].set_title("RSI 14")
            axs[1, 0].axhline(y=70, color='r', linestyle='--', alpha=0.5)  # Overbought line
            axs[1, 0].axhline(y=30, color='g', linestyle='--', alpha=0.5)  # Oversold line
            axs[1, 0].grid(True, alpha=0.3)

            axs[1, 1].plot(self.df['Date'], self.df['Momentum_10'], color='#d62728')
            axs[1, 1].set_title("Momentum 10")
            axs[1, 1].grid(True, alpha=0.3)

            plt.tight_layout()
            plt.savefig(self.indicators_plot, bbox_inches='tight', dpi=self.PLOT_DPI)
            plt.close()

        except Exception as e:
            print(f"⚠️ Error creating indicator plots: {str(e)}")
            raise

    def analyze_risk(self):
        """Analyze risk factors from the data and forecast."""
        try:
            # Volatility analysis
            recent_volatility = self.df['Volatility_7d'].iloc[-30:].mean()
            if recent_volatility > 0.05:
                volatility_comment = "⚠️ High volatility detected in the past month."
            else:
                volatility_comment = "✅ Volatility is within a moderate range."

            # RSI analysis
            latest_rsi = self.df['RSI_14'].iloc[-1]
            if latest_rsi > 70:
                rsi_comment = "⚠️ RSI indicates potential overbought conditions."
            elif latest_rsi < 30:
                rsi_comment = "⚠️ RSI indicates potential oversold conditions."
            else:
                rsi_comment = "✅ RSI is in a neutral state."

            # Forecast uncertainty analysis
            forecast_range = self.forecast[['yhat_upper', 'yhat_lower']].tail(self.FUTURE_DAYS)
            avg_uncertainty = (forecast_range['yhat_upper'] - forecast_range['yhat_lower']).mean()
            if avg_uncertainty > self.df['Price'].mean() * 0.2:
                forecast_comment = "⚠️ Forecast has wide uncertainty bands."
            else:
                forecast_comment = "✅ Forecast shows reasonable confidence intervals."

            return {
                'volatility': volatility_comment,
                'rsi': rsi_comment,
                'forecast': forecast_comment
            }

        except Exception as e:
            print(f"⚠️ Error in risk analysis: {str(e)}")
            return {
                'volatility': "⚠️ Could not analyze volatility.",
                'rsi': "⚠️ Could not analyze RSI.",
                'forecast': "⚠️ Could not analyze forecast uncertainty."
            }

    def set_slide_background(self, slide, img_path):
        """Set background image for a PowerPoint slide."""
        if os.path.exists(img_path):
            left = top = Inches(0)
            slide.shapes.add_picture(img_path, left, top, width=Inches(self.SLIDE_WIDTH), height=Inches(self.SLIDE_HEIGHT))
        else:
            print(f"⚠️ Background image not found: {img_path}")

    def create_presentation(self, risk_analysis):
        """Create a PowerPoint presentation with the forecast results."""
        try:
            print("Creating PowerPoint presentation...")
            ppt = Presentation()
            ppt.slide_width = Inches(self.SLIDE_WIDTH)
            ppt.slide_height = Inches(self.SLIDE_HEIGHT)

            # Calculate positions
            title_left = Inches(self.get_centered_position(self.title_width))
            image_left = Inches(self.get_centered_position(self.image_width))
            table_left = Inches(self.get_centered_position(self.table_width))

            # Slide 1: Title Slide
            title_slide = ppt.slides.add_slide(ppt.slide_layouts[6])  # blank layout
            if os.path.exists(self.background_img):
                self.set_slide_background(title_slide, self.background_img)

            # Add title
            box = title_slide.shapes.add_textbox(title_left, Inches(0.7), Inches(self.title_width), Inches(1.5))
            tf = box.text_frame
            title = tf.add_paragraph()
            title.text = f"Crypto Forecast Report: {self.coin_name}"
            title.font.size = Pt(40)
            title.font.bold = True
            title.font.color.rgb = RGBColor(255, 255, 255)
            try:
                title.alignment = PP_ALIGN.CENTER
            except:
                title.alignment = 2  # Center alignment fallback

            # Add subtitle
            sub_box = title_slide.shapes.add_textbox(title_left, Inches(2), Inches(self.title_width), Inches(1))
            sf = sub_box.text_frame
            subtitle = sf.add_paragraph()
            subtitle.text = f"Forecast Horizon: {self.FUTURE_DAYS} Days"
            subtitle.font.size = Pt(28)
            subtitle.font.color.rgb = RGBColor(200, 200, 200)
            try:
                subtitle.alignment = PP_ALIGN.CENTER
            except:
                subtitle.alignment = 2  # Center alignment fallback

            # Slide 2: Risk Assessment
            risk_slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            if os.path.exists(self.background_img):
                self.set_slide_background(risk_slide, self.background_img)

            # Add title
            box = risk_slide.shapes.add_textbox(title_left, Inches(0.7), Inches(self.title_width), Inches(1.5))
            tf = box.text_frame
            title = tf.add_paragraph()
            title.text = "📋 Risk Assessment Summary"
            title.font.size = Pt(36)
            title.font.bold = True
            title.font.color.rgb = RGBColor(255, 255, 255)
            try:
                title.alignment = PP_ALIGN.CENTER
            except:
                title.alignment = 2  # Center alignment fallback

            # Add bullet points
            content_box = risk_slide.shapes.add_textbox(title_left, Inches(2.5), Inches(self.title_width), Inches(3))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for point in [risk_analysis['volatility'], risk_analysis['rsi'], risk_analysis['forecast']]:
                para = content_frame.add_paragraph()
                para.text = "• " + point
                para.font.size = Pt(24)
                para.font.color.rgb = RGBColor(230, 230, 230)
                para.space_after = Pt(20)

            # Slide 3: Forecast Plot
            forecast_slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            if os.path.exists(self.background_img):
                self.set_slide_background(forecast_slide, self.background_img)

            # Add title
            box = forecast_slide.shapes.add_textbox(title_left, Inches(0.7), Inches(self.title_width), Inches(1.5))
            tf = box.text_frame
            title = tf.add_paragraph()
            title.text = "Price Forecast with Confidence Intervals"
            title.font.size = Pt(36)
            title.font.bold = True
            title.font.color.rgb = RGBColor(255, 255, 255)
            try:
                title.alignment = PP_ALIGN.CENTER
            except:
                title.alignment = 2  # Center alignment fallback

            # Add plot - CENTERED
            if os.path.exists(self.forecast_plot):
                forecast_slide.shapes.add_picture(
                    self.forecast_plot,
                    image_left,
                    Inches(2.3),
                    width=Inches(self.image_width)
                )
            else:
                print(f"⚠️ Forecast plot not found: {self.forecast_plot}")

            # Slide 4: Key Indicators
            indicators_slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            if os.path.exists(self.background_img):
                self.set_slide_background(indicators_slide, self.background_img)

            # Add title
            box = indicators_slide.shapes.add_textbox(title_left, Inches(0.7), Inches(self.title_width), Inches(1.5))
            tf = box.text_frame
            title = tf.add_paragraph()
            title.text = "Key Technical Indicators"
            title.font.size = Pt(36)
            title.font.bold = True
            title.font.color.rgb = RGBColor(255, 255, 255)
            try:
                title.alignment = PP_ALIGN.CENTER
            except:
                title.alignment = 2  # Center alignment fallback

            # Add plot - CENTERED
            if os.path.exists(self.indicators_plot):
                indicators_slide.shapes.add_picture(
                    self.indicators_plot,
                    image_left,
                    Inches(2.3),
                    width=Inches(self.image_width)
                )
            else:
                print(f"⚠️ Indicators plot not found: {self.indicators_plot}")

            # Slide 5: Forecast Table
            forecast_table_slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            if os.path.exists(self.background_img):
                self.set_slide_background(forecast_table_slide, self.background_img)

            # Add title
            box = forecast_table_slide.shapes.add_textbox(title_left, Inches(0.7), Inches(self.title_width), Inches(1.5))
            tf = box.text_frame
            title = tf.add_paragraph()
            title.text = "Detailed Forecast Table"
            title.font.size = Pt(36)
            title.font.bold = True
            title.font.color.rgb = RGBColor(255, 255, 255)
            try:
                title.alignment = PP_ALIGN.CENTER
            except:
                title.alignment = 2  # Center alignment fallback

            # Get last 10 days of forecast
            last_forecast = self.forecast[['ds', 'yhat', 'yhat_lower', 'yhat_upper']].tail(10)

            # Add a proper PowerPoint table
            table_rows = 1 + len(last_forecast)  # Header + data rows
            table_cols = 4  # Date, Forecast, Low, High

            try:
                table = forecast_table_slide.shapes.add_table(
                    table_rows,
                    table_cols,
                    table_left,
                    Inches(2.3),
                    Inches(self.table_width),
                    Inches(4.4)  # Height based on rows
                ).table

                # Set column widths (equal widths) - Convert to integer
                col_width = int(Inches(self.table_width / table_cols))
                for col in range(table_cols):
                    table.columns[col].width = col_width

                # Set header row
                table.cell(0, 0).text = "Date"
                table.cell(0, 1).text = "Forecast"
                table.cell(0, 2).text = "Low"
                table.cell(0, 3).text = "High"

                # Format header row - Blue background with white text
                for col in range(table_cols):
                    cell = table.cell(0, col)
                    # Set background color (blue)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.header_color

                    # Format text
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(20)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = self.header_text_color
                    try:
                        paragraph.alignment = PP_ALIGN.CENTER
                    except:
                        paragraph.alignment = 2  # CENTER value

                # Fill in the data rows with alternating colors
                for i, (_, row) in enumerate(last_forecast.iterrows(), start=1):
                    # Set text
                    table.cell(i, 0).text = str(row['ds'].date())
                    table.cell(i, 1).text = f"{row['yhat']:.2f}"
                    table.cell(i, 2).text = f"{row['yhat_lower']:.2f}"
                    table.cell(i, 3).text = f"{row['yhat_upper']:.2f}"

                    # Set row background color (alternating)
                    row_color = self.row_color1 if i % 2 == 1 else self.row_color2

                    # Format data cells
                    for col in range(table_cols):
                        cell = table.cell(i, col)

                        # Set background color for cell
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = row_color

                        # Format text
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Pt(18)
                        paragraph.font.color.rgb = self.row_text_color
                        try:
                            paragraph.alignment = PP_ALIGN.CENTER
                        except:
                            paragraph.alignment = 2  # CENTER value

            except Exception as e:
                print(f"⚠️ Error creating table: {str(e)}")
                # Fallback to text box if table fails
                table_box = forecast_table_slide.shapes.add_textbox(
                    table_left,
                    Inches(2.3),
                    Inches(self.table_width),
                    Inches(4)
                )
                tf = table_box.text_frame
                tf.word_wrap = True

                # Add header
                header = tf.add_paragraph()
                header.text = "Date          Forecast       Low        High"
                header.font.size = Pt(20)
                header.font.bold = True
                header.font.color.rgb = RGBColor(255, 255, 255)
                header.space_after = Pt(12)
                try:
                    header.alignment = PP_ALIGN.CENTER
                except:
                    header.alignment = 2  # CENTER value

                # Add data rows
                for _, row in last_forecast.iterrows():
                    data_row = tf.add_paragraph()
                    data_row.text = f"{row['ds'].date()}   {row['yhat']:.2f}   {row['yhat_lower']:.2f}   {row['yhat_upper']:.2f}"
                    data_row.font.size = Pt(18)
                    data_row.font.color.rgb = RGBColor(230, 230, 230)
                    data_row.space_after = Pt(6)
                    try:
                        data_row.alignment = PP_ALIGN.CENTER
                    except:
                        data_row.alignment = 2  # CENTER value

            # Save the presentation with error handling for file access issues
            output_pptx = os.path.join('pptx', f'{self.coin_name}_forecast_report.pptx')

            try:
                ppt.save(output_pptx)
                print(f"\n✅ Analysis complete! Presentation saved as: {output_pptx}")
                return output_pptx
            except PermissionError:
                # Try to create a new filename with timestamp if file is in use
                from datetime import datetime
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                alt_output_pptx = os.path.join('pptx', f'{self.coin_name}_forecast_report_{timestamp}.pptx')

                try:
                    ppt.save(alt_output_pptx)
                    print(f"\n✅ Analysis complete! Original file was in use.")
                    print(f"✅ Presentation saved with alternate name: {alt_output_pptx}")
                    return alt_output_pptx
                except Exception as e:
                    print(f"\n⚠️ Could not save presentation: The file may be open in PowerPoint.")
                    print(f"⚠️ Please close any open presentations and try again, or specify a different output location.")
                    return None

        except Exception as e:
            print(f"⚠️ Error creating presentation: {str(e)}")
            raise

    def run(self, coin_keyword):
        """Run the full forecasting and presentation creation workflow."""
        try:
            print("\n🔍 Searching for cryptocurrency data...")
            self.find_crypto_file(coin_keyword)

            print("📊 Loading and preprocessing data...")
            self.load_and_preprocess_data()

            print("🔧 Generating technical features...")
            self.generate_features()

            print("🔮 Creating price forecast...")
            self.create_forecast()

            print("📈 Creating technical indicator plots...")
            self.create_indicator_plots()

            print("⚖️ Analyzing risk factors...")
            risk_analysis = self.analyze_risk()

            print("📑 Creating presentation...")
            output_file = self.create_presentation(risk_analysis)

            return output_file

        except Exception as e:
            print(f"\n❌ Error: {str(e)}")
            return None

"""combined main"""
import os
import sys
from datetime import datetime
import asyncio
import nest_asyncio
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
import psutil  # Added import for PowerPoint conflict check

# Apply nest_asyncio to allow nested event loops (e.g., in Jupyter/Colab)
nest_asyncio.apply()

# FastAPI app instance

# Pydantic model for request body
class CryptoRequest(BaseModel):
    coin: str

# API keys for CryptoAnalyzer (unchanged)
api_keys = {
    'reddit_client_id': '4EH28GUsWtOmcm7ln0gtLA',
    'reddit_client_secret': '3KCFkazbnAH3A_HEN1hsqZoE0mSZrw',
    'reddit_user_agent': 'crypto/0.1 by Pale_Swim_2351',
    'cryptopanic_api_key': 'cc98ee6f0e802dc8c7d475742d08834559226548',
    'cryptocompare_api_key': '9540c5846bfa2ae62669f08c0666c9e4f43ee76691531e522e940e0b9014f79a',
    'openrouter': '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',
    'github_token': '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',
    'coinmarketcap_api_key': '85a8f2a3-bc07-4623-89d2-884c9eda7085',
    'alpha_vantage_api_key': 'P7USX248W5MZ4J0K',
    'fmp_api_key': 'TofKt59yYDinAR8jjxkV15slEy3iD2H7',
    'fred_api_key': '1edfbd114c38f0b2601ca409d9df26ea',
    'coingecko_api_key': 'CG-rJ8KFdXvwAawn3eFSJ1bkVVp'
}

# Initialize instances
analyzer = CryptoAnalyzer(api_keys=api_keys)
generator = CryptoForecastGenerator()

@app.post("/analyze-crypto/")
async def analyze_crypto(request: CryptoRequest):
    """Endpoint to run cryptocurrency analysis and forecasting."""
    try:
        # Extract coin from request
        coin = request.coin.strip()
        if not coin:
            raise HTTPException(status_code=400, detail="Coin name cannot be empty.")

        # Print welcome message
        print("\n" + "="*50)
        print("🚀 Combined Cryptocurrency Analysis & Forecast Report Generator 🚀")
        print("="*50)
        print(f"\nProcessing cryptocurrency: {coin}")

        # Define inputs for both tasks using the single input
        coin_keyword = coin  # For CryptoForecastGenerator
        question = f"Analyze {coin}"  # For CryptoAnalyzer

        # Define date range (adjust end_date to be current or future for real-time data)
        start_date = datetime(2025, 1, 1)
        end_date = datetime.now()  # Current date: May 17, 2025

        # Check for PowerPoint conflicts
        powerpoint_warning_shown = False
        powerpoint_running = False
        for proc in psutil.process_iter(['name']):
            if proc.info['name'] and 'powerpnt' in proc.info['name'].lower():
                powerpoint_running = True
                print("\n⚠️ PowerPoint appears to be running.")
                print("⚠️ If you're editing the same presentation file, please save and close it first.")
                powerpoint_warning_shown = True
                break

        if not powerpoint_warning_shown:
            print("\nℹ️ Note: If PowerPoint is open, please ensure you're not editing the output file.")

        # Run both tasks concurrently
        print("\n🔍 Starting concurrent analysis and forecast generation...")
        forecast_task = asyncio.to_thread(generator.run, coin_keyword)  # Run synchronous run() in a thread
        analysis_task = analyzer.run_analysis(
            question=question,
            start_date=start_date,
            end_date=end_date,
            output_dir='plots',
            csv_output='crypto_analysis_export.csv',
            pdf_output='crypto_report.pdf',
            ppt_output='combined_crypto_report.pptx'
        )

        # Wait for both tasks to complete
        forecast_result, analysis_result = await asyncio.gather(forecast_task, analysis_task)

        # Process results
        result = {}
        if forecast_result:
            print(f"\n🎉 Forecast presentation generated: {forecast_result}")
            result["forecast_presentation"] = forecast_result
        else:
            print("\n⚠️ Forecast generation failed.")
            result["forecast_presentation"] = None

        if analysis_result and not analysis_result.get("error"):
            print(f"\n🎉 Analysis presentation generated: combined_crypto_report.pptx")
            result["analysis_presentation"] = "combined_crypto_report.pptx"
        else:
            print("\n⚠️ Analysis failed:", analysis_result.get("error", "Unknown error"))
            result["analysis_presentation"] = None
            result["analysis_error"] = analysis_result.get("error", "Unknown error")

        # Check if files exist and handle potential conflicts
        forecast_file = forecast_result if forecast_result else None
        analysis_ppt = 'combined_crypto_report.pptx' if analysis_result and not analysis_result.get("error") else None
        if not forecast_file or not os.path.exists(forecast_file):
            print("\n⚠️ Note: Forecast presentation couldn't be saved. Check permissions or close PowerPoint.")
        if not analysis_ppt or not os.path.exists(analysis_ppt):
            print("\n⚠️ Note: Analysis presentation couldn't be saved. Check permissions or close PowerPoint.")

        # Automatically merge the two presentations if both are generated successfully
        if forecast_file and analysis_ppt and os.path.exists(forecast_file) and os.path.exists(analysis_ppt):
            print("\n🔗 Merging the two presentations into one...")
            merged_file = f'pptx/combined_{coin}_report.pptx'
            merge_presentations(
                primary_pptx_path=forecast_file,  # e.g., pptx/Bitcoin_forecast_report.pptx
                secondary_pptx_path=analysis_ppt,  # e.g., combined_crypto_report.pptx
                output_pptx_path=merged_file  # Dynamic output name based on coin
            )
            result["merged_presentation"] = merged_file
        else:
            print("\n⚠️ Cannot merge presentations: One or both files are missing or failed to generate.")
            result["merged_presentation"] = None

        return result

    except Exception as e:
        print(f"\n❌ Process failed: {str(e)}")
        if "Permission denied" in str(e):
            print("\n⚠️ It appears an output file is open in another application (like PowerPoint).")
            print("⚠️ Please close any open presentations and try again.")
        raise HTTPException(status_code=500, detail=f"Process failed: {str(e)}")

@app.get("/download-presentation/")
async def download_presentation(path: str = Query(...)):
    """Endpoint to download a presentation file."""
    try:
        # Validate and sanitize the path to prevent directory traversal attacks
        # This is a simplified version - in production, implement stronger validation
        if ".." in path or path.startswith("/"):
            raise HTTPException(status_code=400, detail="Invalid file path")
        
        # Convert forward slashes to backslashes if on Windows
        # Or vice versa depending on your server OS
        normalized_path = path.replace("/", os.sep)
        
        # Check if file exists
        if not os.path.exists(normalized_path):
            # For demo purposes, create a dummy file if it doesn't exist
            # In production, you would just return an error
            if not os.path.exists(os.path.dirname(normalized_path)):
                os.makedirs(os.path.dirname(normalized_path), exist_ok=True)
            
            # Create a minimal PowerPoint file for testing
            # In a real implementation, this would be handled by your analysis code
            from pptx import Presentation
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = f"Sample Analysis for {os.path.basename(normalized_path).split('_')[1]}"
            subtitle.text = "This is a placeholder presentation for demonstration purposes."
            prs.save(normalized_path)
        
        # Return the file for download
        return FileResponse(
            normalized_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=os.path.basename(normalized_path)
        )
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error downloading file: {str(e)}")

# The merge_presentations function remains unchanged
def merge_presentations(primary_pptx_path, secondary_pptx_path, output_pptx_path):
    """
    Merge slides from secondary_pptx into primary_pptx and save as output_pptx.
    """
    try:
        # Load the primary presentation (Bitcoin_forecast_report.pptx)
        primary_prs = Presentation(primary_pptx_path)
        print(f"Loaded primary presentation with {len(primary_prs.slides)} slides.")

        # Load the secondary presentation (combined_crypto_report.pptx)
        secondary_prs = Presentation(secondary_pptx_path)
        print(f"Loaded secondary presentation with {len(secondary_prs.slides)} slides.")

        # Function to find a matching slide layout in the primary presentation
        def find_matching_layout(secondary_slide, primary_prs):
            secondary_layout = secondary_slide.slide_layout
            for layout in primary_prs.slide_layouts:
                if layout.name == secondary_layout.name:
                    return layout
            for layout in primary_prs.slide_layouts:
                if len(layout.placeholders) == len(secondary_layout.placeholders):
                    return layout
            for layout in primary_prs.slide_layouts:
                if "Blank" in layout.name:
                    return layout
            return primary_prs.slide_layouts[0]

        # Function to safely get RGB color with fallback
        def safe_rgb_color(color, default_rgb=RGBColor(255, 255, 255)):  # Default to white
            return color.rgb if hasattr(color, 'rgb') and color.rgb is not None else default_rgb

        # Function to copy slide background
        def copy_slide_background(secondary_slide, new_slide, prs):
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            fill = secondary_slide.background.fill
            if fill.type == 1:  # Solid fill
                new_slide.background.fill.solid()
                new_slide.background.fill.fore_color.rgb = safe_rgb_color(fill.fore_color)
            elif fill.type == 5:  # BACKGROUND (master-defined)
                print(f"⚠️ Background fill type BACKGROUND (5) detected, copying image if present.")
            elif fill.type in (0, None):  # _NoFill
                pass
            else:
                print(f"⚠️ Unsupported background fill type: {fill.type}, skipping fill.")

            for shape in secondary_slide.shapes:
                if shape.shape_type == 13:  # Picture shape type
                    try:
                        if (shape.left == 0 and shape.top == 0 and 
                            shape.width >= slide_width * 0.9 and 
                            shape.height >= slide_height * 0.9):
                            image_blob = shape.image.blob
                            temp_image_path = "temp_background_image.jpg"
                            with open(temp_image_path, "wb") as f:
                                f.write(image_blob)
                            new_slide.shapes.add_picture(
                                temp_image_path,
                                Inches(0), Inches(0),
                                width=Inches(slide_width / 914400),
                                height=Inches(slide_height / 914400)
                            )
                            os.remove(temp_image_path)
                            print(f"✅ Copied background image to slide.")
                            break
                    except Exception as e:
                        print(f"⚠️ Could not copy background image: {str(e)}")

        # Copy each slide from the secondary presentation to the primary
        for secondary_slide in secondary_prs.slides:
            target_layout = find_matching_layout(secondary_slide, primary_prs)
            new_slide = primary_prs.slides.add_slide(target_layout)

            # Copy the slide background
            copy_slide_background(secondary_slide, new_slide, primary_prs)

            for shape in secondary_slide.shapes:
                if not shape.has_text_frame and not shape.has_table and not shape.has_chart:
                    if shape.shape_type == 13:  # Picture shape type
                        try:
                            image_blob = shape.image.blob
                            temp_image_path = f"temp_image_{shape.element.get('id')}.png"
                            with open(temp_image_path, "wb") as f:
                                f.write(image_blob)
                            new_slide.shapes.add_picture(
                                temp_image_path,
                                shape.left, shape.top,
                                width=shape.width,
                                height=shape.height
                            )
                            os.remove(temp_image_path)
                            print(f"✅ Copied image to slide.")
                        except Exception as e:
                            print(f"⚠️ Could not copy image: {str(e)}")
                    else:
                        try:
                            new_slide.shapes._spTree.append(shape.element)
                        except Exception as e:
                            print(f"⚠️ Could not copy shape: {str(e)}")
                elif shape.has_text_frame:
                    new_shape = new_slide.shapes.add_textbox(
                        left=shape.left, top=shape.top,
                        width=shape.width, height=shape.height
                    )
                    new_shape.text_frame.clear()  # Clear any default text
                    print(f"Copying text frame with {len(shape.text_frame.paragraphs)} paragraphs.")
                    for para in shape.text_frame.paragraphs:
                        new_para = new_shape.text_frame.add_paragraph()
                        new_para.text = para.text
                        new_para.font.size = para.font.size
                        new_para.font.bold = para.font.bold
                        new_para.font.color.rgb = safe_rgb_color(para.font.color)
                        try:
                            new_para.alignment = para.alignment
                        except AttributeError:
                            new_para.alignment = 2  # Fallback to center
                elif shape.has_table:
                    new_table = new_slide.shapes.add_table(
                        rows=len(shape.table.rows),
                        cols=len(shape.table.columns),
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height
                    ).table
                    for row_idx, row in enumerate(shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            new_cell = new_table.cell(row_idx, col_idx)
                            new_cell.text_frame.clear()  # Clear any default text
                            new_cell.text = cell.text
                            new_cell.fill.solid()
                            source_fill = cell.fill
                            if source_fill.type is None or source_fill.type == 'noFill':
                                new_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            else:
                                new_cell.fill.fore_color.rgb = safe_rgb_color(source_fill.fore_color)
                            for para in cell.text_frame.paragraphs:
                                new_para = new_cell.text_frame.paragraphs[0]
                                new_para.text = para.text
                                new_para.font.size = para.font.size
                                new_para.font.bold = para.font.bold
                                new_para.font.color.rgb = safe_rgb_color(para.font.color)
                                try:
                                    new_para.alignment = para.alignment
                                except AttributeError:
                                    new_para.alignment = 2  # Fallback to center

        # Save the merged presentation
        primary_prs.save(output_pptx_path)
        print(f"\n🎉 Merged presentation saved as: {output_pptx_path}")

        # Verify the file exists
        if os.path.exists(output_pptx_path):
            print(f"✅ File successfully created with {len(primary_prs.slides)} slides.")
        else:
            print("\n⚠️ Failed to save the merged presentation. Check permissions or close PowerPoint.")

    except Exception as e:
        print(f"\n❌ Error merging presentations: {str(e)}")
        if "Permission denied" in str(e):
            print("\n⚠️ It appears the output file is open in another application (like PowerPoint).")
            print("⚠️ Please close any open presentations and try again.")

