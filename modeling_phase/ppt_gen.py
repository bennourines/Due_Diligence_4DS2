#version 17 avr avc all data from ines with PowerPoint generation
import os
import numpy as np
import warnings
import requests
import json
import faiss
import pickle
from sentence_transformers import SentenceTransformer
from typing import List, Dict, Any, Optional, Union, Tuple
import glob
import re
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Import functionality from the utility file
from utility_ppt import (
    ppt_keywords_dict, 
    is_presentation_request, 
    extract_presentation_topic,
    apply_theme_to_presentation, 
    create_slide_with_image,
    create_chart_slide
)

# Load environment variables
load_dotenv()
openrouter_api_key = os.getenv("OPENROUTER_API_KEY")

# Suppress warnings
os.environ['TF_ENABLE_ONEDNN_OPTS'] = '0'
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '2'
warnings.filterwarnings('ignore', category=DeprecationWarning)
warnings.filterwarnings('ignore', category=FutureWarning)
warnings.filterwarnings('ignore', module='tensorflow')

# Configuration
MODEL_NAME = "multi-qa-mpnet-base-dot-v1"
OPENROUTER_API_KEY = openrouter_api_key
MISTRAL_MODEL = "meta-llama/llama-4-maverick:free"
FAISS_INDEX_PATH = "./faiss_index_ines/index.faiss"
METADATA_PATH = "./faiss_index_ines/index_metadata.json"
DATA_DIR = r"DATA\cleaned_texts"  # Directory with text files to index
OUTPUT_DIR = "./outputs"  # DWelcome to the RAG and Presentation Generation System

TEMPLATE_PATH = "template_path\Slideworks_due_diligence_OVERVIEW.pptx"



# Helper: sanitize LLM JSON output
def sanitize_presentation(raw: Any) -> Optional[Dict[str,Any]]:
    if not isinstance(raw, dict):
        return None
    slides = raw.get("slides")
    if not isinstance(slides, list):
        return None
    clean = []
    for s in slides:
        if not isinstance(s, dict):
            continue
        title = s.get("title")
        content = s.get("content")
        if title is None or content is None:
            continue
        if isinstance(content, list):
            bullets = [str(x) for x in content if not isinstance(x, type(Ellipsis))]
        else:
            bullets = [str(content)]
        clean.append({"title": str(title), "content": bullets, "notes": str(s.get("notes",""))})
    if not clean:
        return None
    return {"title": str(raw.get("title","Untitled")), "subtitle": str(raw.get("subtitle","")), "slides": clean}




# Load SentenceTransformer model for generating embeddings
def load_embedding_model(model_name: str):
    print(f"Loading SentenceTransformer model: {model_name}")
    return SentenceTransformer(model_name)

# Create FAISS index from documents
def create_faiss_index(data_dir: str, model_name: str, index_path: str, metadata_path: str):
    print(f"Creating FAISS index from documents in {data_dir}...")
    
    # Load model
    model = load_embedding_model(model_name)
    
    # Initialize lists to store document chunks and metadata
    texts = []
    metadata = []
    
    # Process all text files in the data directory
    for file_path in glob.glob(f"{data_dir}/*.txt"):
        file_name = os.path.basename(file_path)
        print(f"Processing {file_name}...")
        
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
            
            # Simple chunking approach - split by paragraphs and limit size
            paragraphs = content.split('\n\n')
            for i, para in enumerate(paragraphs):
                if para.strip():  # Skip empty paragraphs
                    texts.append(para)
                    metadata.append({
                        "text": para,
                        "file_name": file_name,
                        "chunk_id": i
                    })
    
    if not texts:
        print("No texts found to index!")
        return None, None
    
    # Generate embeddings for all chunks
    print(f"Generating embeddings for {len(texts)} chunks...")
    embeddings = model.encode(texts, show_progress_bar=True)
    
    # Create FAISS index - using L2 distance (can be changed to inner product)
    embedding_dim = embeddings.shape[1]
    index = faiss.IndexFlatL2(embedding_dim)
    
    # Add vectors to the index
    index.add(np.array(embeddings).astype('float32'))
    
    # Save index and metadata
    print(f"Saving index to {index_path} and metadata to {metadata_path}")
    os.makedirs(os.path.dirname(index_path), exist_ok=True)
    faiss.write_index(index, index_path)
    
    with open(metadata_path, 'wb') as f:
        pickle.dump(metadata, f)
    
    print(f"Created FAISS index with {index.ntotal} vectors")
    return index, metadata

# Load existing FAISS index
def load_faiss_index(index_path: str, metadata_path: str):
    print(f"Loading FAISS index from {index_path}...")
    
    if not os.path.exists(index_path) or not os.path.exists(metadata_path):
        print("Index or metadata file not found. Creating new index...")
        return None, None
    
    try:
        index = faiss.read_index(index_path)
        
        with open(metadata_path, 'rb') as f:
            metadata = pickle.load(f)
        
        print(f"Loaded index with {index.ntotal} vectors and {len(metadata)} metadata entries")
        return index, metadata
    except Exception as e:
        print(f"Error loading index: {str(e)}")
        return None, None

# Retrieve relevant chunks using FAISS
def retrieve_relevant_chunks(question: str, model, index, metadata, top_k: int = 5):
    try:
        # Generate query embedding
        query_embedding = model.encode([question], show_progress_bar=False)
        
        # Search the index
        distances, indices = index.search(np.array(query_embedding).astype('float32'), top_k)
        
        # Get the relevant chunks
        results = []
        for i, idx in enumerate(indices[0]):
            if idx < len(metadata):
                doc = metadata[idx]
                # Add distance score (convert to similarity score)
                doc_with_score = doc.copy()
                doc_with_score["score"] = float(1.0 / (1.0 + distances[0][i]))  # Convert distance to similarity
                results.append(doc_with_score)
        
        if results:
            print(f"Found {len(results)} relevant chunks using FAISS")
            return results
        else:
            print("No relevant chunks found")
            return []
            
    except Exception as e:
        print(f"Error retrieving chunks: {str(e)}")
        return []

# Generate answer using LLM
def generate_answer_with_llm(question: str, context: str, api_key: str, model: str):
    url = "https://openrouter.ai/api/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    prompt = (
        "You are a highly knowledgeable expert in cryptocurrency, blockchain technology, "
        "and decentralized finance (DeFi), with extensive experience "
        "in both technical and practical aspects of the field. "
        "Your expertise includes understanding of consensus mechanisms, smart contracts, "
        "tokenomics, mining/staking processes, cryptocurrency exchanges, and regulatory "
        "frameworks across different jurisdictions. You will analyze the provided context thoroughly "
        "and respond to questions with precise, factual information supported by the given context. "
        "If the information cannot be found in the context, you will clearly state this limitation rather than making assumptions or providing speculative answers."
        " When answering, use clear, technical language while maintaining accessibility for both beginners and advanced users. "
        "Break down complex concepts when necessary, and provide relevant examples from the context to support your explanations."

        "Use the following context to answer the question. If you cannot find the answer in the context, "
        "explicitly state 'I cannot find this information in the provided context' and do not make up"
        "information that isn't directly supported by the context."

        f"Context: {context}\n\n"
        f"Question: {question}\n\n"
        "Answer: "
    )
    data = {
        "model": model,
        "messages": [
            {"role": "user", "content": prompt}
        ],
        "max_tokens": 5000,
        "temperature": 0.6
    }
    response = requests.post(url, headers=headers, json=data)
    if response.status_code == 200:
        return response.json()["choices"][0]["message"]["content"]
    else:
        raise Exception(f"Error calling LLM: {response.text}")

# Generate structured PowerPoint content from RAG retrieval data
def generate_presentation_from_rag(question: str, relevant_chunks: List[Dict[str, Any]], 
                                  api_key: str, model: str) -> Dict[str, Any]:
    url = "https://openrouter.ai/api/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    
    # Extract context from relevant chunks
    context = "\n\n".join([chunk["text"] for chunk in relevant_chunks])
    
    system_prompt = (
        "You are a professional presentation designer specializing in cryptocurrency and blockchain topics. "
        "Create a structured presentation based on the following question and context. "
        "The presentation should be informative, visually structured, and convey key information "
        "from the context that addresses the question."
        "\n\nFormat your response as a JSON object with the following structure:"
        "\n{"
        "\n  \"title\": \"Main Presentation Title\","
        "\n  \"subtitle\": \"Presentation Subtitle or Author\","
        "\n  \"slides\": ["
        "\n    {"
        "\n      \"title\": \"Slide Title\","
        "\n      \"content\": [\"Bullet point 1\", \"Bullet point 2\", ...],"
        "\n      \"notes\": \"Optional speaker notes with additional context\""
        "\n    },"
        "\n    ..."
        "\n  ]"
        "\n}"
        "\n\nInclude the following slides:"
        "\n1. Title slide (with title and subtitle)"
        "\n2. Introduction/Overview slide (what the presentation will cover)"
        "\n3. 3-5 content slides that answer the question using information from the context"
        "\n4. Conclusion slide with key takeaways"
        "\n\nDon't include any explanations outside the JSON structure."
    )
    
    data = {
        "model": model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Question: {question}\n\nContext from knowledge base:\n{context}"}
        ],
        "max_tokens": 5000,
        "temperature": 0.7
    }
    
    response = requests.post(url, headers=headers, json=data)
    if response.status_code == 200:
        content = response.json()["choices"][0]["message"]["content"]
        # Extract JSON part if there's any text before or after
        try:
            # Find JSON part (between first { and last })
            json_match = re.search(r'({[\s\S]*})', content)
            if json_match:
                content = json_match.group(1)
            return json.loads(content)
        except json.JSONDecodeError as e:
            print(f"Error parsing JSON content: {e}")
            print(f"Raw content: {content}")
            
            # Fallback: Try to create a simple structured presentation
            fallback_content = {
                "title": f"Response to: {question[:50]}...",
                "subtitle": "Generated from Knowledge Base",
                "slides": [
                    {
                        "title": "Question",
                        "content": [question]
                    }
                ]
            }
            
            # Add content from relevant chunks to slides
            for i, chunk in enumerate(relevant_chunks[:5]):
                fallback_content["slides"].append({
                    "title": f"Information {i+1}",
                    "content": [chunk["text"][:200] + "..."] if len(chunk["text"]) > 200 else [chunk["text"]]
                })
                
            fallback_content["slides"].append({
                "title": "Conclusion",
                "content": ["Information retrieved from knowledge base.",
                           "For more detailed information, please review the full text."]
            })
            
            return fallback_content
    else:
        raise Exception(f"Error calling LLM for presentation: {response.text}")

# Create PowerPoint presentation from generated content
def create_presentation(content: Dict[str, Any], template_path: Optional[str] = None) -> Presentation:
    # Load template if provided, otherwise create a new presentation
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = content["title"]
    subtitle.text = content["subtitle"]
    
    # Add content slides
    for slide_content in content["slides"]:
        # Use a layout with a title and content
        bullet_slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Set the slide title
        title = slide.shapes.title
        title.text = slide_content["title"]
        
        # Add bullet points
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        
        # Add content based on type (string or list)
        if isinstance(slide_content["content"], list):
            for i, bullet_text in enumerate(slide_content["content"]):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet_text
                p.level = 0  # Top-level bullet
        else:
            # If content is a string, add it as a single paragraph
            text_frame.text = slide_content["content"]
        
        # Add speaker notes if present
        if "notes" in slide_content and slide_content["notes"]:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_content["notes"]
    
    return prs
# 3) UPDATE: save_presentation to use the actual slide title as filename

def save_presentation(prs: Presentation, presentation_title: str) -> str:
    """
    Saves the Presentation under OUTPUT_DIR using the slide title as filename.
    Invalid filesystem chars get replaced by '_'.
    """
    # replace anything not alphanumeric, space, dash or underscore
    safe_title = re.sub(r'[^A-Za-z0-9 _\-]', '_', presentation_title).strip()
    # collapse multiple spaces or underscores
    safe_title = re.sub(r'[\s_]+', '_', safe_title)
    filename = f"{safe_title}.pptx"

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    path = os.path.join(OUTPUT_DIR, filename)
    prs.save(path)
    print(f"Presentation saved to {path}")
    return path
# Process query with RAG and generate PowerPoint presentation
def process_rag_to_ppt(question: str, template_path: Optional[str] = None) -> Tuple[str, Optional[str]]:
    # Load embedding model
    embedding_model = load_embedding_model(MODEL_NAME)
    
    # Load or create FAISS index
    index, metadata = load_faiss_index(FAISS_INDEX_PATH, METADATA_PATH)
    
    if index is None or metadata is None:
        print("Creating new index...")
        index, metadata = create_faiss_index(DATA_DIR, MODEL_NAME, FAISS_INDEX_PATH, METADATA_PATH)
        if index is None or metadata is None:
            print("Failed to create index. Exiting.")
            return "Failed to create knowledge index. Please try again later.", None
    
    # Retrieve relevant chunks
    relevant_chunks = retrieve_relevant_chunks(question, embedding_model, index, metadata, top_k=5)
    
    if not relevant_chunks:
        print("No relevant chunks found.")
        return "No relevant information found in the knowledge base for your query.", None
    
    print(f"Retrieved {len(relevant_chunks)} relevant chunks for PPT generation")
    
    try:
        # Generate presentation content from RAG results
        presentation_content = generate_presentation_from_rag(
            question, relevant_chunks, OPENROUTER_API_KEY, MISTRAL_MODEL
        )
        
        print("Creating PowerPoint presentation from RAG results...")
        presentation = create_presentation(presentation_content, template_path)
        
        # Generate a filename based on sanitized question
        raw = generate_presentation_from_rag(question, relevant, OPENROUTER_API_KEY, MISTRAL_MODEL)
        pres_data = sanitize_presentation(raw) or {"title": question, "subtitle": "", "slides": []}
        prs = create_presentation(pres_data, template_path)        
        
        # Return a message and the path to the presentation
        ppt_path = save_presentation(prs, pres_data["title"])
        return f"Created presentation based on your query: {question}", ppt_path
        
    except Exception as e:
        print(f"Error generating presentation from RAG: {str(e)}")
        return f"Error generating presentation: {str(e)}", None

# 1) NEW: helper to detect any PPT keyword mention in the user prompt
# -------------------------------------------------------------------
def contains_ppt_keyword(prompt: str) -> bool:
    """
    Return True if any key in ppt_keywords_dict appears as a standalone word
    in the prompt, case‐insensitive.
    """
    for kw in ppt_keywords_dict.keys():
        # use word‐boundary so 'risk' doesn't match 'brisk'
        if re.search(rf'\b{re.escape(kw)}\b', prompt, re.IGNORECASE):
            return True
    return False


# Main function to process the user request
def process_request(prompt: str, template_path: Optional[str] = None) -> Tuple[str, Optional[str]]:
    """
    If the prompt contains one of ppt_keywords_dict OR
    is_presentation_request(prompt) returns True, we generate a PPT.
    Otherwise we return a normal text answer.
    """
    # detect PPT intent
    ppt_intent = contains_ppt_keyword(prompt) or is_presentation_request(prompt)
    if ppt_intent:
        # if is_presentation_request, you may want to strip out any trigger words
        topic = extract_presentation_topic(prompt) if is_presentation_request(prompt) else prompt
        return process_rag_to_ppt(topic, template_path)

    # — otherwise, regular RAG‐driven text answer —
    # load embedding model, index, etc. (your existing code)
    embedding_model = load_embedding_model(MODEL_NAME)
    index, metadata = load_faiss_index(FAISS_INDEX_PATH, METADATA_PATH)
    if index is None or metadata is None:
        index, metadata = create_faiss_index(DATA_DIR, MODEL_NAME, FAISS_INDEX_PATH, METADATA_PATH)
        if index is None or metadata is None:
            return "Failed to create knowledge index. Please try again later.", None

    relevant_chunks = retrieve_relevant_chunks(prompt, embedding_model, index, metadata, top_k=3)
    if not relevant_chunks:
        return "No relevant information found in the knowledge base for your query.", None

    context = "\n".join([chunk["text"] for chunk in relevant_chunks])
    try:
        answer = generate_answer_with_llm(prompt, context, OPENROUTER_API_KEY, MISTRAL_MODEL)
        return answer, None
    except Exception as e:
        return f"Error generating answer: {str(e)}", None

def process_rag_to_ppt(question: str, template_path: Optional[str] = None) -> Tuple[str, Optional[str]]:
    embedding_model = load_embedding_model(MODEL_NAME)
    index, metadata = load_faiss_index(FAISS_INDEX_PATH, METADATA_PATH)
    if index is None:
        index, metadata = create_faiss_index(DATA_DIR, MODEL_NAME, FAISS_INDEX_PATH, METADATA_PATH)
        if index is None:
            return "Failed to create knowledge index.", None
    relevant = retrieve_relevant_chunks(question, embedding_model, index, metadata, top_k=5)
    if not relevant:
        return "No relevant information found.", None
    raw = generate_presentation_from_rag(question, relevant, OPENROUTER_API_KEY, MISTRAL_MODEL)
    pres_data = sanitize_presentation(raw) or {"title": question, "subtitle": "", "slides": []}
    prs = create_presentation(pres_data, template_path)
    ppt_path = save_presentation(prs, pres_data.get("title", question))
    return f"Created presentation: {pres_data.get('title')}", ppt_path

 # Example usage
if __name__ == "__main__":
    print("\nWelcome to the RAG and Presentation Generation System\n")
    user_prompt = input("Enter your question or presentation request:\n> ")

    response, ppt_path = process_request(user_prompt, TEMPLATE_PATH)

    print("\nResponse:", response)
    if ppt_path:
        print(f"\n✅ Your PPTX has been generated here:\n   {ppt_path}")
    print("\nWelcome to the RAG and Presentation Generation System\n")
    query = "Create a presentation about security best practices for cryptocurrency wallets"
    msg, ppt_path = process_rag_to_ppt(query, TEMPLATE_PATH)
    print(msg)
    if ppt_path:
        print(f"Generated PPT at: {ppt_path}")

    # Test standard query
    ans, _ = process_request("How safe and reliable are online and virtual payment and wallet platforms for cryptocurrency transactions?")
    print(ans)

