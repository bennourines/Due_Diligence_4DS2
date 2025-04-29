# pptx_utilities.py
"""
Utility functions for PowerPoint generation, keyword detection, and image generation.
"""

import re
import os
import json
import tempfile
import urllib.request
from PIL import Image, ImageDraw, ImageFont
from typing import Dict, List, Any, Optional, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import numpy as np
import random

# PowerPoint presentation keywords dictionary
ppt_keywords_dict = {
    "create_presentation": [
        "create presentation",
        "generate ppt",
        "generate presentation",
        "make slides",
        "create slides",
        "build slides",
        "design slides",
        "create ppt",
        "make a ppt",
        "generate slides",
        "make a presentation",
        "create slides",
        "build presentation",
        "design powerpoint",
        "prepare slides",
        "make ppt",
        "draft presentation",
        "presentation generation",
        "slide creation",
        "generate PowerPoint",
        "PowerPoint automation"
    ],
    "file_handling": [
        "ppt file",
        "PowerPoint file",
        "save ppt",
        "open pptx",
        "export presentation",
        "download slides",
        "read pptx",
        "load presentation",
        "save as pptx",
        "PowerPoint format"
    ],
    "presentation_elements": [
        "title slide",
        "slide content",
        "add slide",
        "insert image",
        "bullet points",
        "slide layout",
        "animations",
        "charts",
        "text box",
        "presentation theme"
    ],
    "tools_and_libraries": [
        "python-pptx",
        "PowerPoint API",
        "pptx generation",
        "slide automation",
        "presentation tool",
        "template engine",
        "presentation software"
    ],
    "presentation_purposes": [
        "business presentation",
        "project report",
        "academic slides",
        "marketing pitch",
        "final year project presentation",
        "conference slides",
        "demo presentation"
    ]
}

def is_presentation_request(prompt: str) -> bool:
    """
    Check if a user prompt is related to creating a presentation.
    
    Args:
        prompt: The user's input text
        
    Returns:
        Boolean indicating if the prompt is presentation-related
    """
    prompt = prompt.lower()
    # Check if any presentation-related keywords are in the prompt
    for category, keywords in ppt_keywords_dict.items():
        for keyword in keywords:
            if keyword.lower() in prompt:
                return True
    return False

def extract_presentation_topic(prompt: str) -> str:
    """
    Extract the actual presentation topic from a prompt that may contain
    presentation-related keywords.
    
    Args:
        prompt: The user's input text
        
    Returns:
        The cleaned presentation topic
    """
    # Start with the original prompt
    content_prompt = prompt.strip()
    
    # Remove common presentation-related phrases
    for category, keywords in ppt_keywords_dict.items():
        for keyword in keywords:
            content_prompt = content_prompt.replace(keyword, "").strip()
    
    # Handle "about" or "on" phrases often used in presentation requests
    content_prompt = re.sub(r'\s+about\s+', ' ', content_prompt)
    content_prompt = re.sub(r'\s+on\s+', ' ', content_prompt)
    
    # If the prompt is too stripped down, use original
    if len(content_prompt) < 10:
        return prompt
    
    return content_prompt

def apply_theme_to_presentation(presentation: Presentation, theme: Dict[str, Any]) -> None:
    """
    Apply a custom theme to a presentation.
    
    Args:
        presentation: The PowerPoint presentation object
        theme: Dictionary containing theme settings
    """
    # Apply theme colors to slide master
    for slide in presentation.slides:
        # Apply background color if specified
        if "background_color" in theme:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*theme["background_color"])
        
        # Apply text colors and fonts to shapes
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
                
            for paragraph in shape.text_frame.paragraphs:
                if "text_color" in theme:
                    paragraph.font.color.rgb = RGBColor(*theme["text_color"])
                if "font" in theme:
                    paragraph.font.name = theme["font"]
                if "font_size" in theme:
                    paragraph.font.size = Pt(theme["font_size"])

def generate_placeholder_image(
    topic: str, 
    width: int = 800, 
    height: int = 600, 
    image_dir: str = "generated_images"
) -> str:
    """
    Generate a placeholder image with text about the topic.
    
    Args:
        topic: The topic to generate an image for
        width: Image width
        height: Image height
        image_dir: Directory to save generated images
        
    Returns:
        Path to the generated image
    """
    try:
        # Ensure the image directory exists
        os.makedirs(image_dir, exist_ok=True)
        
        # Create a random background color (light shades)
        r = random.randint(200, 240)
        g = random.randint(200, 240)
        b = random.randint(200, 240)
        
        # Create an image with a solid background
        img = Image.new('RGB', (width, height), color=(r, g, b))
        draw = ImageDraw.Draw(img)
        
        # Try to load a font, fallback to default if not available
        try:
            # Adjust font path based on your system
            font_path = None
            # Try some common system fonts
            common_fonts = [
                "arial.ttf", 
                "Arial.ttf",
                "DejaVuSans.ttf", 
                "DejaVuSans-Bold.ttf", 
                "FreeSans.ttf"
            ]
            
            # Look in common font directories
            font_dirs = [
                "/usr/share/fonts/truetype/",
                "C:\\Windows\\Fonts\\",
                "/System/Library/Fonts/"
            ]
            
            for font_dir in font_dirs:
                for font_name in common_fonts:
                    if os.path.exists(os.path.join(font_dir, font_name)):
                        font_path = os.path.join(font_dir, font_name)
                        break
                if font_path:
                    break
            
            if font_path:
                title_font = ImageFont.truetype(font_path, 40)
                text_font = ImageFont.truetype(font_path, 30)
            else:
                # Fallback to default font
                title_font = ImageFont.load_default()
                text_font = ImageFont.load_default()
        except Exception:
            title_font = ImageFont.load_default()
            text_font = ImageFont.load_default()
        
        # Define colors
        title_color = (30, 30, 100)  # Dark blue
        text_color = (50, 50, 50)    # Dark gray
        
        # Draw a border
        border_width = 10
        draw.rectangle(
            [(0, 0), (width-1, height-1)], 
            outline=(r-30, g-30, b-30), 
            width=border_width
        )
        
        # Draw the title
        title = f"Slide Topic: {topic}"
        title_width = draw.textlength(title, font=title_font) if hasattr(draw, 'textlength') else len(title) * 15
        title_position = ((width - title_width) / 2, 100)
        draw.text(title_position, title, fill=title_color, font=title_font)
        
        # Draw some decorative elements
        # Horizontal line under title
        draw.line([(width/4, 180), (width*3/4, 180)], fill=title_color, width=3)
        
        # Generate some bullet points based on the topic
        bullet_points = [
            f"Key information about {topic}",
            f"Important aspects to consider",
            f"Relevant data and statistics"
        ]
        
        # Draw bullet points
        for i, point in enumerate(bullet_points):
            y_position = 250 + i * 60
            draw.text((width/4, y_position), "• " + point, fill=text_color, font=text_font)
        
        # Add a footer
        footer_text = "Generated for presentation purposes"
        draw.text((width/2 - 150, height - 50), footer_text, fill=text_color, font=text_font)
        
        # Create a unique filename
        safe_topic = re.sub(r'[^\w\- ]', '_', topic).strip().replace(' ', '_')
        filename = f"{safe_topic}_{random.randint(1000, 9999)}.png"
        filepath = os.path.join(image_dir, filename)
        
        # Save the image
        img.save(filepath)
        print(f"[✓] Generated placeholder image: {filepath}")
        
        return filepath
    except Exception as e:
        print(f"[⚠️] Image generation failed: {e}")
        return None

def fetch_image_to_temp(url: str) -> str:
    """
    Download an image URL to a temporary file and return its local path.
    """
    suffix = os.path.splitext(url)[1] or ".png"
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    try:
        urllib.request.urlretrieve(url, temp_file.name)
        # Check if the downloaded file is a valid image
        from PIL import Image
        try:
            Image.open(temp_file.name).verify()
        except Exception as e:
            print(f"[⚠️] Downloaded file is not a valid image: {url} - {e}")
            os.remove(temp_file.name)  # Remove invalid file
            return None
        return temp_file.name
    except Exception as e:
        print(f"[⚠️] Failed to download image {url}: {e}")
        if os.path.exists(temp_file.name):
            try:
                os.remove(temp_file.name)  # Clean up if download failed
            except Exception:
                pass  # Ignore if file cannot be deleted
        return None

def create_slide_with_image(
    presentation: Presentation,
    title: str,
    image_source: str = None,
    caption: Optional[str] = None,
    generate_if_missing: bool = True
) -> None:
    """
    Add a slide with an image (local, remote, or generated) to the presentation.
    
    Args:
        presentation: The PowerPoint presentation object
        title: The slide title
        image_source: URL or path to image (optional)
        caption: Optional caption for the image
        generate_if_missing: Generate an image if source is missing or invalid
    """
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    image_path = None
    
    # Try to use provided image source if it exists
    if image_source:
        if image_source.startswith('http'):
            try:
                image_path = fetch_image_to_temp(image_source)
            except Exception as e:
                print(f"[⚠️] Failed to download image {image_source}: {e}")
        elif os.path.exists(image_source):
            image_path = image_source
    
    # Generate a placeholder image if needed and requested
    if not image_path and generate_if_missing:
        print(f"[🖼️] Generating placeholder image for: {title}")
        image_path = generate_placeholder_image(title)
    
    # Add the image to the slide if we have one
    if image_path and os.path.exists(image_path):
        left, top = Inches(2), Inches(2)
        slide.shapes.add_picture(image_path, left, top)
        if caption:
            left, top, width, height = Inches(2), Inches(5), Inches(6), Inches(1)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text_frame.text = caption
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    else:
        print(f"[⚠️] No image available for slide: {title}")


def create_chart_slide(
    presentation: Presentation,
    title: str,
    chart_info: Dict[str, Any]
) -> None:
    """
    Add a slide with a chart based on provided data.
    """
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title

    # Prepare chart data
    data = ChartData()
    labels = chart_info.get("data", {}).get("labels", [])
    values = chart_info.get("data", {}).get("values", [])
    data.categories = labels
    data.add_series(chart_info.get("title", "Series 1"), values)

    # Map type
    type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE
    }
    chart_type = type_map.get(chart_info.get("type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)

    # Insert chart
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
    graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, data)
    chart = graphic_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_info.get("title", "")


    # --- New Utility Functions for Evaluation ---

def calculate_schema_compliance_rate(results: List[Dict[str, Any]], required_fields: List[str]) -> float:
    """
    Calculate how many generated results comply with a required JSON schema.

    Args:
        results: List of JSON objects generated by the model
        required_fields: List of fields that must be present

    Returns:
        Schema-Compliance Rate as a float between 0 and 1
    """
    if not results:
        return 0.0

    compliant = 0
    for result in results:
        if all(field in result for field in required_fields):
            compliant += 1

    compliance_rate = compliant / len(results)
    print(f"[📊] Schema-Compliance Rate: {compliance_rate:.2%}")
    return compliance_rate


def calculate_precision_at_k(retrieved_chunks: List[Any], relevant_chunks: List[Any], k: int = 5) -> float:
    """
    Calculate Precision@K for retrieved content.

    Args:
        retrieved_chunks: List of retrieved chunks
        relevant_chunks: List of ground-truth relevant chunks
        k: How many top results to consider

    Returns:
        Precision@K as a float between 0 and 1
    """
    if not retrieved_chunks or not relevant_chunks:
        return 0.0

    top_k_retrieved = retrieved_chunks[:k]
    relevant_retrieved = [chunk for chunk in top_k_retrieved if chunk in relevant_chunks]

    precision = len(relevant_retrieved) / k
    print(f"[📈] Precision@{k}: {precision:.2%}")
    return precision
