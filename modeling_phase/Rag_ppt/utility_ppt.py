# pptx_utilities.py
"""
Utility functions for PowerPoint generation and keyword detection.
"""

import re
import os
from typing import Dict, List, Any, Optional, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# PowerPoint presentation keywords dictionary
ppt_keywords_dict = {
    "create_presentation": [
        "generate ppt",
        "make a presentation",
        "create slides",
        "build presentation",
        "design powerpoint",
        "prepare slides",
        "make ppt",
        "draft presentation",
        "presentation generation",
        "slide creation",
        "generate PowerPoint",
        "PowerPoint automation"
    ],
    "file_handling": [
        "ppt file",
        "PowerPoint file",
        "save ppt",
        "open pptx",
        "export presentation",
        "download slides",
        "read pptx",
        "load presentation",
        "save as pptx",
        "PowerPoint format"
    ],
    "presentation_elements": [
        "title slide",
        "slide content",
        "add slide",
        "insert image",
        "bullet points",
        "slide layout",
        "animations",
        "charts",
        "text box",
        "presentation theme"
    ],
    "tools_and_libraries": [
        "python-pptx",
        "PowerPoint API",
        "pptx generation",
        "slide automation",
        "presentation tool",
        "template engine",
        "presentation software"
    ],
    "presentation_purposes": [
        "business presentation",
        "project report",
        "academic slides",
        "marketing pitch",
        "final year project presentation",
        "conference slides",
        "demo presentation"
    ]
}

def is_presentation_request(prompt: str) -> bool:
    """
    Check if a user prompt is related to creating a presentation.
    
    Args:
        prompt: The user's input text
        
    Returns:
        Boolean indicating if the prompt is presentation-related
    """
    prompt = prompt.lower()
    # Check if any presentation-related keywords are in the prompt
    for category, keywords in ppt_keywords_dict.items():
        for keyword in keywords:
            if keyword.lower() in prompt:
                return True
    return False

def extract_presentation_topic(prompt: str) -> str:
    """
    Extract the actual presentation topic from a prompt that may contain
    presentation-related keywords.
    
    Args:
        prompt: The user's input text
        
    Returns:
        The cleaned presentation topic
    """
    # Start with the original prompt
    content_prompt = prompt.strip()
    
    # Remove common presentation-related phrases
    for category, keywords in ppt_keywords_dict.items():
        for keyword in keywords:
            content_prompt = content_prompt.replace(keyword, "").strip()
    
    # Handle "about" or "on" phrases often used in presentation requests
    content_prompt = re.sub(r'\s+about\s+', ' ', content_prompt)
    content_prompt = re.sub(r'\s+on\s+', ' ', content_prompt)
    
    # If the prompt is too stripped down, use original
    if len(content_prompt) < 10:
        return prompt
    
    return content_prompt

def apply_theme_to_presentation(presentation: Presentation, theme: Dict[str, Any]) -> None:
    """
    Apply a custom theme to a presentation.
    
    Args:
        presentation: The PowerPoint presentation object
        theme: Dictionary containing theme settings
    """
    # Apply theme colors to slide master
    for slide in presentation.slides:
        # Apply background color if specified
        if "background_color" in theme:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*theme["background_color"])
        
        # Apply text colors and fonts to shapes
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
                
            for paragraph in shape.text_frame.paragraphs:
                if "text_color" in theme:
                    paragraph.font.color.rgb = RGBColor(*theme["text_color"])
                if "font" in theme:
                    paragraph.font.name = theme["font"]
                if "font_size" in theme:
                    paragraph.font.size = Pt(theme["font_size"])

def create_slide_with_image(presentation: Presentation, title: str, image_path: str, 
                           caption: Optional[str] = None) -> None:
    """
    Add a slide with an image to the presentation.
    
    Args:
        presentation: The PowerPoint presentation object
        title: Title of the slide
        image_path: Path to the image file
        caption: Optional caption text for the image
    """
    # Use a layout with title and content
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    
    # Set the slide title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add the image if the file exists
    if os.path.exists(image_path):
        left = Inches(2)
        top = Inches(2)
        slide.shapes.add_picture(image_path, left, top)
        
        # Add caption if provided
        if caption:
            left = Inches(2)
            top = Inches(5)
            width = Inches(6)
            height = Inches(1)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text_frame.text = caption
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_chart_slide(presentation: Presentation, title: str, chart_data: Dict[str, Any]) -> None:
    """
    Add a slide with a chart to the presentation.
    
    Args:
        presentation: The PowerPoint presentation object
        title: Title of the slide
        chart_data: Dictionary containing chart data and configuration
    """
    # Use a layout with title and content
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    
    # Set the slide title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Define chart position and size
    x, y = Inches(2), Inches(2)
    cx, cy = Inches(6), Inches(4.5)
    
    # Add chart based on chart type
    chart_type = chart_data.get("type", "bar")
    
    if chart_type == "bar":
        chart = slide.shapes.add_chart(
            chart_type=1,  # 1 is for BAR_CLUSTERED
            x=x, y=y, cx=cx, cy=cy
        ).chart
    elif chart_type == "line":
        chart = slide.shapes.add_chart(
            chart_type=4,  # 4 is for LINE_MARKERS
            x=x, y=y, cx=cx, cy=cy
        ).chart
    elif chart_type == "pie":
        chart = slide.shapes.add_chart(
            chart_type=5,  # 5 is for PIE
            x=x, y=y, cx=cx, cy=cy
        ).chart
    else:
        # Default to column chart
        chart = slide.shapes.add_chart(
            chart_type=0,  # 0 is for COLUMN_CLUSTERED
            x=x, y=y, cx=cx, cy=cy
        ).chart
    
    # Set chart data and labels
    # Note: This is simplified; in a real application you would need to
    # populate chart data based on chart_data dictionary contents
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_data.get("chart_title", "Chart")