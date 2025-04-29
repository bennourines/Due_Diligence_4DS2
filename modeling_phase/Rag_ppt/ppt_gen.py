import os
import re
import json
import argparse
import requests
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv

# Load API key from .env
load_dotenv()
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
OPENROUTER_API_URL = "https://openrouter.ai/api/v1/chat/completions"

# Defaults
DEFAULT_TEMPLATE = os.path.join(os.path.dirname(__file__), "Minimalist light sales pitch.pptx")
OUTPUT_DIR = "outputs"
DEFAULT_FONT_SIZE = Pt(32)  # Large font for Pecha Kucha style
TITLE_FONT_SIZE = Pt(40)

# Colors
TITLE_COLOR = RGBColor(44, 62, 80)  # Dark blue
TEXT_COLOR = RGBColor(52, 73, 94)   # Slate

SYSTEM_PROMPT = (
    "You are a Pecha Kucha presentation specialist."
    " Create a visually striking, minimalist slide deck with exactly one key idea per slide."
    " Total slides: exactly 8 slides (including title and conclusion)."
    " Each slide should have ONLY ONE short, powerful statement of 3-5 words."
    " No bullet points. Just one impactful statement per slide."
    " Add brief speaker notes with details for each slide."
    " Use formal but impactful business language.\n"
    "Return ONLY valid JSON with the structure:\n"
    "{\n"
    "  \"title\": \"Presentation Title\",\n"
    "  \"subtitle\": \"Subtitle or Presenter Name\",\n"
    "  \"slides\": [\n"
    "    { \"title\": \"Single Powerful Statement\", \"notes\": \"Speaker notes here\" }\n"
    "  ]\n"
    "}"
)

def setup_args():
    parser = argparse.ArgumentParser(
        description="Generate Pecha Kucha-style PowerPoint via AI"
    )
    parser.add_argument('-t', '--topic', required=True, help='Presentation topic')
    parser.add_argument('-c', '--context', default='', help='Optional background/context')
    parser.add_argument('-m', '--model', default='anthropic/claude-3-haiku', help='AI model')
    parser.add_argument('-p', '--template', default=DEFAULT_TEMPLATE,
                        help='Path to PPTX template')
    return parser.parse_args()

def sanitize_filename(name: str) -> str:
    clean = re.sub(r'[\\/:*?"<>|]', '_', name)
    return clean.strip().replace(' ', '_')

def query_ai(prompt: str, user_text: str, model: str) -> str:
    if not OPENROUTER_API_KEY:
        raise RuntimeError("Missing OPENROUTER_API_KEY in environment.")
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": prompt},
            {"role": "user", "content": user_text}
        ]
    }
    resp = requests.post(OPENROUTER_API_URL, headers=headers, json=payload)
    resp.raise_for_status()
    return resp.json()["choices"][0]["message"]["content"]

def parse_json(resp: str) -> dict:
    match = re.search(r'```(?:json)?\s*([\s\S]*?)```', resp)
    raw = match.group(1) if match else resp
    return json.loads(raw)

def clear_slides(prs: Presentation):
    sld_ids = list(prs.slides._sldIdLst)
    for sld_id in sld_ids:
        prs.slides._sldIdLst.remove(sld_id)

def clean_slide(slide):
    """Remove all shapes except the title placeholder from a slide"""
    shapes_to_remove = []
    for shape in slide.shapes:
        # Keep only the title placeholder
        if shape.name != 'Title' and not shape.has_title:
            shapes_to_remove.append(shape)
    
    # Remove shapes
    for shape in shapes_to_remove:
        sp = shape._sp
        sp.getparent().remove(sp)

def build_presentation(data: dict, template: str) -> Presentation:
    if not os.path.isfile(template):
        raise FileNotFoundError(f"Template not found: {template}")

    prs = Presentation(template)
    clear_slides(prs)

    # Get layouts
    title_layout = None
    content_layout = None
    
    # Find appropriate layouts
    for layout in prs.slide_layouts:
        if layout.name.lower().find('title') >= 0 and title_layout is None:
            title_layout = layout
        if layout.name.lower().find('content') >= 0 and content_layout is None:
            content_layout = layout

    # Fallbacks if specific layouts not found
    if title_layout is None:
        title_layout = prs.slide_layouts[0]
    if content_layout is None:
        content_layout = prs.slide_layouts[1]

    # Title slide
    title_slide = prs.slides.add_slide(title_layout)
    
    # Set the main title
    title_shape = None
    for shape in title_slide.shapes:
        if shape.has_text_frame and shape.name == 'Title':
            title_shape = shape
            break
    
    if title_shape:
        title_shape.text = data.get('title', '')
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = TITLE_FONT_SIZE
                run.font.bold = True
                run.font.color.rgb = TITLE_COLOR
    
    # Process content slides
    slides = data.get('slides', [])[:8]  # Limit to 8 slides
    
    for slide_data in slides:
        sl = prs.slides.add_slide(content_layout)
        
        # Clean the slide of any unwanted content
        clean_slide(sl)
        
        # Find title placeholder
        title_shape = None
        for shape in sl.shapes:
            if shape.has_text_frame and shape.name == 'Title':
                title_shape = shape
                break
        
        # If found, set title
        if title_shape:
            title_shape.text = slide_data.get('title', '')
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = DEFAULT_FONT_SIZE
                    run.font.bold = True
                    run.font.color.rgb = TEXT_COLOR
        
        # Add a central text box for the statement if there's no title shape
        if not title_shape:
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(2)
            
            txbox = sl.shapes.add_textbox(left, top, width, height)
            tf = txbox.text_frame
            tf.word_wrap = True
            
            p = tf.add_paragraph()
            p.text = slide_data.get('title', '')
            p.alignment = PP_ALIGN.CENTER
            
            for run in p.runs:
                run.font.size = DEFAULT_FONT_SIZE
                run.font.bold = True
                run.font.color.rgb = TEXT_COLOR
        
        # Add notes if available
        if slide_data.get('notes'):
            try:
                sl.notes_slide.notes_text_frame.text = slide_data['notes']
            except Exception:
                pass

    return prs

def main():
    args = setup_args()
    user_text = f"Topic: {args.topic}"
    if args.context:
        user_text += f"\nContext: {args.context}"

    print("[AI] Generating Pecha Kucha slide structure...")
    ai_resp = query_ai(SYSTEM_PROMPT, user_text, args.model)
    
    try:
        data = parse_json(ai_resp)
    except json.JSONDecodeError:
        print("[ERROR] Failed to parse AI response as JSON. Using fallback structure.")
        data = {
            "title": args.topic,
            "subtitle": "",
            "slides": [
                {"title": "Key Point 1", "notes": "Expand on key point 1"},
                {"title": "Key Point 2", "notes": "Expand on key point 2"},
                {"title": "Key Point 3", "notes": "Expand on key point 3"},
                {"title": "Key Point 4", "notes": "Expand on key point 4"},
                {"title": "Key Point 5", "notes": "Expand on key point 5"},
                {"title": "Key Point 6", "notes": "Expand on key point 6"},
                {"title": "Conclusion", "notes": "Summarize the key points"}
            ]
        }

    print(f"[IO] Building Pecha Kucha presentation with template: {args.template}")
    prs = build_presentation(data, args.template)

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    fn = sanitize_filename(data.get('title', 'presentation')) + '.pptx'
    out_path = os.path.join(OUTPUT_DIR, fn)
    prs.save(out_path)
    print(f"[IO] Saved Pecha Kucha deck: {out_path}")

if __name__ == '__main__':
    main()