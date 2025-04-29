import os
import numpy as np
import warnings
import requests
import json
import faiss
import pickle
from sentence_transformers import SentenceTransformer
from typing import Any, Dict, List, Optional, Tuple, Union
import glob
import re
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Pt
import tempfile
import urllib.request
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


from metrics import evaluate_generated_slides


# Import utility functions
from utility_ppt_image import (
    ppt_keywords_dict,
    is_presentation_request,
    extract_presentation_topic,
    apply_theme_to_presentation,
    create_slide_with_image,
    create_chart_slide,
    
)

# Load environment variables
load_dotenv()
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")

# Suppress TensorFlow warnings
os.environ.setdefault('TF_ENABLE_ONEDNN_OPTS', '0')
os.environ.setdefault('TF_CPP_MIN_LOG_LEVEL', '2')
warnings.filterwarnings('ignore', category=DeprecationWarning)
warnings.filterwarnings('ignore', category=FutureWarning)

# Configuration
MODEL_NAME = "multi-qa-mpnet-base-dot-v1"
MISTRAL_MODEL = "meta-llama/llama-4-maverick:free"
FAISS_INDEX_PATH = r"G:\Ela\study\4-Ds\S2\PIDS\modeling_phase\Rag_gen\faiss_index_ines\index.faiss"
DATA_DIR = r"G:\Ela\study\4-Ds\S2\PIDS\modeling_phase\data\cleaned_texts"
METADATA_PATH = r"G:\Ela\study\4-Ds\S2\PIDS\modeling_phase\Rag_gen\faiss_index_ines\index_metadata.pkl"
OUTPUT_DIR = r"G:\Ela\study\4-Ds\S2\PIDS\modeling_phase\Rag_gen\outputs"
TEMPLATE_PATH = r"template_path/Slideworks_due_diligence_OVERVIEW.pptx"


def sanitize_presentation(raw: Any) -> Optional[Dict[str, Any]]:
    if not isinstance(raw, dict):
        return None
    slides = raw.get("slides")
    if not isinstance(slides, list):
        return None
    clean_slides = []
    for s in slides:
        if not isinstance(s, dict):
            continue
        # Must have title & content
        title = s.get("title")
        content = s.get("content")
        if title is None or content is None:
            continue

        clean_slide = {
            "title":   str(title),
            "content": [str(b) for b in (content if isinstance(content, list) else [content])],
            "notes":   str(s.get("notes", "")),
        }
        # Preserve the RAG-generated visuals
        if s.get("image"):
            clean_slide["image"] = s["image"]
        if s.get("chart"):
            clean_slide["chart"] = s["chart"]

        clean_slides.append(clean_slide)

    if not clean_slides:
        return None

    return {
        "title":    str(raw.get("title", "Untitled Presentation")),
        "subtitle": str(raw.get("subtitle", "")),
        "slides":   clean_slides,
    }


def load_embedding_model(name: str) -> SentenceTransformer:
    print(f"Loading embedding model '{name}'...")
    return SentenceTransformer(name)


def create_faiss_index(
    data_dir: str,
    model_name: str,
    index_path: str,
    metadata_path: str,
) -> Tuple[Optional[faiss.Index], Optional[List[Dict[str, Any]]]]:
    model = load_embedding_model(model_name)
    texts, meta = [], []
    for fp in glob.glob(os.path.join(data_dir, "*.txt")):
        with open(fp, encoding='utf-8') as f:
            paragraphs = f.read().split('\n\n')
        for i, p in enumerate(paragraphs):
            if p.strip():
                texts.append(p)
                meta.append({"text": p, "source": os.path.basename(fp), "chunk_id": i})
    if not texts:
        print("No documents to index.")
        return None, None
    embeds = model.encode(texts, show_progress_bar=True)
    dim = embeds.shape[1]
    index = faiss.IndexFlatL2(dim)
    index.add(np.array(embeds, dtype='float32'))
    os.makedirs(os.path.dirname(index_path), exist_ok=True)
    faiss.write_index(index, index_path)
    with open(metadata_path, 'wb') as f:
        pickle.dump(meta, f)
    print(f"Built FAISS index with {index.ntotal} vectors.")
    return index, meta


def load_faiss_index(
    index_path: str,
    metadata_path: str,
) -> Tuple[Optional[faiss.Index], Optional[List[Dict[str, Any]]]]:
    if not os.path.exists(index_path) or not os.path.exists(metadata_path):
        return None, None
    index = faiss.read_index(index_path)
    with open(metadata_path, 'rb') as f:
        meta = pickle.load(f)
    print(f"Loaded FAISS index ({index.ntotal} vectors).")
    return index, meta


def retrieve_relevant_chunks(
    question: str,
    model: SentenceTransformer,
    index: faiss.Index,
    metadata: List[Dict[str, Any]],
    top_k: int = 5,
) -> List[Dict[str, Any]]:
    q_emb = model.encode([question], show_progress_bar=False)
    dists, idxs = index.search(np.array(q_emb, dtype='float32'), top_k)
    results = []
    for dist, idx in zip(dists[0], idxs[0]):
        if idx < len(metadata):
            record = metadata[idx].copy()
            record['score'] = float(1 / (1 + dist))
            results.append(record)
    print(f"Retrieved {len(results)} relevant chunks.")
    return results

def generate_presentation_from_rag(
    question: str,
    chunks: List[Dict[str, Any]],
    api_key: str,
    model_name: str,
) -> Dict[str, Any]:
    # Combine context
    context = "\n\n".join(c['text'] for c in chunks)
    # Updated schema: include optional image and chart fields
    system_prompt = (
    "You are a professional presentation designer. "
    "For the given question and context, output ONLY valid JSON with this EXACT schema:\n"
    "  {\n"
    "    title: str,\n"
    "    subtitle: str,\n"
    "    slides: [\n"
    "      {\n"
    "        title: str,\n"
    "        content: [str],\n"
    "        notes: str,\n"
    "        image: str,        # REQUIRED – a URL to a representative image\n"
    "        chart: {           # REQUIRED – at least one chart\n"
    "          type: str,       # 'bar', 'line', or 'pie'\n"
    "          title: str,\n"
    "          data: {\n"
    "            labels: [str],\n"
    "            values: [float]\n"
    "          }\n"
    "        }\n"
    "      }, ...\n"
    "    ]\n"
    "  }\n"
    "Ensure every slide has BOTH `image` and `chart` fields populated."
)

    payload = {
        "model": model_name,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Question: {question}\nContext: {context}"},
        ],
        "max_tokens": 1800,
    }
    resp = requests.post(
        "https://openrouter.ai/api/v1/chat/completions",
        headers={"Authorization": f"Bearer {api_key}"},
        json=payload,
    )
    resp.raise_for_status()
    text = resp.json()['choices'][0]['message']['content']
    match = re.search(r'({[\s\S]*})', text)
    data = json.loads(match.group(1).strip()) if match else json.loads(text.strip())
    return data


def create_presentation(
    content: Dict[str, Any],
    template_path: Optional[str] = None,
) -> Presentation:
    # Use template if available
    prs = Presentation(template_path) if template_path and os.path.exists(template_path) else Presentation()

    # Title slide
    slide0 = prs.slides.add_slide(prs.slide_layouts[0])
    slide0.shapes.title.text = content['title']
    if content.get('subtitle'):
        slide0.placeholders[1].text = content['subtitle']

    # Content slides
    for s in content['slides']:
        # Regular content slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s['title']
        tf = slide.placeholders[1].text_frame
        
        # Add bullet text
        for i, bullet in enumerate(s.get('content', [])):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
        # Notes
        if s.get('notes'):
            slide.notes_slide.notes_text_frame.text = s['notes']
            
        # Add chart if provided
        if s.get('chart'):
            chart_info = s['chart']
            create_chart_slide(prs, chart_info.get('title', s['title']), chart_info)
            
        # Add image slide
        image_source = s.get('image')
        # Skip example.com URLs as they'll always fail
        if image_source and "example.com" in image_source:
            image_source = None
            
        # Create slide with image - will generate a placeholder if needed
        create_slide_with_image(
            prs,
            s['title'],
            image_source=image_source,
            caption=None,
            generate_if_missing=True
        )

    return prs

def save_presentation(prs: Presentation, title: str) -> str:
    safe = re.sub(r'[^\w\- ]', '_', title).strip().replace(' ', '_')
    fname = f"{safe}.pptx"
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    path = os.path.join(OUTPUT_DIR, fname)
    prs.save(path)
    print(f"Presentation saved to {path}")
    return path


def process_rag_to_ppt(
    question: str,
    template: Optional[str] = None,
) -> Tuple[str, Optional[str]]:
    model = load_embedding_model(MODEL_NAME)
    index, meta = load_faiss_index(FAISS_INDEX_PATH, METADATA_PATH)
    if index is None or meta is None:
        index, meta = create_faiss_index(DATA_DIR, MODEL_NAME, FAISS_INDEX_PATH, METADATA_PATH)
        if index is None:
            return "Failed to build index.", None
    chunks = retrieve_relevant_chunks(question, model, index, meta)
    if not chunks:
        return "No relevant data found.", None
    raw = generate_presentation_from_rag(question, chunks, OPENROUTER_API_KEY, MISTRAL_MODEL)
    print("⏺ Raw JSON from LLM:\n", json.dumps(raw, indent=2))
    pres_data = sanitize_presentation(raw) or {"title": question, "subtitle": "", "slides": []}
    
    # Process images before creating presentation
    for slide in pres_data['slides']:
        # Try to use the image URL if provided
        if slide.get('image') and slide['image'].startswith('http'):
            # We'll keep the URL and let create_slide_with_image handle downloading
            pass
        else:
            # Generate an image for this slide topic
            slide['image'] = None  # Clear any invalid or example.com image URL
    
    prs = create_presentation(pres_data, template)
    path = save_presentation(prs, pres_data['title'])
    return f"Created presentation: {pres_data['title']}", path


def contains_ppt_keyword(prompt: str) -> bool:
    return any(re.search(rf"\b{re.escape(k)}\b", prompt, re.IGNORECASE) for k in ppt_keywords_dict)


def process_request(
    prompt: str,
    template: Optional[str] = None,
) -> Tuple[str, Optional[str]]:
    if contains_ppt_keyword(prompt) or is_presentation_request(prompt):
        topic = extract_presentation_topic(prompt) if is_presentation_request(prompt) else prompt
        return process_rag_to_ppt(topic, template)
    # Fallback to normal Q&A...
    return "Only PPT mode supported in this script.", None


if __name__ == '__main__':
    user_prompt = input("Enter question or PPT request:\n> ")
    msg, ppt_path = None, None
    
    model = load_embedding_model(MODEL_NAME)
    index, meta = load_faiss_index(FAISS_INDEX_PATH, METADATA_PATH)
    
    if index is None or meta is None:
        index, meta = create_faiss_index(DATA_DIR, MODEL_NAME, FAISS_INDEX_PATH, METADATA_PATH)
    
    # Retrieve relevant chunks
    retrieved = retrieve_relevant_chunks(user_prompt, model, index, meta)
    
    if not retrieved:
        print("No relevant data found.")
    else:
        # Generate the presentation
        raw = generate_presentation_from_rag(user_prompt, retrieved, OPENROUTER_API_KEY, MISTRAL_MODEL)
        print("⏺ Raw JSON from LLM:\n", json.dumps(raw, indent=2))
        
        pres_data = sanitize_presentation(raw) or {"title": user_prompt, "subtitle": "", "slides": []}
        
        prs = create_presentation(pres_data, TEMPLATE_PATH)
        ppt_path = save_presentation(prs, pres_data['title'])
        
        print(f"✅ Created presentation: {pres_data['title']}")
        print(f"✅ PPT saved at: {ppt_path}")
        
        # Now call evaluation properly
        generated_slides = pres_data['slides']
        ground_truth_answers = []  # <-- If you have ground truth answers, load them here (currently empty)
        retrieved_chunks = [item['text'] for item in retrieved]
        relevant_chunks = [item['text'] for item in retrieved]  # <-- Assuming retrieved ones are relevant for now
        
        evaluate_generated_slides(
            generated_slides=generated_slides,
            ground_truth_answers=ground_truth_answers,
            retrieved_chunks=retrieved_chunks,
            relevant_chunks=relevant_chunks,
            k=5
        )
